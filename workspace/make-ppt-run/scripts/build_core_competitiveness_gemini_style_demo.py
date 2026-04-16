from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    latest_matching_path,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import BLUE, NAVY, WHITE, load_font
from build_neural_params_controlled_demo import draw_centered


DESKTOP = Path("/Users/mac/Desktop")
SOURCE_DOC = DESKTOP / "项目核心竞争力分析.md"
ARCHIVE_DIR = latest_archive_dir()
STYLE_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    DESKTOP / "参考图-大模型能力突然变强.png",
)
BACKGROUND = first_existing_path(
    latest_matching_path(repo_runtime_dir("single-slide-demo"), f"*-神经网络学习-gemini参考图版/神经网络学习的重要参数-Gemini参考图版-底图.png"),
    latest_matching_path(ARCHIVE_DIR, "**/神经网络学习的重要参数-Gemini参考图版-底图.png") if ARCHIVE_DIR else None,
    DESKTOP / "神经网络学习的重要参数-Gemini参考图版-底图.png",
    STYLE_REFERENCE,
)
OUTPUT_ROOT = Path("/Users/mac/Desktop/Make PPt/runtime/image-ppt-demo")

WIDTH = 1920
HEIGHT = 1080

TITLE_FONT = load_font(58, heavy=True)
TITLE_SMALL_FONT = load_font(50, heavy=True)
BODY_FONT = load_font(30)
BODY_SMALL_FONT = load_font(26)
SCREEN_FONT = load_font(28, heavy=True)
SCREEN_AI_FONT = load_font(54, heavy=True)
SPEECH_FONT = load_font(28)
CAPTION_FONT = load_font(19, heavy=True)

SCREEN_CENTERS = [(748, 562), (1168, 562), (1590, 562)]
SPEECH_CENTER = (1704, 340)


@dataclass(frozen=True)
class SlidePage:
    title: str
    bullets: str
    screens: list[str]
    speech: str
    filename: str


PAGES = [
    SlidePage(
        title="项目核心竞争力分析",
        bullets="• 不是做一个 Skill，而是在运营一个 Field\n• 真正的护城河，藏在系统如何成为“生命体”",
        screens=["Skill\n可复制", "Field\n会进化", "护城河\n复制不了"],
        speech="核心差异",
        filename="P1-封面-Field母版风格.png",
    ),
    SlidePage(
        title="引言：我们构建的不是“技能”，而是“场域”",
        bullets="• Skill：输入 → 处理 → 输出，一次性交付\n• Field：执行 → 反馈 → 学习 → 增强，持续进化",
        screens=["线性\n工具", "反馈\n学习", "生命体"],
        speech="越跑越强",
        filename="P2-引言-Field不是Skill.png",
    ),
    SlidePage(
        title="第一层护城河：闭环自进化系统",
        bullets="• 对手是静态工具，我们是动态系统\n• 每一次市场投放，都会反哺下一次决策",
        screens=["投放", "量化\n反馈", "下一轮\n更强"],
        speech="闭环",
        filename="P3-第一层-闭环自进化.png",
    ),
    SlidePage(
        title="第二层护城河：垄断性数据资产",
        bullets="• 对手看到一个点，我们看到整片区域\n• 多点运营最终沉淀为“游客欲望图谱”",
        screens=["单点\n数据", "欲望\n图谱", "定义\n风向"],
        speech="上帝视角",
        filename="P4-第二层-欲望图谱.png",
    ),
    SlidePage(
        title="第三层护城河：人机混合智能",
        bullets="• AI 负责 99% 的重复执行与规模动作\n• 人类保留 1% 的直觉判断与战略选择",
        screens=["AI执行\n99%", "人类判断\n1%", "组合\n智能"],
        speech="10选1",
        filename="P5-第三层-人机混智.png",
    ),
    SlidePage(
        title="第四层护城河：平台化扩张能力",
        bullets="• 对手做单点定制，我们做可复制平台\n• 把成功模式从 1 家样板快速复制到 100 家",
        screens=["1家\n样板", "复制\n100家", "网络\n扩张"],
        speech="零边际",
        filename="P6-第四层-平台扩张.png",
    ),
    SlidePage(
        title="最终结论：What 能复制，How 复制不了",
        bullets="• 别人可以模仿我们做什么，却模仿不了我们如何成为\n• 我们运营的不是技能，而是会学习、会垄断、会扩张的生命体",
        screens=["What", "How", "生命体"],
        speech="隐形上帝",
        filename="P7-结论-What和How.png",
    ),
]


def open_background() -> Image.Image:
    image = Image.open(BACKGROUND).convert("RGB").resize((WIDTH, HEIGHT))
    return Image.blend(image, Image.new("RGB", (WIDTH, HEIGHT), WHITE), 0.10).filter(ImageFilter.SMOOTH)


def draw_screen_text(draw: ImageDraw.ImageDraw, center: tuple[int, int], text: str, is_ai: bool = False) -> None:
    font = SCREEN_AI_FONT if is_ai else SCREEN_FONT
    spacing = 3
    box = draw.multiline_textbbox((0, 0), text, font=font, spacing=spacing, align="center")
    x = center[0] - (box[2] - box[0]) / 2
    y = center[1] - (box[3] - box[1]) / 2
    draw.multiline_text((x, y), text, font=font, fill=BLUE if is_ai else NAVY, spacing=spacing, align="center")


def draw_speech_text(draw: ImageDraw.ImageDraw, text: str) -> None:
    box = draw.textbbox((0, 0), text, font=SPEECH_FONT)
    x = SPEECH_CENTER[0] - (box[2] - box[0]) / 2
    y = SPEECH_CENTER[1] - (box[3] - box[1]) / 2
    draw.text((x, y), text, font=SPEECH_FONT, fill=NAVY)


def add_footer(draw: ImageDraw.ImageDraw, page: int, total: int, caption: str) -> None:
    draw.text((80, 1014), caption, font=CAPTION_FONT, fill="#64748B")
    draw.text((1742, 1012), f"{page:02d}/{total:02d}", font=CAPTION_FONT, fill=BLUE)


def render_page(page: SlidePage, index: int, total: int) -> Image.Image:
    image = open_background()
    draw = ImageDraw.Draw(image)

    draw_centered(draw, page.title, 74, TITLE_SMALL_FONT if len(page.title) > 22 else TITLE_FONT, BLUE)
    draw_centered(draw, page.bullets, 198, BODY_SMALL_FONT if len(page.bullets) > 38 else BODY_FONT, NAVY, max_width=1180, spacing=16)

    for screen_index, (center, text) in enumerate(zip(SCREEN_CENTERS, page.screens)):
        draw_screen_text(draw, center, text, is_ai=(screen_index == 2 and text == "AI"))

    draw_speech_text(draw, page.speech)
    add_footer(draw, index + 1, total, f"项目核心竞争力分析 · 参考图风格页 {index + 1}")
    return image


def create_contact_sheet(images: list[Path], output_path: Path) -> None:
    columns = 3
    rows = (len(images) + columns - 1) // columns
    thumb_w = 480
    thumb_h = 270
    padding = 28
    label_h = 56
    sheet_w = padding + columns * (thumb_w + padding)
    sheet_h = padding + rows * (thumb_h + label_h + padding)
    sheet = Image.new("RGB", (sheet_w, sheet_h), "#111827")
    draw = ImageDraw.Draw(sheet)
    font = load_font(20, heavy=True)

    for idx, image_path in enumerate(images):
        image = Image.open(image_path).convert("RGB").resize((thumb_w, thumb_h))
        row = idx // columns
        col = idx % columns
        x = padding + col * (thumb_w + padding)
        y = padding + row * (thumb_h + label_h + padding)
        sheet.paste(image, (x, y))
        draw.rounded_rectangle((x - 2, y - 2, x + thumb_w + 2, y + thumb_h + 2), radius=8, outline="#475569", width=2)
        draw.text((x, y + thumb_h + 14), image_path.stem, font=font, fill="#FFFFFF")

    sheet.save(output_path, quality=92)


def save_ppt(images: list[Path], output_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output_path))


def write_readme(output_dir: Path) -> None:
    lines = [
        "这是根据《项目核心竞争力分析.md》提炼的 7 页 Gemini 参考图风格图片型 PPT。",
        "生成模式：复用参考图空白底图，仅本地叠加中文标题、屏幕字和气泡字。",
        f"源文档：{SOURCE_DOC}",
        f"参考底图：{BACKGROUND}",
        f"结果目录：{output_dir}",
    ]
    (output_dir / "README.txt").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    args = parse_args()
    if not SOURCE_DOC.exists():
        raise RuntimeError(f"Source markdown not found: {SOURCE_DOC}")
    if BACKGROUND is None or not BACKGROUND.exists():
        raise RuntimeError(f"Reference background not found: {BACKGROUND}")

    output_dir = create_timestamped_output_dir(args.output_root, "项目核心竞争力分析-gemini参考图风格七页样张")

    image_paths: list[Path] = []
    total = len(PAGES)
    for idx, page in enumerate(PAGES):
        image = render_page(page, idx, total)
        path = output_dir / page.filename
        image.save(path, quality=95)
        image_paths.append(path)

    pptx_path = output_dir / "项目核心竞争力分析-Gemini参考图风格-七页样张.pptx"
    save_ppt(image_paths, pptx_path)
    contact_sheet = output_dir / "contact-sheet.png"
    if args.with_contact_sheet:
        create_contact_sheet(image_paths, contact_sheet)
    write_readme(output_dir)
    write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(image_paths),
        sourceDoc=str(SOURCE_DOC),
        background=str(BACKGROUND),
        readme=str(output_dir / "README.txt"),
        contactSheet=str(contact_sheet) if args.with_contact_sheet else "",
        imageDirectory=str(output_dir),
        mode="gemini-reference-local-overlay",
    )
    print(pptx_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the Gemini-reference-style core-competitiveness PPT deck.")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=OUTPUT_ROOT,
        help="Directory where timestamped local result folders are created.",
    )
    parser.add_argument(
        "--with-contact-sheet",
        action="store_true",
        help="Also write a local contact-sheet preview image.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    main()

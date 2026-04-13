import pptx
import shutil
import re

template_path = '/Users/mac/Desktop/副本AI介绍 漫画.pptx'
out_path = '/Users/mac/Desktop/ClawLink_完美复刻版.pptx'

shutil.copyfile(template_path, out_path)
prs = pptx.Presentation(out_path)

replacements = [
    ("时代的叩问", "ClawLink AI 产业融合架构"),
    ("旧地图，找不到新大陆", "我们不是在进化，我们是在被替换"),
    ("一人公司 = 人 \+ 系统", "努力=收入？技能=价值？"),
    ("组织被压缩，能力被释放", "AI 不是帮你做得更快，而是直接跳过你。"),
    ("极少数人，撬动绝大多数结果", "杠杆与放大器：极少数人撬动绝大部分结果"),
    ("你不需要干活，你需要“调度”", "一人公司 ≠ 一个人单干，组织被压缩！"),
    ("识别机会 → 配置系统 → 收割结果", "硅基特种部队：不需要工资的数字劳动力"),
    ("智能体调度力", "从本金消耗者 变成 本金投资者"),
    ("从“看不懂的文化”到“可计算的资产”", "未来核心能力只有一个：智能体调度力"),
    ("订单式种植，农业的工业化与金融化！", "我们不做单一产品，做产业系统底座"),
    ("我们不是做旅居，我们在控制“人去哪里”", "非遗+AI=资产重构，从无形到可计算"),
    ("谁掌握算力，谁就掌握产业的控制权", "智慧农业：从看天吃饭到算法种地"),
    ("从赚辛苦钱，到赚系统的钱", "AI旅居：控制流量来源、定价、分配"),
    ("这不是快一点，这是快一个维度", "把分散算力变成统一生产力"),
    ("你，准备好做指挥官了吗", "我们不是参与分配，我们在重写分配"),
    ("努力", "ClawLink"),
    ("技能", "AI"),
    ("系统", "智能体")
]

# Quick text replace function
def replace_text_in_shape(shape):
    if hasattr(shape, "has_text_frame") and not shape.has_text_frame:
        return
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old_text, new_text in replacements:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

# Iterate over all shapes in all slides
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_in_shape(shape)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_text_in_shape(cell)
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                replace_text_in_shape(child_shape)

prs.save(out_path)
print(f"Saved {out_path}")

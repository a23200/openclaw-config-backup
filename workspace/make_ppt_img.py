import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

out_path = '/Users/mac/Desktop/ClawLink_产业融合架构_漫画精美版.pptx'
prs = pptx.Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

content = [
    ("ClawLink AI 产业融合架构", "我们不是在进化，我们是在“被替换”", "/Users/mac/.openclaw/media/tool-image-generation/slide1_cyber_city---21d4afea-c9bc-4609-8783-6fd5f39f8a31.png"),
    ("残酷的真相", "AI 不是帮你做得更快。\n它是直接问你：这事还需要你做吗？\n你不会被淘汰，你会被直接跳过。", "/Users/mac/.openclaw/media/tool-image-generation/slide3_robot_hand---bf0977f8-5ed2-45dd-b234-b204b00ce9d4.png"),
    ("一人公司的觉醒", "一人公司 ≠ 一个人单干\n一人公司 = 人 + 系统\n组织被压缩，能力被释放！", "/Users/mac/.openclaw/media/tool-image-generation/slide5_one_man_company---37fafefb-92aa-44f9-882d-f3f67820f1ed.png"),
    ("硅基员工系统", "ClawLink 硅基特种部队：\n不需要工资、休息的数字化劳动力。\n你不需要干活，你需要“调度”。", "/Users/mac/.openclaw/media/tool-image-generation/slide6_robot_army---ff0c1003-0f04-498a-9bf8-105deded9981.jpg")
]

blank_layout = prs.slide_layouts[6] # Blank layout

for title, body, img_path in content:
    slide = prs.slides.add_slide(blank_layout)
    
    # Set background image
    try:
        slide.shapes.add_picture(img_path, 0, 0, width=Inches(13.33), height=Inches(7.5))
    except Exception as e:
        print(f"Failed to load image {img_path}: {e}")

    # Add dark semi-transparent overlay for text readability
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.transparency = 0.6  # 60% transparent
    overlay.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(64)
    p.font.color.rgb = RGBColor(255, 204, 0) # Cyberpunk Yellow
    p.font.name = "Arial Black"

    # Body
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for line in body.split('\n'):
        p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.bold = True
        p2.font.size = Pt(40)
        p2.font.color.rgb = RGBColor(255, 255, 255) # White
        p2.font.name = "Arial"

prs.save(out_path)
print(f"Saved {out_path}")

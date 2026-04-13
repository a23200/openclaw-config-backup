import sys

try:
    import docx
    import pptx
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx", "python-pptx", "openai", "markdown2"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    import docx
    import pptx

def extract_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extract_pptx(path):
    prs = pptx.Presentation(path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_text = f"Slide {i+1}:\n"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        slide_text += run.text.strip() + " "
        text_runs.append(slide_text)
    return "\n\n".join(text_runs)

try:
    print("--- DOCX CONTENT ---")
    print(extract_docx('/Users/mac/Desktop/副本ClawLink AI产业融合架构.docx')[:3000])
    print("\n--- PPTX CONTENT ---")
    print(extract_pptx('/Users/mac/Desktop/副本AI介绍 漫画.pptx')[:3000])
except Exception as e:
    print(f"Error: {e}")

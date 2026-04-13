import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor

# 模板路径
template_path = '/Users/mac/Desktop/副本AI介绍 漫画.pptx'
out_path = '/Users/mac/Desktop/ClawLink_产业融合架构_漫画版.pptx'

# 加载模板
prs = pptx.Presentation(template_path)

# 我们直接在最后追加新的幻灯片，或者使用第一种布局创建新幻灯片
# 为了简单起见，我们将之前生成的大纲内容写入使用第一种布局的新幻灯片，
# 并在最后删除原有的幻灯片（可选），或者只是把新内容加进去。

# 清空原有幻灯片可能比较麻烦（python-pptx 不直接支持删除），
# 所以我们最好基于它的 master 创建一个新的演示文稿。
# 其实 prs = pptx.Presentation(template_path) 就是基于它创建的。

content = [
    ("ClawLink AI 产业融合架构", "我们不是在进化，我们是在“被替换”"),
    ("时代的叩问", "努力 = 收入？\n技能 = 价值？\n旧地图，找不到新大陆。"),
    ("残酷的真相", "AI 不是帮你做得更快。\n它是直接问你：这事还需要你做吗？\n你不会被淘汰，你会被直接跳过。"),
    ("新的财富法则", "杠杆与放大器：\n极少数人，撬动绝大多数结果。"),
    ("一人公司的觉醒", "一人公司 ≠ 一个人单干\n一人公司 = 人 + 系统\n组织被压缩，能力被释放！"),
    ("硅基员工系统", "ClawLink 硅基特种部队：\n不需要工资、不需要休息的数字化劳动力。\n你不需要干活，你需要“调度”。"),
    ("身份的跃迁", "从“本金消耗者”\n变成“本金投资者”\n识别机会 → 配置系统 → 收割结果"),
    ("决战：从做事到闭环", "未来核心能力只有一个：\n智能体调度力。"),
    ("ClawLink 的野望", "我们不做单一产品，\n我们做产业操作系统的“底座”。"),
    ("降维打击 1 - 非遗文化资产化", "非遗 + AI = 资产重构\n从“看不懂的文化”到“可计算的资产”"),
    ("降维打击 2 - 智慧农业", "从“看天吃饭”到“算法种地”\n订单式种植，农业的工业化与金融化！"),
    ("降维打击 3 - AI 旅居流量控制", "控制三权：流量来源、定价、分配\n我们不是做旅居，我们在控制“人去哪里”。"),
    ("终极底牌 - 分布式算力网", "把分散算力，变成统一生产力。\n谁掌握算力，谁就掌握产业的控制权。"),
    ("赚钱逻辑的颠覆", "我们不是参与分配，我们在重写分配。\n从赚辛苦钱，到赚系统的钱。"),
    ("唯一的护城河", "天下财富，唯快能取。\n这不是快一点，这是快一个维度。"),
    ("终局之战", "当别人还在点灯找路，我们已让风起云翻；\n当 AI 握着超级算力下场，人只剩选择站边。\n你，准备好做指挥官了吗？")
]

# 尽量找一个比较空白的版式（通常版式 1 或 6）
title_slide_layout = prs.slide_layouts[0]
try:
    content_slide_layout = prs.slide_layouts[1]
except IndexError:
    content_slide_layout = prs.slide_layouts[0]

for title, body in content:
    slide = prs.slides.add_slide(content_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # 找到正文框
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body_shape = shape
            break
            
    if body_shape and body_shape.has_text_frame:
        body_shape.text_frame.text = body
    else:
        # 如果没有占位符，手动加一个文本框
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.text = body

# 由于 python-pptx 不能直接删除原始幻灯片，我们会把新幻灯片加在后面。
# 如果你想清理前面的旧幻灯片，我们需要操作底层的 XML。
# 下面是一个清理旧幻灯片的 hack：
xml_slides = prs.slides._sldIdLst
slides_to_keep = list(xml_slides)[-len(content):] # 刚刚加的这几页
for sldId in list(xml_slides)[:-len(content)]:
    xml_slides.remove(sldId)

prs.save(out_path)
print(f"Saved {out_path}")

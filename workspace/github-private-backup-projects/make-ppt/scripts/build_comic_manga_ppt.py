from __future__ import annotations

import argparse
import base64
import binascii
import ipaddress
import json
import os
import random
import re
import socket
import subprocess
import textwrap
import time
import urllib.parse
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    latest_matching_path,
    load_env_candidates,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import BLUE, NAVY, WHITE, load_font, wrap_text


DEFAULT_GEMINI_IMAGE_MODEL = "gemini-3-pro-image-preview"
DEFAULT_OPENAI_IMAGE_MODEL = "gpt-image-2"
DEFAULT_COMIC_REFERENCE_MODE = "full"
DEFAULT_GEMINI_IMAGE_MAX_ATTEMPTS = 4
DEFAULT_OPENAI_IMAGE_MAX_ATTEMPTS = 8
DEFAULT_GEMINI_RETRY_BASE_DELAY_SECONDS = 2.0
DEFAULT_GEMINI_RETRY_MAX_DELAY_SECONDS = 18.0
DEFAULT_GEMINI_RETRY_JITTER_SECONDS = 0.8
IMAGE_REQUEST_TIMEOUT_SECONDS = 300
OPENAI_COMPATIBLE_IMAGE_REQUEST_TIMEOUT_SECONDS = 900
DEFAULT_IMAGE_CONNECT_TIMEOUT_SECONDS = 20
DEFAULT_IMAGE_LOW_SPEED_TIME_SECONDS = 45
DEFAULT_OPENAI_COMPATIBLE_LOW_SPEED_TIME_SECONDS = 240
DEFAULT_IMAGE_LOW_SPEED_LIMIT_BYTES = 1
DEFAULT_OPENAI_COMPATIBLE_GENERATION_RESPONSE_FORMAT = "url"
DEFAULT_OPENAI_COMPATIBLE_DNS_RESOLVERS = ("1.1.1.1", "8.8.8.8", "192.168.1.1")
WIDTH = 1920
HEIGHT = 1080
ARCHIVE_DIR = latest_archive_dir()
OUTPUT_ROOT = repo_runtime_dir("image-ppt-demo")

STYLE_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    Path("/Users/mac/Desktop/参考图-大模型能力突然变强.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png"),
)
LAYOUT_REFERENCE = first_existing_path(
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini原生出字生图*页样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini图文融合生图*页样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png"),
)

TITLE_FONT = load_font(58, heavy=True)
BODY_FONT = load_font(30)
SMALL_FONT = load_font(24)
OVERLAY_TITLE_FONT = load_font(54, heavy=True)
OVERLAY_TITLE_SMALL_FONT = load_font(44, heavy=True)
OVERLAY_BODY_FONT = load_font(31)
OVERLAY_BODY_SMALL_FONT = load_font(28)
OVERLAY_PAGE_FONT = load_font(22, heavy=True)

IGNORED_SECTION_TITLES = {
    "目录",
    "contents",
    "附录",
    "参考资料",
    "致谢",
}

STRONG_POINT_PATTERNS = [
    r"核心",
    r"本质",
    r"真正",
    r"目标",
    r"意味着",
    r"关键能力",
    r"能力",
    r"模型",
    r"系统",
    r"结构",
    r"升级",
    r"完成",
    r"形成",
    r"路径",
    r"战略",
    r"竞争力",
    r"增长",
    r"承接",
    r"资源",
    r"收益",
    r"控制",
]

WEAK_POINT_PATTERNS = [
    r"^当前，",
    r"^当前,",
    r"^这是一次",
    r"^这不是一次简单的",
    r"^我们正在推进的",
    r"^如果你是这么理解的",
    r"^因为",
    r"^所以",
]

GENERIC_SCENE_KEYWORDS = {
    "capital": "abstract strategic structure scene with layered bridge blocks, upward arrows, executive geometry and clean white space",
    "ai_system": "abstract command center scene with one glowing core, surrounding modules, curved connection lines and calm white space",
    "industry_control": "abstract multi-sector hub scene with four linked zones, one central platform, converging paths and clean white space",
    "profit_engine": "abstract business engine scene with circular loops, layered pathways, scaling arrows and clean white space",
    "cover": "abstract premium cover scene with old structure collapsing, a new core awakening, strong arrows and large white space",
    "general": "abstract editorial business scene with clean geometry, structured arrows, modular shapes and large white space",
}


@dataclass(frozen=True)
class MarkdownSection:
    index: int
    level: int
    title: str
    lines: list[str]


@dataclass(frozen=True)
class SlidePage:
    page: int
    title: str
    bullets: list[str]
    context_title: str
    context_points: list[str]
    scene_prompt: str
    filename: str


class GeminiGenerationError(RuntimeError):
    def __init__(self, message: str, *, retryable: bool = False):
        super().__init__(message)
        self.retryable = retryable


class OpenAICompatibleImageError(RuntimeError):
    def __init__(self, message: str, *, retryable: bool = False):
        super().__init__(message)
        self.retryable = retryable


def emit_progress(stage: str, message: str, percent: int, **extra: object) -> None:
    payload: dict[str, object] = {
        "stage": stage,
        "message": message,
        "percent": max(0, min(100, int(percent))),
    }
    payload.update(extra)
    print(f"__PROGRESS__ {json.dumps(payload, ensure_ascii=False)}", flush=True)


def emit_log(kind: str, title: str, body: str, **extra: object) -> None:
    payload: dict[str, object] = {
        "kind": kind,
        "title": title,
        "body": body,
    }
    if extra:
        payload["extra"] = extra
    print(f"__LOG__ {json.dumps(payload, ensure_ascii=False)}", flush=True)


def clean_markdown_text(text: str) -> str:
    cleaned = str(text or "")
    cleaned = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"`([^`]+)`", r"\1", cleaned)
    cleaned = re.sub(r"\*\*([^*]+)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"\*([^*]+)\*", r"\1", cleaned)
    cleaned = re.sub(r"^\s*[-*+]\s*", "", cleaned.strip())
    cleaned = re.sub(r"^\s*\d+[.)]\s*", "", cleaned.strip())
    cleaned = re.sub(r"^\s*[●•▪◆◇■□]+\s*", "", cleaned.strip())
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip(" ：:|")


def read_positive_int_env(name: str, default: int) -> int:
    raw = str(os.environ.get(name, "")).strip()
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return value if value > 0 else default


def read_non_negative_float_env(name: str, default: float) -> float:
    raw = str(os.environ.get(name, "")).strip()
    if not raw:
        return default
    try:
        value = float(raw)
    except ValueError:
        return default
    return value if value >= 0 else default


def read_csv_env(name: str, default: tuple[str, ...] = ()) -> list[str]:
    raw = str(os.environ.get(name, "")).strip()
    if not raw:
        return [item for item in default if item]
    return [item.strip() for item in re.split(r"[\s,;]+", raw) if item.strip()]


def shorten(text: str, limit: int) -> str:
    cleaned = clean_markdown_text(text)
    return cleaned if len(cleaned) <= limit else f"{cleaned[: limit - 1]}…"


def sanitize_filename_part(text: str) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|`]+", "", text)
    cleaned = re.sub(r"\s+", "", cleaned)
    return cleaned[:24] or "页面"


def mime_type_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    return "application/octet-stream"


def get_api_key() -> str:
    load_env_candidates()
    for key in ["GEMINI_API_KEY", "GOOGLE_API_KEY", "GOOGLE_GENAI_API_KEY"]:
        value = os.environ.get(key)
        if value:
            return value
    raise RuntimeError("Gemini API key not found in env.")


def get_openai_compatible_api_key() -> str:
    load_env_candidates()
    for key in ["COMIC_OPENAI_API_KEY", "COMIC_GPT_IMAGE_API_KEY", "OPENAI_API_KEY", "GPT_IMAGE_API_KEY"]:
        value = os.environ.get(key)
        if value:
            return value
    raise RuntimeError("OpenAI-compatible image API key not found in env.")


def normalize_image_provider(provider: str | None) -> str:
    normalized = str(provider or "").strip().lower().replace("-", "_")
    if normalized in {"gpt", "openai", "openai_compatible", "gpt_image", "gpt_image_2"}:
        return "openai"
    return "gemini"


def normalize_reference_mode(reference_mode: str | None) -> str:
    normalized = str(reference_mode or "").strip().lower().replace("-", "_")
    if normalized in {"prompt", "prompt_only", "pure_prompt", "no_reference", "none"}:
        return "prompt_only"
    if normalized in {"style", "style_only", "style_reference"}:
        return "style_only"
    return "full"


def reference_mode_label(reference_mode: str) -> str:
    normalized = normalize_reference_mode(reference_mode)
    if normalized == "prompt_only":
        return "纯提示词"
    if normalized == "style_only":
        return "仅风格参考"
    return "完整参考"


def openai_compatible_base_url() -> str:
    base_url = (
        os.environ.get("COMIC_OPENAI_BASE_URL")
        or os.environ.get("COMIC_GPT_IMAGE_BASE_URL")
        or os.environ.get("OPENAI_BASE_URL")
        or "https://api.openai.com/v1"
    ).rstrip("/")
    parsed = urllib.parse.urlparse(base_url)
    if parsed.path.rstrip("/").endswith("/v1"):
        return base_url
    return f"{base_url}/v1"


def parse_markdown_sections(markdown: str, source_path: Path) -> tuple[str, list[MarkdownSection]]:
    deck_title = source_path.stem
    sections: list[MarkdownSection] = []
    current_level = 0
    current_title = ""
    current_lines: list[str] = []
    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*$")

    for raw_line in markdown.splitlines():
        stripped = raw_line.strip()
        match = heading_re.match(stripped)
        if match:
            if current_title:
                sections.append(
                    MarkdownSection(
                        index=len(sections),
                        level=current_level,
                        title=current_title,
                        lines=current_lines,
                    )
                )
            current_level = len(match.group(1))
            current_title = clean_markdown_text(match.group(2))
            current_lines = []
            if current_level == 1 and current_title:
                deck_title = current_title
            continue
        if current_title:
            current_lines.append(raw_line)

    if current_title:
        sections.append(
            MarkdownSection(
                index=len(sections),
                level=current_level,
                title=current_title,
                lines=current_lines,
            )
        )

    return deck_title or source_path.stem, sections


def extract_cover_lead_lines(markdown: str, deck_title: str, max_points: int = 2) -> list[str]:
    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
    cover_title = clean_markdown_text(deck_title.split("：", 1)[0])
    full_title = clean_markdown_text(deck_title)
    subtitle = clean_markdown_text(deck_title.split("：", 1)[1]) if "：" in deck_title else ""
    candidates: list[str] = [subtitle] if subtitle else []

    for raw_line in markdown.splitlines():
        stripped = raw_line.strip()
        if heading_re.match(stripped):
            break
        if not stripped:
            continue
        cleaned = clean_markdown_text(re.sub(r"^[—–-]{2,}\s*", "", stripped))
        if not cleaned:
            continue
        if cleaned.lower() in {cover_title.lower(), full_title.lower()}:
            continue
        candidates.append(cleaned)

    deduped: list[str] = []
    seen: set[str] = set()
    for item in candidates:
        normalized = item.lower()
        if normalized in seen:
            continue
        seen.add(normalized)
        deduped.append(shorten(item, 28))
        if len(deduped) >= max_points:
            break
    return deduped


def extract_section_points(section: MarkdownSection, max_points: int = 4) -> list[str]:
    candidates: list[str] = []
    for raw_line in section.lines:
        stripped = raw_line.strip()
        if not stripped or stripped == "---" or stripped.startswith("```"):
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            line = clean_markdown_text(stripped.replace("|", " | "))
            if line:
                candidates.append(line)
            continue

        if re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+[.)]\s+", stripped):
            line = clean_markdown_text(stripped)
            if line:
                candidates.append(line)
            continue

        line = clean_markdown_text(stripped)
        if not line:
            continue
        parts = re.split(r"[。；;]", line)
        for part in parts:
            part = clean_markdown_text(part)
            if len(part) >= 8:
                candidates.append(part)

    deduped: list[str] = []
    for item in candidates:
        normalized = item.lower()
        if not item or normalized in {existing.lower() for existing in deduped}:
            continue
        deduped.append(item)
        if len(deduped) >= max_points:
            break

    return deduped


def related_sections(sections: list[MarkdownSection], anchor: MarkdownSection) -> list[MarkdownSection]:
    start_index = next((index for index, section in enumerate(sections) if section.index == anchor.index), None)
    if start_index is None:
        return [anchor]
    collected = [sections[start_index]]
    for section in sections[start_index + 1 :]:
        if section.level <= anchor.level:
            break
        collected.append(section)
    return collected


def heading_like_text(text: str) -> bool:
    cleaned = clean_markdown_text(text)
    return bool(
        re.match(r"^\d+(?:\.\d+)*\s+", cleaned)
        or re.match(r"^[一二三四五六七八九十]+[、.]\s*", cleaned)
        or re.match(r"^第[一二三四五六七八九十\d]+[章节部分步条项]", cleaned)
    )


def sentence_priority(text: str, *, is_bullet: bool = False, is_emphasis: bool = False, title_hint: str = "") -> int:
    score = 0
    cleaned = clean_markdown_text(text)
    if not cleaned:
        return -999
    if is_bullet:
        score += 4
    if is_emphasis:
        score += 6
    if 10 <= len(cleaned) <= 34:
        score += 3
    if re.search(r"\d|%|万|亿|年|月|日", cleaned):
        score += 2
    if any(re.search(pattern, cleaned) for pattern in STRONG_POINT_PATTERNS):
        score += 5
    if title_hint and any(re.search(pattern, title_hint) for pattern in STRONG_POINT_PATTERNS):
        score += 1
    if any(re.search(pattern, cleaned) for pattern in WEAK_POINT_PATTERNS):
        score -= 5
    if len(cleaned) > 42:
        score -= 2
    return score


def split_line_fragments(raw_line: str) -> list[str]:
    normalized = str(raw_line or "").replace("●", "\n●").replace("•", "\n•").replace("▪", "\n▪").replace("◆", "\n◆")
    fragments: list[str] = []
    for block in normalized.splitlines():
        stripped = block.strip()
        if not stripped or stripped == "---" or stripped.startswith("```"):
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            cleaned = clean_markdown_text(stripped.replace("|", " | "))
            if cleaned:
                fragments.append(cleaned)
            continue
        if re.match(r"^\s*[-*+●•▪◆◇■□]\s*", stripped) or re.match(r"^\s*\d+[.)]\s*", stripped):
            cleaned = clean_markdown_text(stripped)
            if cleaned:
                fragments.append(cleaned)
            continue
        line = clean_markdown_text(stripped)
        if not line:
            continue
        parts = re.split(r"[。；;]", line)
        for part in parts:
            cleaned = clean_markdown_text(part)
            if len(cleaned) >= 8:
                fragments.append(cleaned)
    return fragments


def extract_related_section_points(sections: list[MarkdownSection], anchor: MarkdownSection, max_points: int = 4) -> list[str]:
    ranked: list[tuple[int, int, str]] = []
    order = 0
    for section in related_sections(sections, anchor):
        if section.index != anchor.index and section.title:
            cleaned_title = clean_markdown_text(section.title)
            if len(cleaned_title) >= 6:
                ranked.append((sentence_priority(cleaned_title, is_emphasis=True, title_hint=anchor.title), order, cleaned_title))
                order += 1
        for raw_line in section.lines:
            stripped = raw_line.strip()
            if not stripped:
                continue
            bold_hits = re.findall(r"\*\*([^*]+)\*\*", raw_line)
            for hit in bold_hits:
                cleaned = clean_markdown_text(hit)
                if len(cleaned) >= 8:
                    ranked.append((sentence_priority(cleaned, is_emphasis=True, title_hint=anchor.title), order, cleaned))
                    order += 1
            is_bullet = bool(re.match(r"^\s*[-*+●•▪◆◇■□]\s*", stripped) or re.match(r"^\s*\d+[.)]\s*", stripped))
            for fragment in split_line_fragments(raw_line):
                ranked.append((sentence_priority(fragment, is_bullet=is_bullet, title_hint=anchor.title), order, fragment))
                order += 1

    deduped: list[str] = []
    seen: set[str] = set()
    for _, _, item in sorted(ranked, key=lambda row: (-row[0], row[1], len(row[2]))):
        normalized = item.lower()
        if normalized in seen:
            continue
        seen.add(normalized)
        deduped.append(item)
        if len(deduped) >= max_points:
            break

    if deduped:
        return deduped
    return extract_section_points(anchor, max_points=max_points)


def extract_focus_points(sections: list[MarkdownSection], anchor: MarkdownSection, max_points: int = 4) -> list[str]:
    def collect(from_sections: list[MarkdownSection]) -> list[str]:
        ranked: list[tuple[int, int, str]] = []
        order = 0
        for depth, section in enumerate(from_sections):
            section_bonus = max(0, 4 - depth)
            for raw_line in section.lines:
                stripped = raw_line.strip()
                if not stripped:
                    continue
                bold_hits = re.findall(r"\*\*([^*]+)\*\*", raw_line)
                for hit in bold_hits:
                    cleaned = clean_markdown_text(hit)
                    if len(cleaned) >= 8 and not heading_like_text(cleaned):
                        ranked.append((sentence_priority(cleaned, is_emphasis=True, title_hint=anchor.title) + section_bonus + 3, order, cleaned))
                        order += 1
                is_bullet = bool(re.match(r"^\s*[-*+●•▪◆◇■□]\s*", stripped) or re.match(r"^\s*\d+[.)]\s*", stripped))
                for fragment in split_line_fragments(raw_line):
                    if heading_like_text(fragment):
                        continue
                    ranked.append((sentence_priority(fragment, is_bullet=is_bullet, title_hint=anchor.title) + section_bonus, order, fragment))
                    order += 1

        deduped: list[str] = []
        seen: set[str] = set()
        for _, _, item in sorted(ranked, key=lambda row: (-row[0], row[1], len(row[2]))):
            normalized = item.lower()
            if normalized in seen:
                continue
            seen.add(normalized)
            deduped.append(item)
            if len(deduped) >= max_points:
                break
        return deduped

    anchor_only = collect([anchor])
    if len(anchor_only) >= min(2, max_points):
        return anchor_only[:max_points]

    ranked: list[tuple[int, int, str]] = []
    order = 0
    for depth, section in enumerate(related_sections(sections, anchor)):
        section_bonus = max(0, 4 - depth)
        for raw_line in section.lines:
            stripped = raw_line.strip()
            if not stripped:
                continue
            bold_hits = re.findall(r"\*\*([^*]+)\*\*", raw_line)
            for hit in bold_hits:
                cleaned = clean_markdown_text(hit)
                if len(cleaned) >= 8 and not heading_like_text(cleaned):
                    ranked.append((sentence_priority(cleaned, is_emphasis=True, title_hint=anchor.title) + section_bonus + 3, order, cleaned))
                    order += 1
            is_bullet = bool(re.match(r"^\s*[-*+●•▪◆◇■□]\s*", stripped) or re.match(r"^\s*\d+[.)]\s*", stripped))
            for fragment in split_line_fragments(raw_line):
                if heading_like_text(fragment):
                    continue
                ranked.append((sentence_priority(fragment, is_bullet=is_bullet, title_hint=anchor.title) + section_bonus, order, fragment))
                order += 1

    deduped: list[str] = []
    seen: set[str] = set()
    for _, _, item in sorted(ranked, key=lambda row: (-row[0], row[1], len(row[2]))):
        normalized = item.lower()
        if normalized in seen:
            continue
        seen.add(normalized)
        deduped.append(item)
        if len(deduped) >= max_points:
            break

    if deduped:
        return deduped
    return extract_related_section_points(sections, anchor, max_points=max_points)


def score_section(section: MarkdownSection, points: list[str]) -> int:
    score = len(points) * 6
    if section.level == 2:
        score += 8
    if re.search(r"(总结|结论|建议|展望|下一步|规划|策略|方案|路径|能力|优势|竞争力|洞察|分析|架构|系统|数据)", section.title):
        score += 10
    if re.search(r"\d|%|万|亿|年|月|日", section.title):
        score += 4
    if any(re.search(r"\d|%|万|亿|年|月|日", point) for point in points):
        score += 6
    if len(section.title) >= 6:
        score += 2
    return score


def build_global_points(sections: list[MarkdownSection]) -> list[str]:
    collected: list[str] = []
    for section in sections:
        points = extract_focus_points(sections, section, max_points=2)
        for point in points:
            short_point = shorten(point, 22)
            if short_point and short_point not in collected:
                collected.append(short_point)
            if len(collected) >= 2:
                return collected
    return collected or ["核心观点拆解", "关键动作落地"]


def build_cover_points(markdown: str, deck_title: str, story_sections: list[MarkdownSection], sections: list[MarkdownSection]) -> list[str]:
    lead_lines = extract_cover_lead_lines(markdown, deck_title, max_points=2)
    if lead_lines:
        return lead_lines

    section_titles: list[str] = []
    for section in story_sections or sections:
        cleaned = clean_markdown_text(section.title)
        if not cleaned:
            continue
        cleaned = re.sub(r"^[一二三四五六七八九十]+[、.]\s*", "", cleaned)
        cleaned = re.sub(r"^\d+(?:\.\d+)*\s*", "", cleaned)
        cleaned = re.sub(r"^第[一二三四五六七八九十\d]+[章节部分步条项]\s*", "", cleaned)
        cleaned = clean_markdown_text(cleaned)
        if cleaned:
            section_titles.append(shorten(cleaned, 24))
        if len(section_titles) >= 2:
            break

    if section_titles:
        return section_titles
    return build_global_points(story_sections or sections)


def preferred_sections(sections: list[MarkdownSection]) -> list[MarkdownSection]:
    level2 = [section for section in sections if section.level == 2]
    if level2:
        return level2
    level3 = [section for section in sections if section.level == 3]
    if level3:
        return level3
    return [section for section in sections if section.level >= 1]


def select_story_sections(sections: list[MarkdownSection], max_sections: int = 8) -> list[MarkdownSection]:
    candidates = [
        section
        for section in preferred_sections(sections)
        if section.title and section.title.lower() not in IGNORED_SECTION_TITLES
    ]
    if not candidates:
        return []

    intro_like = [section for section in candidates if re.search(r"(引言|概述|摘要|总览|背景)", section.title)]
    closing_like = [section for section in candidates if re.search(r"(结论|总结|建议|展望|下一步)", section.title)]
    excluded = {section.index for section in intro_like[:1]}
    middle_pool = [section for section in candidates if section.index not in excluded]

    if len(middle_pool) <= max_sections:
        return middle_pool

    reserved_ids: set[int] = set()
    if closing_like:
        reserved_ids.add(closing_like[0].index)

    slots = max_sections - len(reserved_ids)
    ranked = sorted(
        (
            (score_section(section, extract_related_section_points(sections, section, max_points=4)), section)
            for section in middle_pool
            if section.index not in reserved_ids
        ),
        key=lambda item: (-item[0], item[1].index),
    )
    picked_ids = {section.index for _, section in ranked[: max(0, slots)]}
    picked_ids.update(reserved_ids)
    selected = [section for section in middle_pool if section.index in picked_ids]
    return selected[:max_sections]


def visual_scene_prompt(title: str, deck_title: str) -> str:
    if title == deck_title:
        return GENERIC_SCENE_KEYWORDS["cover"]
    if re.search(r"(资本|路径|战略|规划)", title):
        return GENERIC_SCENE_KEYWORDS["capital"]
    if re.search(r"(企业集群|行业表象|权力接管|区域|产业)", title):
        return GENERIC_SCENE_KEYWORDS["industry_control"]
    if re.search(r"(AI应用体系|平台|系统|ClawLink)", title):
        return GENERIC_SCENE_KEYWORDS["ai_system"]
    if re.search(r"(盈利|收益|分配|模型)", title):
        return GENERIC_SCENE_KEYWORDS["profit_engine"]
    return GENERIC_SCENE_KEYWORDS["general"]


def build_pages_from_markdown(markdown: str, source_path: Path) -> tuple[str, list[SlidePage]]:
    deck_title, sections = parse_markdown_sections(markdown, source_path)
    story_sections = select_story_sections(sections, max_sections=8)
    pages: list[SlidePage] = []

    cover_points = build_cover_points(markdown, deck_title, story_sections, sections)
    cover_title = shorten(deck_title.split("：", 1)[0], 18)
    pages.append(
        SlidePage(
            page=1,
            title=cover_title or "漫画总览",
            bullets=cover_points[:2],
            context_title=deck_title,
            context_points=cover_points[:4],
            scene_prompt=visual_scene_prompt(deck_title, deck_title),
            filename=f"P1-{sanitize_filename_part(cover_title or deck_title)}.png",
        )
    )

    for index, section in enumerate(story_sections, start=2):
        context_points = extract_focus_points(sections, section, max_points=4) or [section.title]
        render_title = section.title
        render_bullets = context_points[:2]
        pages.append(
            SlidePage(
                page=index,
                title=render_title,
                bullets=render_bullets,
                context_title=section.title,
                context_points=context_points,
                scene_prompt=visual_scene_prompt(section.title, deck_title),
                filename=f"P{index}-{sanitize_filename_part(shorten(render_title, 18))}.png",
            )
        )

    return deck_title, pages


def build_analysis_summary(markdown: str, source_path: Path, deck_title: str, pages: list[SlidePage]) -> str:
    _, sections = parse_markdown_sections(markdown, source_path)
    section_titles = [section.title for section in sections[:6] if section.title]
    selected_titles = [page.context_title for page in pages[1:5]]
    return (
        f"识别文档标题《{deck_title}》，共解析 {len(sections)} 个章节；"
        f"本次选取 {len(pages)} 页内容进行漫画化表达。"
        + (f" 重点章节：{' / '.join(section_titles)}。" if section_titles else "")
        + (f" 当前入选页面：{' / '.join(selected_titles)}。" if selected_titles else "")
    )


def build_design_summary(pages: list[SlidePage]) -> str:
    slide_titles = " / ".join(page.title for page in pages[:4])
    return (
        "设计上采用商业漫画信息图路线：白底留白、蓝橙强调色、粗描边、卡片化信息区，"
        "并让标题与要点直接融合进插画。"
        + (f" 首页到前几页的叙事顺序为：{slide_titles}。" if slide_titles else "")
    )


def page_architecture_rules(page: SlidePage) -> str:
    if page.page == 1:
        return textwrap.dedent(
            """
            - cover layout: dark blue editorial title plaque in the upper-left
            - left or lower-left: one collapsing old structure or old-order symbol
            - center-right: one awakened core or circular engine as the main visual anchor
            - multiple thick blue/orange arrows must connect old structure to new core and surrounding modules
            - right side: dense supporting business modules, terminals, dashboards or symbolic devices around the core
            """
        ).strip()
    if page.page == 2:
        return textwrap.dedent(
            """
            - upper-left: large title-and-bullets panel
            - lower-left: one compact mini-diagram showing cause/problem/source structure
            - center: one transformation hub with gears, modules or bridge blocks
            - right half: one staged roadmap with multi-step business path and thick directional arrows
            - overall feel: left-to-right strategic progression board with dense panels and milestone rhythm
            """
        ).strip()
    if page.page == 3:
        return textwrap.dedent(
            """
            - upper-left: compact title panel
            - left side: one before-state cluster or old-boundary motif
            - center: one circular system loop or command-center ring as the dominant structure
            - one strong horizontal arrow should push from left toward the central system loop
            - right side: awakening or AI-output symbols, capability icons and outward signal marks
            """
        ).strip()
    if page.page == 4:
        return textwrap.dedent(
            """
            - top-left: long horizontal title panel with bullets
            - upper-right: one compact profitability or control wheel/dashboard cluster
            - middle and bottom: wide ecosystem blueprint with many linked modules, stations, pipelines and containers
            - use a layered industrial or service-chain feel with left-to-right operational flow
            - keep the page dense and panoramic rather than centered like a poster
            """
        ).strip()
    if page.page == 5:
        return textwrap.dedent(
            """
            - upper-left: title panel
            - left-to-center: one core control path with control box, money or distribution box and strong connecting arrow
            - right half: stacked blocks for revenue structure, system engine and scaling path
            - include mini grids, charts, money icons, replication or scaling modules and execution nodes
            - overall feel: modular operating-model board with several boxed subsystems arranged in a clean grid
            """
        ).strip()
    return textwrap.dedent(
        """
        - use a management-infographic page architecture with one title card, one main system hub and several supporting cards connected by arrows
        - keep the composition diagrammatic and dense, not poster-like
        """
    ).strip()


def build_prompt(
    page: SlidePage,
    native_text: bool = False,
    *,
    image_provider: str = "gemini",
    has_style_reference: bool = False,
    has_layout_reference: bool = False,
    has_page_reference: bool = False,
) -> str:
    provider = normalize_image_provider(image_provider)
    reference_intro_lines: list[str] = []
    if has_style_reference:
        if has_layout_reference or has_page_reference:
            reference_intro_lines.append("The first reference image defines the target visual language.")
        else:
            reference_intro_lines.append(
                "If a style reference image is provided, use it only for overall visual language such as palette, line weight, panel feel and cleanliness."
            )
            reference_intro_lines.append("Do not copy its composition, page structure, text landing points or specific content.")
    if has_layout_reference:
        reference_intro_lines.append("If a layout reference image is provided, use it for composition rhythm, text integration method and page density.")
    if has_page_reference:
        reference_intro_lines.append("If a page-specific reference image is provided, use it as the exact layout, text-placement and image-text-fusion reference for the same page.")
    if not reference_intro_lines:
        reference_intro_lines.append("Rely entirely on the prompt below to invent a fresh composition.")
        reference_intro_lines.append("Do not imitate any fixed reference page, previous deck layout or template screenshot.")
    reference_intro = "\n".join(reference_intro_lines)
    page_reference_rules = ""
    if has_page_reference:
        page_reference_rules = textwrap.dedent(
            """
            Page-specific reference rules:
            - treat the last input reference image as a strict composition blueprint, not as content to copy
            - preserve the reference page's macro layout: title/card location, flow direction, diagram density, arrow paths, dashboard/card clusters, corner decorations and white-space rhythm
            - preserve the reference page's text landing points: where the title block sits, how large it is, how bullet lines sit under it, and how the text block connects to surrounding graphics
            - preserve the image-text fusion method: text must live inside designed blue/white/orange panels that are visually connected to the business diagram, never as a flat overlay
            - replace every readable word from the reference page with only the exact lines requested below; never copy reference wording, old labels, page numbers or logos
            - keep small UI/detail areas as abstract bars, dots, icons or unreadable micro-marks unless their text is one of the exact requested lines
            """
        ).strip()
    provider_rules = ""
    if provider == "openai":
        provider_rules = (
            textwrap.dedent(
                """
                GPT image route priorities:
                - follow the reference layout more strongly than generic poster composition
                - keep the old fusion-version feel: dense management infographic, structured panels, connected arrows, mini dashboards, business icons and blueprint-like wiring
                - render Chinese with simple bold sans-serif strokes; prioritize exact readable characters over decorative distortion
                - if text accuracy conflicts with visual density, simplify nearby decoration instead of inventing or corrupting text
                - avoid a clean local-overlay look; make the text blocks feel natively designed into the illustration
                """
            ).strip()
            if has_layout_reference or has_page_reference
            else textwrap.dedent(
                """
                GPT image route priorities:
                - create a fresh dense management-infographic composition from the content itself, not a generic poster
                - keep the old fusion-version feel: structured panels, connected arrows, mini dashboards, business icons and blueprint-like wiring
                - render Chinese with simple bold sans-serif strokes; prioritize exact readable characters over decorative distortion
                - if text accuracy conflicts with visual density, simplify nearby decoration instead of inventing or corrupting text
                - make the text blocks feel natively designed into the illustration, not pasted on later
                """
            ).strip()
        )
    architecture_rules = page_architecture_rules(page)
    if native_text:
        if page.page == 1:
            exact_lines = "\n".join([page.title, *page.bullets[:2]])
            text_layout_rules = textwrap.dedent(
                """
                - render the first line as the main cover title
                - render the remaining lines as subtitle or tagline lines, without bullet dots unless they are already present
                - place the text in one integrated editorial block near the upper-left or upper-middle area
                - the text zone must feel like part of the illustration, such as a designed information card or embedded panel
                - keep each line clearly separated, balanced and easy to read
                """
            ).strip()
        else:
            exact_lines = "\n".join([page.title, *[f"• {point}" for point in page.bullets[:2]]])
            text_layout_rules = textwrap.dedent(
                """
                - place all required text in one natural editorial zone near the upper-left or upper-middle area
                - render the title larger and the two bullet lines smaller below it
                - keep the text block integrated with the illustration like a designed information panel
                - keep each bullet on its own line with clear spacing
                - keep the text short, balanced and visually fused with the illustration
                """
            ).strip()
        return textwrap.dedent(
            f"""
            Create one finished 16:9 slide illustration for a Chinese business comic PowerPoint.
            {reference_intro}

            {page_reference_rules}

            {provider_rules}

            Visual style requirements:
            - premium editorial comic infographic, not childish anime, not photorealistic, not 3D
            - white background, large clean negative space, crisp black outlines
            - deep blue accents, orange accents, light blue and light orange blocks
            - business storyboard feeling, designed for management presentation instead of entertainment poster
            - text must feel embedded inside the design, not pasted on top later
            - prefer diagrams, cards, dashboards, arrows and abstract symbols over detailed people
            - avoid cute mascots and avoid malformed human figures

            Page architecture rules:
            {architecture_rules}

            Main scene to generate:
            {page.scene_prompt}

            You must render these exact lines directly inside the image, with correct wording and clean line separation:
            {exact_lines}

            Text layout rules:
            {text_layout_rules}

            Hard constraints:
            - do not add any extra text beyond the exact lines above
            - no wrong characters, no garbled Chinese, no fake text, no pseudo text
            - no extra English or numbers beyond the exact lines above
            - no watermark, no logo, no page number
            - no giant disconnected title banner across the whole page
            - no creepy faces, no extra hands or legs
            """
        ).strip()
    page_reference_rules_text_free = ""
    if has_page_reference:
        page_reference_rules_text_free = textwrap.dedent(
            """
            Page-specific reference rules:
            - treat the last input reference image as a strict composition blueprint, not as content to copy
            - preserve its macro layout, diagram density, arrow paths, dashboard/card clusters, corner decorations and white-space rhythm
            - reserve the same title/card landing area as clean empty space for later editable text
            - every title area, label area and diagram-caption area from the reference must become empty panels, abstract icon chips, bars, dots or unreadable micro-marks
            - do not copy any readable wording, labels, page numbers or logos from the reference image
            """
        ).strip()
    return textwrap.dedent(
        f"""
        Create one finished 16:9 text-free slide illustration for a Chinese business comic PowerPoint.
        {reference_intro}

        {page_reference_rules_text_free}

        {provider_rules}

        Visual style requirements:
        - premium editorial comic infographic, not childish anime, not photorealistic, not 3D
        - white background, large clean negative space, crisp black outlines
        - deep blue accents, orange accents, light blue and light orange blocks
        - business storyboard feeling, designed for management presentation instead of entertainment poster
        - prefer diagrams, cards, dashboards, arrows and abstract symbols over detailed people
        - avoid cute mascots and avoid malformed human figures
        - reserve a clean empty typography zone in the upper-left area, about 45% width and 34% height

        Page architecture rules:
        {architecture_rules}

        Narrative context:
        - page focus: {page.context_title}
        - the page should communicate the section meaning visually, without rendering any readable words

        Main scene to generate:
        {page.scene_prompt}

        Hard constraints:
        - absolutely no readable text, no Chinese, no English, no numbers, no labels, no captions, no axis words
        - no fake text, no pseudo text, no watermark, no logo
        - convert any would-be label area into blank chips, blocks, bars, ticks or abstract icon placeholders
        - keep the upper-left typography zone clean and mostly empty for external text overlay
        - no creepy faces, no extra hands or legs
        """
    ).strip()


def summarize_http_error(error_body: str) -> str:
    try:
        payload = json.loads(error_body)
        message = payload.get("error", {}).get("message")
        status = payload.get("error", {}).get("status")
        if message and status:
            return f"{status}: {message}"
        if message:
            return message
    except Exception:
        pass
    return error_body[:480]


def extract_gemini_error_details(response_payload: dict[str, object]) -> tuple[str, str, int | None]:
    error = response_payload.get("error")
    if not isinstance(error, dict):
        return "", "", None
    status = str(error.get("status") or "").strip()
    message = str(error.get("message") or "").strip()
    code = error.get("code")
    return status, message, code if isinstance(code, int) else None


def should_retry_gemini_error(status: str, message: str, code: int | None) -> bool:
    normalized_status = status.upper()
    normalized_message = message.lower()
    if normalized_status in {"UNAVAILABLE", "DEADLINE_EXCEEDED", "INTERNAL", "ABORTED", "UNKNOWN"}:
        return True
    if code in {429, 500, 502, 503, 504}:
        if "quota" in normalized_message or "billing" in normalized_message:
            return False
        return True
    return any(
        marker in normalized_message
        for marker in [
            "high demand",
            "try again later",
            "temporarily unavailable",
            "deadline exceeded",
            "timed out",
            "timeout",
            "busy",
            "overloaded",
            "overload",
            "高需求",
            "繁忙",
            "过载",
            "稍后再试",
        ]
    )


def should_retry_openai_compatible_error(message: str, status_code: int | None = None) -> bool:
    normalized_message = message.lower()
    if status_code in {408, 409, 429, 500, 502, 503, 504}:
        if "quota" in normalized_message or "billing" in normalized_message:
            return False
        return True
    return any(
        marker in normalized_message
        for marker in [
            "high demand",
            "try again later",
            "temporarily unavailable",
            "deadline exceeded",
            "timed out",
            "timeout",
            "rate limit",
            "server error",
            "empty reply from server",
            "ssl_error_syscall",
            "ssl_connect",
            "connection reset",
            "connection refused",
            "connection aborted",
            "bad gateway",
            "gateway timeout",
            "upstream",
            "overloaded",
            "overload",
            "busy",
            "负载",
            "饱和",
            "上游",
            "繁忙",
            "过载",
            "稍后再试",
        ]
    )


def compute_retry_delay_seconds(attempt: int, base_delay_seconds: float, max_delay_seconds: float, jitter_seconds: float) -> float:
    delay_seconds = min(max_delay_seconds, base_delay_seconds * (2 ** max(attempt - 1, 0)))
    if jitter_seconds > 0:
        delay_seconds += random.uniform(0, jitter_seconds)
    return delay_seconds


def request_gemini_image_once(url: str, api_key: str, payload_json: bytes, output_path: Path) -> str:
    request_timeout_seconds = read_positive_int_env("GEMINI_IMAGE_REQUEST_TIMEOUT_SECONDS", IMAGE_REQUEST_TIMEOUT_SECONDS)
    connect_timeout_seconds = read_positive_int_env("GEMINI_IMAGE_CONNECT_TIMEOUT_SECONDS", DEFAULT_IMAGE_CONNECT_TIMEOUT_SECONDS)
    low_speed_time_seconds = read_positive_int_env("GEMINI_IMAGE_LOW_SPEED_TIME_SECONDS", DEFAULT_IMAGE_LOW_SPEED_TIME_SECONDS)
    low_speed_limit_bytes = read_positive_int_env("GEMINI_IMAGE_LOW_SPEED_LIMIT_BYTES", DEFAULT_IMAGE_LOW_SPEED_LIMIT_BYTES)
    try:
        response = subprocess.run(
            [
                "curl",
                "-sS",
                "--http1.1",
                "--connect-timeout",
                str(connect_timeout_seconds),
                "--max-time",
                str(request_timeout_seconds),
                "--speed-time",
                str(low_speed_time_seconds),
                "--speed-limit",
                str(low_speed_limit_bytes),
                "-X",
                "POST",
                url,
                "-H",
                f"x-goog-api-key: {api_key}",
                "-H",
                "Content-Type: application/json",
                "--data-binary",
                "@-",
            ],
            input=payload_json,
            capture_output=True,
            timeout=request_timeout_seconds + 10,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise GeminiGenerationError(
            f"Gemini image request timed out after {request_timeout_seconds} seconds.",
            retryable=True,
        ) from exc

    if response.returncode != 0:
        raise GeminiGenerationError(
            response.stderr.decode("utf-8", errors="replace") or f"curl exited {response.returncode}",
            retryable=True,
        )

    try:
        response_payload = json.loads(response.stdout.decode("utf-8", errors="replace"))
    except json.JSONDecodeError as exc:
        raise GeminiGenerationError(response.stdout.decode("utf-8", errors="replace")[:480]) from exc

    if response_payload.get("error"):
        error_message = summarize_http_error(json.dumps(response_payload, ensure_ascii=False))
        error_status, raw_message, error_code = extract_gemini_error_details(response_payload)
        raise GeminiGenerationError(
            error_message,
            retryable=should_retry_gemini_error(error_status, raw_message, error_code),
        )

    parts = response_payload.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    for part in parts:
        inline_data = part.get("inlineData") or part.get("inline_data")
        if inline_data and inline_data.get("data"):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(base64.b64decode(inline_data["data"]))
            return "gemini-api"

    text_parts = [part.get("text", "") for part in parts if part.get("text")]
    raise GeminiGenerationError("；".join(text_parts) if text_parts else "Gemini response did not contain inline image data.")


def reference_paths(
    *,
    style_reference: Path | None = None,
    layout_reference: Path | None = None,
    page_reference: Path | None = None,
) -> list[Path]:
    references: list[Path] = []
    for candidate in [style_reference, layout_reference, page_reference]:
        if candidate and candidate.exists() and candidate not in references:
            references.append(candidate)
    return references


def parse_ip_address(value: str | None) -> ipaddress._BaseAddress | None:
    try:
        return ipaddress.ip_address(str(value or "").strip())
    except ValueError:
        return None


def is_fake_dns_ip(value: str | None) -> bool:
    address = parse_ip_address(value)
    return bool(address and address in ipaddress.ip_network("198.18.0.0/15"))


def curl_resolve_args_for_url(url: str) -> list[str]:
    if str(os.environ.get("COMIC_OPENAI_DISABLE_DNS_BYPASS", "")).strip() == "1":
        return []
    parsed = urllib.parse.urlparse(str(url or "").strip())
    host = parsed.hostname
    if not host or parse_ip_address(host):
        return []
    port = parsed.port or (443 if parsed.scheme == "https" else 80)
    address = resolve_host_for_curl(host)
    if not address:
        return []
    return ["--resolve", f"{host}:{port}:{address}"]


def resolve_host_for_curl(host: str) -> str:
    public_address = resolve_host_via_public_dns(host)
    if public_address and not is_fake_dns_ip(public_address):
        return public_address
    try:
        resolved = socket.gethostbyname(host)
    except OSError:
        return public_address
    if resolved and not is_fake_dns_ip(resolved):
        return resolved
    return public_address or ""


def resolve_host_via_public_dns(host: str) -> str:
    resolvers = read_csv_env("COMIC_OPENAI_DNS_RESOLVERS", DEFAULT_OPENAI_COMPATIBLE_DNS_RESOLVERS)
    for resolver in resolvers:
        try:
            response = subprocess.run(
                [
                    "dig",
                    "+short",
                    "+time=2",
                    "+tries=1",
                    f"@{resolver}",
                    host,
                    "A",
                ],
                capture_output=True,
                timeout=5,
                check=False,
            )
        except (subprocess.TimeoutExpired, FileNotFoundError):
            continue
        if response.returncode != 0:
            continue
        for line in response.stdout.decode("utf-8", errors="replace").splitlines():
            candidate = line.strip()
            address = parse_ip_address(candidate)
            if address and address.version == 4:
                return candidate
    return ""


def decode_base64_image_data(value: object) -> bytes:
    data = str(value or "").strip()
    if data.startswith("data:") and "," in data:
        data = data.split(",", 1)[1]
    data = re.sub(r"\s+", "", data)
    if not data:
        raise ValueError("empty b64_json")
    remainder = len(data) % 4
    if remainder == 1:
        raise ValueError(f"invalid base64 length {len(data)}")
    padded = data + ("=" * ((4 - remainder) % 4))
    try:
        return base64.b64decode(padded, altchars=b"-_", validate=True)
    except (binascii.Error, ValueError) as exc:
        raise ValueError(str(exc)) from exc


def request_openai_compatible_image_once(
    base_url: str,
    api_key: str,
    model: str,
    prompt: str,
    references: list[Path],
    output_path: Path,
) -> str:
    request_timeout_seconds = read_positive_int_env("COMIC_OPENAI_IMAGE_REQUEST_TIMEOUT_SECONDS", OPENAI_COMPATIBLE_IMAGE_REQUEST_TIMEOUT_SECONDS)
    connect_timeout_seconds = read_positive_int_env("COMIC_OPENAI_IMAGE_CONNECT_TIMEOUT_SECONDS", DEFAULT_IMAGE_CONNECT_TIMEOUT_SECONDS)
    low_speed_time_seconds = read_positive_int_env("COMIC_OPENAI_IMAGE_LOW_SPEED_TIME_SECONDS", DEFAULT_OPENAI_COMPATIBLE_LOW_SPEED_TIME_SECONDS)
    low_speed_limit_bytes = read_positive_int_env("COMIC_OPENAI_IMAGE_LOW_SPEED_LIMIT_BYTES", DEFAULT_IMAGE_LOW_SPEED_LIMIT_BYTES)
    size = str(os.environ.get("COMIC_OPENAI_IMAGE_SIZE") or os.environ.get("COMIC_GPT_IMAGE_SIZE") or "auto").strip()
    image_field = str(os.environ.get("COMIC_OPENAI_IMAGE_FIELD") or "image[]").strip() or "image[]"
    generation_response_format = str(
        os.environ.get("COMIC_OPENAI_IMAGE_GENERATION_RESPONSE_FORMAT")
        or os.environ.get("COMIC_OPENAI_IMAGE_RESPONSE_FORMAT")
        or DEFAULT_OPENAI_COMPATIBLE_GENERATION_RESPONSE_FORMAT
    ).strip() or DEFAULT_OPENAI_COMPATIBLE_GENERATION_RESPONSE_FORMAT
    command = [
        "curl",
        "-sS",
        "--http1.1",
        "--connect-timeout",
        str(connect_timeout_seconds),
        "--max-time",
        str(request_timeout_seconds),
        "--speed-time",
        str(low_speed_time_seconds),
        "--speed-limit",
        str(low_speed_limit_bytes),
        "-X",
        "POST",
    ]

    request_url = f"{base_url}/images/edits" if references else f"{base_url}/images/generations"
    request_resolve_args = curl_resolve_args_for_url(request_url)

    if references:
        command.extend(
            [
                *request_resolve_args,
                request_url,
                "-H",
                f"Authorization: Bearer {api_key}",
                "-F",
                f"model={model}",
                "-F",
                f"prompt={prompt}",
            ]
        )
        if size:
            command.extend(["-F", f"size={size}"])
        for reference in references:
            command.extend(["-F", f"{image_field}=@{reference};type={mime_type_for_path(reference)}"])
    else:
        payload: dict[str, object] = {
            "model": model,
            "prompt": prompt,
            "response_format": generation_response_format,
        }
        if size:
            payload["size"] = size
        command.extend(
            [
                *request_resolve_args,
                request_url,
                "-H",
                f"Authorization: Bearer {api_key}",
                "-H",
                "Content-Type: application/json",
                "--data-binary",
                json.dumps(payload, ensure_ascii=False),
            ]
        )

    try:
        response = subprocess.run(
            command,
            capture_output=True,
            timeout=request_timeout_seconds + 10,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise OpenAICompatibleImageError(
            f"OpenAI-compatible image request timed out after {request_timeout_seconds} seconds.",
            retryable=True,
        ) from exc

    if response.returncode != 0:
        raise OpenAICompatibleImageError(
            response.stderr.decode("utf-8", errors="replace") or f"curl exited {response.returncode}",
            retryable=True,
        )

    decoded = response.stdout.decode("utf-8", errors="replace")
    try:
        response_payload = json.loads(decoded)
    except json.JSONDecodeError as exc:
        raise OpenAICompatibleImageError(decoded[:480]) from exc

    if response_payload.get("error"):
        error = response_payload.get("error") if isinstance(response_payload.get("error"), dict) else {}
        message = str(error.get("message") or response_payload.get("error") or "").strip()
        status = error.get("code")
        raise OpenAICompatibleImageError(
            summarize_http_error(json.dumps(response_payload, ensure_ascii=False)),
            retryable=should_retry_openai_compatible_error(message, status if isinstance(status, int) else None),
        )

    data_items = response_payload.get("data")
    if isinstance(data_items, list) and data_items:
        first_item = data_items[0]
        if isinstance(first_item, dict):
            if first_item.get("b64_json"):
                try:
                    image_bytes = decode_base64_image_data(first_item["b64_json"])
                except ValueError as exc:
                    raise OpenAICompatibleImageError(
                        f"OpenAI-compatible response contained invalid b64_json image data: {exc}",
                        retryable=True,
                    ) from exc
                output_path.parent.mkdir(parents=True, exist_ok=True)
                output_path.write_bytes(image_bytes)
                return "openai-compatible-api"
            if first_item.get("url"):
                output_path.parent.mkdir(parents=True, exist_ok=True)
                download = subprocess.run(
                    [
                        "curl",
                        "-sS",
                        "--location",
                        *curl_resolve_args_for_url(str(first_item["url"])),
                        "--connect-timeout",
                        str(connect_timeout_seconds),
                        "--max-time",
                        str(request_timeout_seconds),
                        str(first_item["url"]),
                        "-o",
                        str(output_path),
                    ],
                    capture_output=True,
                    timeout=request_timeout_seconds + 10,
                    check=False,
                )
                if download.returncode == 0 and output_path.exists() and output_path.stat().st_size > 0:
                    return "openai-compatible-api"
                raise OpenAICompatibleImageError(
                    download.stderr.decode("utf-8", errors="replace") or "Failed to download generated image URL.",
                    retryable=True,
                )

    raise OpenAICompatibleImageError("OpenAI-compatible response did not contain image data.")


def generate_gemini_image(
    page: SlidePage,
    output_path: Path,
    native_text: bool = False,
    *,
    style_reference: Path | None = None,
    layout_reference: Path | None = None,
    page_reference: Path | None = None,
) -> str:
    api_key = get_api_key()
    model = urllib.parse.quote(os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL), safe="")
    base_url = os.environ.get("GEMINI_IMAGE_BASE_URL", "https://generativelanguage.googleapis.com/v1beta").rstrip("/")
    url = f"{base_url}/models/{model}:generateContent"
    max_attempts = read_positive_int_env("GEMINI_IMAGE_MAX_ATTEMPTS", DEFAULT_GEMINI_IMAGE_MAX_ATTEMPTS)
    base_delay_seconds = read_non_negative_float_env("GEMINI_IMAGE_RETRY_BASE_DELAY_SECONDS", DEFAULT_GEMINI_RETRY_BASE_DELAY_SECONDS)
    max_delay_seconds = read_non_negative_float_env("GEMINI_IMAGE_RETRY_MAX_DELAY_SECONDS", DEFAULT_GEMINI_RETRY_MAX_DELAY_SECONDS)
    jitter_seconds = read_non_negative_float_env("GEMINI_IMAGE_RETRY_JITTER_SECONDS", DEFAULT_GEMINI_RETRY_JITTER_SECONDS)

    references = reference_paths(
        style_reference=style_reference,
        layout_reference=layout_reference,
        page_reference=page_reference,
    )
    has_style_reference = style_reference is not None and style_reference.exists()
    has_layout_reference = layout_reference is not None and layout_reference.exists()
    has_page_reference = page_reference is not None and page_reference.exists()

    parts: list[dict[str, object]] = [
        {
            "text": build_prompt(
                page,
                native_text=native_text,
                image_provider="gemini",
                has_style_reference=has_style_reference,
                has_layout_reference=has_layout_reference,
                has_page_reference=has_page_reference,
            )
        }
    ]
    for reference in references:
        parts.append(
            {
                "inline_data": {
                    "mime_type": mime_type_for_path(reference),
                    "data": base64.b64encode(reference.read_bytes()).decode("utf-8"),
                }
            }
        )

    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
            "imageConfig": {"aspectRatio": "16:9", "imageSize": "2K"},
        },
    }
    payload_json = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    last_error: GeminiGenerationError | None = None

    for attempt in range(1, max_attempts + 1):
        try:
            if attempt > 1:
                emit_log(
                    "retry",
                    "Gemini 生图重试请求已发起",
                    f"正在发起第 {attempt}/{max_attempts} 次请求；如果上游处理较慢，下一条日志可能需要等待数分钟。",
                    attempt=attempt,
                    maxAttempts=max_attempts,
                    outputPath=str(output_path),
                )
            return request_gemini_image_once(url, api_key, payload_json, output_path)
        except GeminiGenerationError as exc:
            last_error = exc
            if not exc.retryable or attempt >= max_attempts:
                raise RuntimeError(str(exc)) from exc
            delay_seconds = compute_retry_delay_seconds(attempt, base_delay_seconds, max_delay_seconds, jitter_seconds)
            emit_log(
                "retry",
                "Gemini 生图重试",
                f"第 {attempt}/{max_attempts} 次请求失败：{exc}；将在 {delay_seconds:.1f} 秒后自动重试。",
                attempt=attempt,
                maxAttempts=max_attempts,
                retryInSeconds=round(delay_seconds, 2),
                outputPath=str(output_path),
            )
            time.sleep(delay_seconds)

    if last_error is not None:
        raise RuntimeError(str(last_error)) from last_error
    raise RuntimeError("Gemini image request failed before any attempt was made.")


def generate_openai_compatible_image(
    page: SlidePage,
    output_path: Path,
    native_text: bool = False,
    *,
    style_reference: Path | None = None,
    layout_reference: Path | None = None,
    page_reference: Path | None = None,
) -> str:
    api_key = get_openai_compatible_api_key()
    model = os.environ.get("COMIC_OPENAI_IMAGE_MODEL") or os.environ.get("COMIC_GPT_IMAGE_MODEL") or DEFAULT_OPENAI_IMAGE_MODEL
    base_url = openai_compatible_base_url()
    max_attempts = read_positive_int_env("COMIC_OPENAI_IMAGE_MAX_ATTEMPTS", DEFAULT_OPENAI_IMAGE_MAX_ATTEMPTS)
    base_delay_seconds = read_non_negative_float_env("COMIC_OPENAI_IMAGE_RETRY_BASE_DELAY_SECONDS", DEFAULT_GEMINI_RETRY_BASE_DELAY_SECONDS)
    max_delay_seconds = read_non_negative_float_env("COMIC_OPENAI_IMAGE_RETRY_MAX_DELAY_SECONDS", DEFAULT_GEMINI_RETRY_MAX_DELAY_SECONDS)
    jitter_seconds = read_non_negative_float_env("COMIC_OPENAI_IMAGE_RETRY_JITTER_SECONDS", DEFAULT_GEMINI_RETRY_JITTER_SECONDS)
    references = reference_paths(
        style_reference=style_reference,
        layout_reference=layout_reference,
        page_reference=page_reference,
    )
    has_style_reference = style_reference is not None and style_reference.exists()
    has_layout_reference = layout_reference is not None and layout_reference.exists()
    has_page_reference = page_reference is not None and page_reference.exists()
    prompt = build_prompt(
        page,
        native_text=native_text,
        image_provider="openai",
        has_style_reference=has_style_reference,
        has_layout_reference=has_layout_reference,
        has_page_reference=has_page_reference,
    )
    last_error: OpenAICompatibleImageError | None = None

    for attempt in range(1, max_attempts + 1):
        try:
            if attempt > 1:
                emit_log(
                    "retry",
                    "GPT 生图重试请求已发起",
                    f"正在发起第 {attempt}/{max_attempts} 次请求；如果上游处理较慢，下一条日志可能需要等待数分钟。",
                    attempt=attempt,
                    maxAttempts=max_attempts,
                    outputPath=str(output_path),
                )
            return request_openai_compatible_image_once(base_url, api_key, model, prompt, references, output_path)
        except OpenAICompatibleImageError as exc:
            last_error = exc
            if not exc.retryable or attempt >= max_attempts:
                raise RuntimeError(str(exc)) from exc
            delay_seconds = compute_retry_delay_seconds(attempt, base_delay_seconds, max_delay_seconds, jitter_seconds)
            emit_log(
                "retry",
                "GPT 生图重试",
                f"第 {attempt}/{max_attempts} 次请求失败：{exc}；将在 {delay_seconds:.1f} 秒后自动重试。",
                attempt=attempt,
                maxAttempts=max_attempts,
                retryInSeconds=round(delay_seconds, 2),
                outputPath=str(output_path),
            )
            time.sleep(delay_seconds)

    if last_error is not None:
        raise RuntimeError(str(last_error)) from last_error
    raise RuntimeError("OpenAI-compatible image request failed before any attempt was made.")


def text_block_height(draw: ImageDraw.ImageDraw, text: str, font: object, spacing: int) -> int:
    box = draw.multiline_textbbox((0, 0), text, font=font, spacing=spacing)
    return int(box[3] - box[1])


def overlay_card_geometry(page: SlidePage) -> tuple[int, int, int]:
    presets = {
        1: (34, 24, 900),
        2: (28, 20, 1140),
        3: (34, 26, 860),
        4: (26, 22, 1220),
        5: (30, 26, 920),
    }
    if page.page in presets:
        return presets[page.page]
    return (66, 46, 900 if len(page.title) > 22 else 820)


def render_native_text_overlay(image: Image.Image, page: SlidePage) -> Image.Image:
    base = image.convert("RGBA").resize((WIDTH, HEIGHT))
    overlay = Image.new("RGBA", (WIDTH, HEIGHT), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    card_left, card_top, card_width = overlay_card_geometry(page)
    title_font = OVERLAY_TITLE_SMALL_FONT if len(page.title) > 22 else OVERLAY_TITLE_FONT
    bullet_font = OVERLAY_BODY_SMALL_FONT if any(len(bullet) > 32 for bullet in page.bullets[:2]) else OVERLAY_BODY_FONT
    title_spacing = 10
    bullet_spacing = 8

    title_text = wrap_text(page.title, title_font, card_width - 104)
    bullet_texts = [wrap_text(bullet, bullet_font, card_width - 160) for bullet in page.bullets[:2]]

    title_height = text_block_height(draw, title_text, title_font, title_spacing)
    bullet_heights = [text_block_height(draw, bullet_text, bullet_font, bullet_spacing) for bullet_text in bullet_texts]
    title_panel_height = max(96, title_height + 42)
    body_height = sum(bullet_heights) + max(0, len(bullet_heights) - 1) * 24 + 54
    card_height = min(540, 28 + title_panel_height + body_height + 24)

    cover_box = (card_left - 18, card_top - 14, card_left + card_width + 18, card_top + card_height + 16)
    shadow_box = (card_left + 10, card_top + 12, card_left + card_width + 10, card_top + card_height + 12)
    card_box = (card_left, card_top, card_left + card_width, card_top + card_height)
    draw.rounded_rectangle(cover_box, radius=36, fill=(255, 255, 255, 250))
    draw.rounded_rectangle(shadow_box, radius=32, fill=(15, 23, 42, 34))
    draw.rounded_rectangle(card_box, radius=32, fill=(255, 255, 255, 244), outline=(15, 45, 75, 230), width=5)

    title_box = (card_left + 24, card_top + 24, card_left + card_width - 24, card_top + 24 + title_panel_height)
    draw.rounded_rectangle(title_box, radius=22, fill=(220, 238, 254, 236), outline=(31, 111, 178, 210), width=3)
    draw.rounded_rectangle((title_box[0], title_box[1], title_box[0] + 14, title_box[3]), radius=7, fill=(242, 154, 46, 255))
    draw.multiline_text((title_box[0] + 36, title_box[1] + 20), title_text, font=title_font, fill=NAVY, spacing=title_spacing)

    body_top = title_box[3] + 18
    body_box = (card_left + 24, body_top, card_left + card_width - 24, card_top + card_height - 24)
    draw.rounded_rectangle(body_box, radius=20, fill=(255, 248, 239, 228), outline=(242, 154, 46, 180), width=2)

    text_x = body_box[0] + 64
    dot_x = body_box[0] + 36
    current_y = body_box[1] + 24
    for bullet_text, bullet_height in zip(bullet_texts, bullet_heights):
        dot_y = current_y + 17
        draw.ellipse((dot_x - 8, dot_y - 8, dot_x + 8, dot_y + 8), fill=(242, 154, 46, 255))
        draw.multiline_text((text_x, current_y), bullet_text, font=bullet_font, fill=NAVY, spacing=bullet_spacing)
        current_y += bullet_height + 24

    draw.text((WIDTH - 112, HEIGHT - 72), f"{page.page:02d}", font=OVERLAY_PAGE_FONT, fill=(31, 111, 178, 210))
    return Image.alpha_composite(base, overlay).convert("RGB")


def render_page_image(raw_path: Path, final_path: Path, page: SlidePage, native_text: bool = False) -> None:
    image = Image.open(raw_path).convert("RGB").resize((WIDTH, HEIGHT))
    if not native_text:
        image = render_native_text_overlay(image, page)
    image.save(final_path, quality=95)


def generate_dry_run_image(page: SlidePage, output_path: Path) -> str:
    image = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((150, 130, 170, 380), radius=10, fill="#F29A2E")
    draw.rounded_rectangle((960, 120, 1760, 920), radius=52, fill="#F5FAFF", outline="#E2E8F0", width=3)
    draw.ellipse((1120, 240, 1540, 660), fill="#DCEEFE", outline="#1F6FB2", width=6)
    draw.rounded_rectangle((1120, 700, 1600, 820), radius=32, fill="#FDE9CF", outline="#F29A2E", width=4)
    draw.line((980, 760, 1100, 760), fill="#1F6FB2", width=8)
    draw.line((1088, 746, 1120, 760), fill="#1F6FB2", width=8)
    draw.line((1088, 774, 1120, 760), fill="#1F6FB2", width=8)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, quality=94)
    return "dry-run"


def build_outline_markdown(deck_title: str, pages: list[SlidePage]) -> str:
    lines = [f"# {deck_title}", ""]
    for page in pages:
        lines.append(f"## {page.page:02d}. {page.context_title}")
        for bullet in page.context_points[:4]:
            lines.append(f"- {bullet}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def save_ppt(images: list[Path], output_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def write_readme(
    output_dir: Path,
    source_path: Path,
    deck_title: str,
    pages: list[SlidePage],
    raw_paths: list[Path],
    outline_path: Path,
    pptx_path: Path,
    style_reference: Path | None,
    layout_reference: Path | None,
    page_reference_dir: Path | None,
    reference_mode: str,
    image_provider: str,
    native_text: bool = False,
) -> None:
    provider_label = "GPT Image / OpenAI-compatible" if image_provider == "openai" else "Google Gemini"
    normalized_reference_mode = normalize_reference_mode(reference_mode)
    if normalized_reference_mode == "prompt_only":
        generation_mode = (
            f"生成模式：{provider_label} 真生图 + 纯提示词构图 + 中文文字直接融合进画面。"
            if native_text
            else f"生成模式：{provider_label} 真生图 + 纯提示词构图 + 本地系统字体后加覆盖。"
        )
    elif normalized_reference_mode == "style_only":
        generation_mode = (
            f"生成模式：{provider_label} 真生图 + 风格参考图 + 中文文字直接融合进画面。"
            if native_text
            else f"生成模式：{provider_label} 真生图 + 风格参考图 + 本地系统字体后加覆盖。"
        )
    else:
        generation_mode = (
            f"生成模式：{provider_label} 真生图 + 参考图喂模 + 融合版图内文字直出。"
            if native_text and (layout_reference or page_reference_dir)
            else f"生成模式：{provider_label} 真生图 + 参考图喂模 + 中文文字由模型直接在图内生成。"
            if native_text
            else f"生成模式：{provider_label} 真生图 + 参考图喂模 + 本地系统字体后加覆盖。"
        )
    lines = [
        f"这是根据《{source_path.name}》自动分析后生成的漫画风图片型 PPT，共 {len(pages)} 页。",
        generation_mode,
        f"文档标题：{deck_title}",
        f"源文档：{source_path}",
        f"参考图模式：{reference_mode_label(normalized_reference_mode)}",
        f"风格参考图：{style_reference}" if style_reference else "风格参考图：未使用",
        f"布局参考图：{layout_reference}" if layout_reference else "布局参考图：未额外使用",
        f"按页融合参考目录：{page_reference_dir}" if page_reference_dir else "按页融合参考目录：未额外使用",
        f"大纲：{outline_path}",
        f"PPT：{pptx_path}",
        "",
        "页面规划：",
        *[f"P{page.page}: {page.context_title}" for page in pages],
        "",
        "原始生图底图：",
        *[str(path) for path in raw_paths],
    ]
    (output_dir / "README.txt").write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a comic-style PPT deck from Markdown using Gemini image generation.")
    parser.add_argument("--source", type=Path, required=True, help="Markdown source path.")
    parser.add_argument("--project-name", type=str, default="", help="Optional project title override.")
    parser.add_argument("--image-provider", type=str, default="", choices=["", "gemini", "gpt", "openai", "openai-compatible", "gpt-image-2"], help="Image generation provider route.")
    parser.add_argument("--reference-mode", type=str, default="", choices=["", "full", "style_only", "prompt_only", "style-only", "prompt-only"], help="Reference usage mode.")
    parser.add_argument("--style-reference", type=Path, default=None, help="Optional style reference image override.")
    parser.add_argument("--layout-reference", type=Path, default=None, help="Optional layout/fusion reference image override.")
    parser.add_argument("--page-reference-dir", type=Path, default=None, help="Optional directory of per-page fusion reference images.")
    parser.add_argument("--output-dir", type=Path, default=None, help="Exact output directory for generated assets.")
    parser.add_argument("--output-root", type=Path, default=OUTPUT_ROOT, help="Fallback root when output-dir is omitted.")
    parser.add_argument("--pptx-path", type=Path, default=None, help="Exact final PPTX output path.")
    parser.add_argument("--outline-path", type=Path, default=None, help="Exact outline markdown output path.")
    parser.add_argument("--dry-run", action="store_true", help="Skip Gemini and render local placeholders for integration testing.")
    parser.add_argument("--native-text", action="store_true", help="Ask Gemini to render the Chinese text directly in the image.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    source_path = args.source.expanduser().resolve()
    if not source_path.exists():
        raise RuntimeError(f"Source markdown not found: {source_path}")
    image_provider = normalize_image_provider(args.image_provider or os.environ.get("COMIC_IMAGE_PROVIDER") or os.environ.get("APP_FACTORY_COMIC_IMAGE_PROVIDER"))
    reference_mode = normalize_reference_mode(
        args.reference_mode
        or os.environ.get("COMIC_REFERENCE_MODE")
        or os.environ.get("APP_FACTORY_COMIC_REFERENCE_MODE")
        or DEFAULT_COMIC_REFERENCE_MODE
    )
    style_reference = args.style_reference.expanduser().resolve() if args.style_reference else STYLE_REFERENCE
    layout_reference = args.layout_reference.expanduser().resolve() if args.layout_reference else LAYOUT_REFERENCE
    page_reference_dir = args.page_reference_dir.expanduser().resolve() if args.page_reference_dir else None
    if reference_mode == "prompt_only":
        style_reference = None
        layout_reference = None
        page_reference_dir = None
    elif reference_mode == "style_only":
        layout_reference = None
        page_reference_dir = None
    if page_reference_dir and not page_reference_dir.exists():
        raise RuntimeError(f"Page reference directory not found: {page_reference_dir}")

    markdown = source_path.read_text(encoding="utf-8")
    output_dir = args.output_dir.resolve() if args.output_dir else create_timestamped_output_dir(args.output_root, f"{sanitize_filename_part(source_path.stem)}-漫画风PPT")
    output_dir.mkdir(parents=True, exist_ok=True)
    raw_dir = output_dir / "gemini-backdrops"
    raw_dir.mkdir(parents=True, exist_ok=True)

    emit_progress("prepare_markdown", "正在解析文档结构", 8)
    deck_title, pages = build_pages_from_markdown(markdown, source_path)
    if args.project_name.strip():
        deck_title = args.project_name.strip()
    emit_log("analysis", "文档分析结果", build_analysis_summary(markdown, source_path, deck_title, pages))
    emit_progress("plan_story", f"已规划 {len(pages)} 页漫画分镜", 18, total=len(pages))
    emit_log("design", "设计思路", build_design_summary(pages))

    prompt_rows: list[dict[str, object]] = []
    image_paths: list[Path] = []
    raw_paths: list[Path] = []
    total = len(pages)

    for index, page in enumerate(pages, start=1):
        percent = 18 + int(index / max(total, 1) * 66)
        emit_log(
            "page",
            f"第 {index}/{total} 页设计",
            f"页面主题《{page.context_title}》；画面重点：{page.scene_prompt}；文案：{' / '.join(page.bullets[:2])}",
            page=index,
            total=total,
        )
        emit_progress(
            "generate_pages",
            f"正在生成第 {index}/{total} 页漫画插画",
            percent,
            current=index,
            total=total,
        )
        raw_path = raw_dir / f"P{index}-gemini-backdrop.png"
        final_path = output_dir / page.filename
        page_reference = None
        if page_reference_dir:
            candidate = page_reference_dir / page.filename
            if candidate.exists():
                page_reference = candidate
            else:
                matches = sorted(page_reference_dir.glob(f"P{page.page}-*.png"))
                if matches:
                    page_reference = matches[0]
        generate_image = generate_openai_compatible_image if image_provider == "openai" else generate_gemini_image
        mode = (
            generate_dry_run_image(page, raw_path)
            if args.dry_run
            else generate_image(
                page,
                raw_path,
                native_text=args.native_text,
                style_reference=style_reference,
                layout_reference=layout_reference,
                page_reference=page_reference,
            )
        )
        render_page_image(raw_path, final_path, page, native_text=args.native_text)
        image_paths.append(final_path)
        raw_paths.append(raw_path)
        prompt_rows.append(
            {
                "page": index,
                "title": page.title,
                "bullets": page.bullets,
                "contextTitle": page.context_title,
                "contextPoints": page.context_points,
                "scenePrompt": page.scene_prompt,
                "prompt": build_prompt(
                    page,
                    native_text=args.native_text,
                    image_provider=image_provider,
                    has_style_reference=style_reference is not None and style_reference.exists(),
                    has_layout_reference=layout_reference is not None and layout_reference.exists(),
                    has_page_reference=page_reference is not None and page_reference.exists(),
                ),
                "imageProvider": image_provider,
                "referenceMode": reference_mode,
                "styleReference": str(style_reference or ""),
                "layoutReference": str(layout_reference or ""),
                "pageReference": str(page_reference or ""),
                "rawBackdrop": str(raw_path),
                "mode": mode,
            }
        )

    outline_path = args.outline_path.resolve() if args.outline_path else output_dir / "outline.md"
    outline_path.parent.mkdir(parents=True, exist_ok=True)
    outline_path.write_text(build_outline_markdown(deck_title, pages), encoding="utf-8")

    pptx_path = args.pptx_path.resolve() if args.pptx_path else output_dir / f"{sanitize_filename_part(deck_title)}-漫画风PPT.pptx"
    emit_progress("package_ppt", "正在打包 PPT 文件", 92)
    emit_log("progress", "打包导出", "全部页面已完成，正在打包成最终 PPT 文件。")
    save_ppt(image_paths, pptx_path)

    slide_plan_path = output_dir / "slide-plan.json"
    slide_plan_path.write_text(
        json.dumps(
            {
                "deckTitle": deck_title,
                "slideCount": len(pages),
                "slides": [
                    {
                        "page": page.page,
                        "title": page.title,
                        "contextTitle": page.context_title,
                        "bullets": page.context_points,
                        "scenePrompt": page.scene_prompt,
                        "filename": page.filename,
                    }
                    for page in pages
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    prompt_catalog_path = output_dir / "gemini-prompts.json"
    prompt_catalog_path.write_text(json.dumps(prompt_rows, ensure_ascii=False, indent=2), encoding="utf-8")
    write_readme(
        output_dir,
        source_path,
        deck_title,
        pages,
        raw_paths,
        outline_path,
        pptx_path,
        style_reference,
        layout_reference,
        page_reference_dir,
        reference_mode,
        image_provider,
        native_text=args.native_text,
    )
    manifest_path = write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(pages),
        deckTitle=deck_title,
        sourceDoc=str(source_path),
        outlinePath=str(outline_path),
        slidePlan=str(slide_plan_path),
        promptCatalog=str(prompt_catalog_path),
        imageDirectory=str(output_dir),
        geminiBackdropDirectory=str(raw_dir),
        imageProvider=image_provider,
        referenceMode=reference_mode,
        styleReference=str(style_reference or ""),
        layoutReference=str(layout_reference or ""),
        pageReferenceDirectory=str(page_reference_dir or ""),
        mode=f"comic-{'gpt-image' if image_provider == 'openai' else 'gemini'}-{reference_mode.replace('_', '-')}-{'direct-text' if args.native_text else 'overlay-text'}",
        dryRun=args.dry_run,
    )

    result_payload = {
        "deckTitle": deck_title,
        "slideCount": len(pages),
        "pptxPath": str(pptx_path),
        "outlinePath": str(outline_path),
        "outputDir": str(output_dir),
        "slidePlanPath": str(slide_plan_path),
        "promptCatalogPath": str(prompt_catalog_path),
        "manifestPath": str(manifest_path),
        "dryRun": args.dry_run,
    }
    emit_progress("completed", "漫画风 PPT 已生成完成", 100, total=len(pages))
    emit_log("success", "完成", f"已输出 {len(pages)} 页 PPT，结果目录：{output_dir}")
    print(f"__RESULT__ {json.dumps(result_payload, ensure_ascii=False)}", flush=True)


if __name__ == "__main__":
    main()

from __future__ import annotations

import base64
import argparse
import json
import os
import subprocess
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    latest_matching_path,
    load_env_candidates,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_neural_params_controlled_demo import (
    WIDTH,
    HEIGHT,
    WHITE,
    BLUE,
    ORANGE,
    NAVY,
    BODY_FONT,
    TITLE_FONT,
    SCREEN_FONT,
    SCREEN_AI_FONT,
    draw_centered,
)

DESKTOP = Path("/Users/mac/Desktop")
ARCHIVE_DIR = latest_archive_dir()
REF_IMAGE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    DESKTOP / "参考图-大模型能力突然变强.png",
)
TEMPLATE_PPTX = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "神经网络学习的重要参数-Gemini参考图版.pptx" if ARCHIVE_DIR else None,
    DESKTOP / "神经网络学习的重要参数-Gemini参考图版.pptx",
)
MODEL = os.environ.get("GEMINI_IMAGE_MODEL", "gemini-3-pro-image-preview")
OUTPUT_ROOT = repo_runtime_dir("single-slide-demo")
OUTPUT_LABEL = "神经网络学习-gemini参考图版"


def get_api_key() -> str:
    load_env_candidates()
    for key in ["GEMINI_API_KEY", "GOOGLE_API_KEY", "GOOGLE_GENAI_API_KEY"]:
        value = os.environ.get(key)
        if value:
            return value
    raise RuntimeError("Gemini API key not found in env.")


def build_prompt() -> str:
    return "\n".join(
        [
            "Use the provided image strictly as a visual style reference, not as content to copy.",
            "Create a 16:9 educational infographic slide in the same visual language:",
            "- white background",
            "- the top headline will be added later as plain deep-blue text, so keep the top area empty",
            "- clean blue and orange accent colors",
            "- thick black outlines",
            "- rounded comic speech bubble style",
            "- subtle geometric corner decorations inspired by the reference image",
            "- polished classroom infographic quality",
            "",
            "Very important composition rules:",
            "1. Reserve clean empty white space across the top for a large Chinese title and two short bullet lines to be added later.",
            "2. In the lower half, place exactly these subjects from left to right on one horizontal line:",
            "   - one seated husky dog, clearly recognizable as a husky, tail wagging",
            "   - one computer-head stick figure facing forward",
            "   - one computer-head stick figure facing forward",
            "   - one computer-head stick figure facing forward",
            "3. Connect the four subjects with orange arrows.",
            "4. Leave the three computer screens empty and clean for later text overlay.",
            "5. Leave one empty speech bubble above the rightmost computer-head figure for later text overlay.",
            "",
            "Hard constraints:",
            "- no readable text",
            "- no numbers",
            "- no formulas",
            "- no books",
            "- no servers",
            "- no extra characters",
            "- no extra icons except subtle corner ornaments",
            "- no filled title band",
            "- no large colored title banner or header block",
            "- no watermark",
            "- keep the layout clean and spacious",
        ]
    )


def generate_reference_image(out_raw: Path, out_prompt: Path) -> str:
    prompt = build_prompt()
    if REF_IMAGE is None:
        raise RuntimeError("Reference image not found in archive or desktop.")
    out_prompt.write_text(prompt, encoding="utf-8")
    try:
        api_key = get_api_key()
    except RuntimeError:
        out_raw.write_bytes(REF_IMAGE.read_bytes())
        return "archived-reference-fallback"
    image_b64 = base64.b64encode(REF_IMAGE.read_bytes()).decode("utf-8")
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": prompt},
                    {"inline_data": {"mime_type": "image/png", "data": image_b64}},
                ]
            }
        ],
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
            "imageConfig": {"aspectRatio": "16:9", "imageSize": "2K"},
        },
    }
    response = subprocess.run(
        [
            "curl",
            "-sS",
            "-X",
            "POST",
            f"https://generativelanguage.googleapis.com/v1beta/models/{MODEL}:generateContent",
            "-H",
            f"x-goog-api-key: {api_key}",
            "-H",
            "Content-Type: application/json",
            "--data-binary",
            "@-",
        ],
        input=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        capture_output=True,
        timeout=300,
        check=False,
    )
    if response.returncode != 0:
        raise RuntimeError(response.stderr.decode("utf-8", errors="replace") or f"curl exited {response.returncode}")
    data = json.loads(response.stdout.decode("utf-8", errors="replace"))
    if data.get("error"):
        raise RuntimeError(json.dumps(data["error"], ensure_ascii=False))
    parts = data.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    for part in parts:
        inline_data = part.get("inlineData") or part.get("inline_data")
        if inline_data and inline_data.get("data"):
            out_raw.write_bytes(base64.b64decode(inline_data["data"]))
            return "gemini-api"
    raise RuntimeError("Gemini did not return image data.")


def overlay_text(out_raw: Path, out_final: Path) -> None:
    image = Image.open(out_raw).convert("RGB").resize((WIDTH, HEIGHT))
    image = Image.blend(image, Image.new("RGB", (WIDTH, HEIGHT), WHITE), 0.10).filter(ImageFilter.SMOOTH)
    draw = ImageDraw.Draw(image)

    draw_centered(draw, "神经网络学习的重要参数", 74, TITLE_FONT, BLUE)
    draw_centered(
        draw,
        "• 激活函数：决定神经元是否“激活”\n• 常见函数：ReLU，Sigmoid，Tanh",
        198,
        BODY_FONT,
        NAVY,
        max_width=1100,
        spacing=14,
    )

    positions = [(748, 562), (1168, 562), (1590, 562)]
    texts = ["耳朵*权重+\n颜色*权重+\n...=2.7", "激活函数*\n(2.7)=1", "AI"]
    fonts = [SCREEN_FONT, SCREEN_FONT, SCREEN_AI_FONT]
    fills = [NAVY, NAVY, BLUE]
    for (cx, cy), text, font, fill in zip(positions, texts, fonts, fills):
        box = draw.multiline_textbbox((0, 0), text, font=font, spacing=3, align="center")
        draw.multiline_text((cx - (box[2] - box[0]) / 2, cy - (box[3] - box[1]) / 2), text, font=font, fill=fill, spacing=3, align="center")

    speech = "这是狗"
    box = draw.textbbox((0, 0), speech, font=BODY_FONT)
    draw.text((1704 - (box[2] - box[0]) / 2, 340 - (box[3] - box[1]) / 2), speech, font=BODY_FONT, fill=NAVY)
    image.save(out_final, quality=95)


def save_ppt(out_final: Path, out_pptx: Path) -> None:
    if Presentation is not None and Inches is not None:
        prs = Presentation()
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(out_final), 0, 0, width=prs.slide_width, height=prs.slide_height)
        prs.save(str(out_pptx))
        return

    if TEMPLATE_PPTX is None or not TEMPLATE_PPTX.exists():
        raise RuntimeError(f"python-pptx is not installed and PPTX template is missing: {TEMPLATE_PPTX}")

    temp_pptx = out_pptx.with_suffix(".tmp.pptx")
    replaced = False
    with zipfile.ZipFile(TEMPLATE_PPTX, "r") as source, zipfile.ZipFile(temp_pptx, "w") as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == "ppt/media/image1.png":
                data = out_final.read_bytes()
                replaced = True
            target.writestr(item, data)

    if not replaced:
        temp_pptx.unlink(missing_ok=True)
        raise RuntimeError(f"No replaceable slide image found in PPTX template: {TEMPLATE_PPTX}")
    temp_pptx.replace(out_pptx)


def find_latest_output_dir(output_root: Path) -> Path | None:
    if not output_root.exists():
        return None
    candidates = [path for path in output_root.iterdir() if path.is_dir() and path.name.endswith(OUTPUT_LABEL)]
    if not candidates:
        return None
    return sorted(candidates)[-1]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the Gemini-reference neural-parameters slide.")
    parser.add_argument(
        "--reuse-raw",
        action="store_true",
        help="Reuse the existing Gemini raw background instead of calling the API again.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=OUTPUT_ROOT,
        help="Directory where timestamped local result folders are created.",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=None,
        help="Reuse or write into a specific local result directory.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if args.output_dir is not None:
        output_dir = args.output_dir
        output_dir.mkdir(parents=True, exist_ok=True)
    elif args.reuse_raw:
        output_dir = find_latest_output_dir(args.output_root) or args.output_root / OUTPUT_LABEL
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        output_dir = create_timestamped_output_dir(args.output_root, OUTPUT_LABEL)
    out_raw = output_dir / "神经网络学习的重要参数-Gemini参考图版-底图.png"
    out_final = output_dir / "神经网络学习的重要参数-Gemini参考图版-v2.png"
    out_pptx = output_dir / "神经网络学习的重要参数-Gemini参考图版-v2.pptx"
    out_prompt = output_dir / "神经网络学习-Gemini参考图提示词.txt"
    if args.reuse_raw:
        if not out_raw.exists():
            raise RuntimeError(f"Raw Gemini background not found: {out_raw}")
        background_mode = "reuse-raw"
    else:
        background_mode = generate_reference_image(out_raw, out_prompt)
    overlay_text(out_raw, out_final)
    save_ppt(out_final, out_pptx)
    write_result_manifest(
        output_dir,
        out_pptx,
        previewImage=str(out_final),
        rawBackground=str(out_raw),
        promptFile=str(out_prompt),
        referenceImage=str(REF_IMAGE) if REF_IMAGE else "",
        backgroundMode=background_mode,
        mode="gemini-reference-local-overlay",
    )
    print(out_pptx)


if __name__ == "__main__":
    main()

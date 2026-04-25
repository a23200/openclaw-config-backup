from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw
from local_result import create_timestamped_output_dir, write_result_manifest

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import (
    AI_FONT,
    BLUE,
    BLUE_LIGHT,
    BODY_BOLD_FONT,
    BODY_FONT,
    DEFAULT_OUTPUT_ROOT,
    GRAY,
    GRAY_LIGHT,
    GREEN,
    GREEN_LIGHT,
    HEIGHT,
    HUGE_FONT,
    INK,
    LABEL_FONT,
    NAVY,
    ORANGE,
    ORANGE_LIGHT,
    RED,
    RED_LIGHT,
    SECTION_FONT,
    SMALL_BOLD_FONT,
    SMALL_FONT,
    TINY_BOLD_FONT,
    WHITE,
    WIDTH,
    create_contact_sheet,
    draw_ai_person,
    draw_annotation,
    draw_arrow,
    draw_bottom_rule,
    draw_card,
    draw_chart,
    draw_chip,
    draw_clipboard,
    draw_cloud,
    draw_dashboard,
    draw_gear,
    draw_inn,
    draw_metric_tile,
    draw_sketch_round_rect,
    draw_slide_background,
    draw_speech,
    draw_stick_person,
    draw_text_in_box,
    draw_title_block,
    load_font,
)


DESKTOP = Path("/Users/mac/Desktop")
SOURCE_DOC = DESKTOP / "项目核心竞争力分析.md"
OUTPUT_ROOT = DEFAULT_OUTPUT_ROOT


@dataclass(frozen=True)
class LayerData:
    title: str
    rival: str
    system: str
    conclusion: str
    examples: list[str]


@dataclass(frozen=True)
class StoryData:
    title: str
    intro_heading: str
    intro_body: str
    layers: list[LayerData]
    final_conclusion: str


def clean_text(text: str) -> str:
    text = text.replace("`", "")
    text = text.replace("**", "")
    return re.sub(r"\s+", " ", text).strip(" -\n")


def parse_markdown_story(path: Path) -> StoryData:
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()

    title = ""
    intro_heading = ""
    intro_lines: list[str] = []
    final_lines: list[str] = []
    layers: list[tuple[str, list[str]]] = []

    current_layer_title: str | None = None
    current_layer_lines: list[str] = []
    current_h2 = ""

    for raw_line in lines:
        line = raw_line.rstrip()
        if line.startswith("# "):
            title = clean_text(line[2:])
            continue
        if line.startswith("## "):
            if current_layer_title:
                layers.append((current_layer_title, current_layer_lines[:]))
                current_layer_title = None
                current_layer_lines = []
            current_h2 = clean_text(line[3:])
            if current_h2.startswith("引言"):
                intro_heading = current_h2
            continue
        if line.startswith("### "):
            if current_layer_title:
                layers.append((current_layer_title, current_layer_lines[:]))
            current_layer_title = clean_text(line[4:])
            current_layer_lines = []
            continue

        if current_layer_title:
            current_layer_lines.append(line)
        elif current_h2.startswith("引言"):
            if line.strip() and line.strip() != "---":
                intro_lines.append(clean_text(line))
        elif current_h2.startswith("最终结论"):
            if line.strip() and line.strip() != "---":
                final_lines.append(clean_text(line))

    if current_layer_title:
        layers.append((current_layer_title, current_layer_lines))

    parsed_layers: list[LayerData] = []
    for title_text, body_lines in layers:
        rival = ""
        system = ""
        conclusion = ""
        examples: list[str] = []
        current_target: str | None = None
        capture_examples = False

        for raw in body_lines:
            stripped = raw.strip()
            if not stripped:
                continue
            if "对手的“Skill”" in stripped:
                rival = clean_text(stripped.split("：", 1)[1])
                current_target = "rival"
                capture_examples = False
                continue
            if "我们的“系统”" in stripped:
                system = clean_text(stripped.split("：", 1)[1])
                current_target = "system"
                capture_examples = False
                continue
            if "我们将知道" in stripped:
                capture_examples = True
                current_target = None
                continue
            if "结论" in stripped:
                conclusion = clean_text(stripped.split("：", 1)[1])
                current_target = "conclusion"
                capture_examples = False
                continue
            if capture_examples and stripped.startswith("-"):
                examples.append(clean_text(stripped.lstrip("- ")))
                continue
            if current_target == "rival":
                rival = f"{rival} {clean_text(stripped)}".strip()
            elif current_target == "system":
                system = f"{system} {clean_text(stripped)}".strip()
            elif current_target == "conclusion":
                conclusion = f"{conclusion} {clean_text(stripped)}".strip()

        parsed_layers.append(
            LayerData(
                title=title_text,
                rival=rival,
                system=system,
                conclusion=conclusion,
                examples=examples,
            )
        )

    return StoryData(
        title=title or "项目核心竞争力分析",
        intro_heading=intro_heading or "引言：我们构建的不是“技能”，而是“场域”",
        intro_body=" ".join(intro_lines),
        layers=parsed_layers,
        final_conclusion=" ".join(final_lines),
    )


def draw_center_text(draw: ImageDraw.ImageDraw, text: str, center_x: int, y: int, font, fill: str) -> None:
    box = draw.multiline_textbbox((0, 0), text, font=font, spacing=6, align="center")
    draw.multiline_text((center_x - (box[2] - box[0]) / 2, y), text, font=font, fill=fill, spacing=6, align="center")


def draw_pill(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, color: str, fill: str | None = None, w: int = 214) -> None:
    inner_fill = fill or (BLUE_LIGHT if color == BLUE else ORANGE_LIGHT if color == ORANGE else GREEN_LIGHT)
    draw.rounded_rectangle((x, y, x + w, y + 46), radius=23, fill=inner_fill, outline=color, width=3)
    box = draw.textbbox((0, 0), text, font=LABEL_FONT)
    draw.text((x + (w - (box[2] - box[0])) / 2, y + 10), text, font=LABEL_FONT, fill=color)


def draw_stage_box(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, color: str, body: str, w: int = 190, h: int = 110) -> None:
    fill = BLUE_LIGHT if color == BLUE else ORANGE_LIGHT if color == ORANGE else GREEN_LIGHT
    draw_sketch_round_rect(draw, (x, y, x + w, y + h), radius=20, fill=fill, outline=color, width=3)
    draw.text((x + 18, y + 16), title, font=SMALL_BOLD_FONT, fill=color)
    draw_text_in_box(draw, (x + 18, y + 52), body, SMALL_FONT, NAVY, w - 36, spacing=5)


def draw_core_circle(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, subtitle: str, accent: str = BLUE) -> None:
    draw.ellipse((x - 160, y - 160, x + 160, y + 160), fill=WHITE, outline=accent, width=6)
    draw.ellipse((x - 132, y - 132, x + 132, y + 132), fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT, outline=accent, width=3)
    draw_center_text(draw, title, x, y - 44, HUGE_FONT, accent)
    draw_center_text(draw, subtitle, x, y + 40, SMALL_BOLD_FONT, NAVY)


def draw_quote_band(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, text: str, color: str = ORANGE) -> None:
    fill = ORANGE_LIGHT if color == ORANGE else BLUE_LIGHT
    draw_sketch_round_rect(draw, (x, y, x + w, y + 94), radius=26, fill=fill, outline=color, width=4)
    draw_text_in_box(draw, (x + 26, y + 22), f"“{text}”", BODY_FONT, NAVY, w - 52, spacing=7)


def draw_point_node(draw: ImageDraw.ImageDraw, x: int, y: int, label: str, accent: str = BLUE, secondary: str = "") -> None:
    draw.ellipse((x - 54, y - 54, x + 54, y + 54), fill=WHITE, outline=accent, width=4)
    draw.ellipse((x - 36, y - 36, x + 36, y + 36), fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT, outline=accent, width=3)
    draw_center_text(draw, label, x, y - 14, LABEL_FONT, accent)
    if secondary:
        draw_center_text(draw, secondary, x, y + 20, TINY_BOLD_FONT, NAVY)


def draw_small_success_node(draw: ImageDraw.ImageDraw, x: int, y: int, label: str, accent: str = BLUE) -> None:
    draw.rounded_rectangle((x, y, x + 118, y + 68), radius=20, fill=WHITE, outline=accent, width=3)
    draw.rounded_rectangle((x + 12, y + 12, x + 46, y + 56), radius=12, fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT, outline=accent, width=2)
    draw.text((x + 58, y + 22), label, font=LABEL_FONT, fill=accent)


def build_cover(story: StoryData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "项目核心竞争力分析", "我们真正的护城河，不是 Skill，而是 Field", small=False)

    draw_core_circle(draw, 960, 530, "FIELD", "场域 / 权力系统", accent=BLUE)
    card_specs = [
        (216, 336, "L1", "闭环进化", "每运行一次，就更强一次", BLUE),
        (1308, 336, "L2", "欲望图谱", "从单点视角升级到市场定义权", ORANGE),
        (216, 650, "L3", "人机混智", "1%战略洞察 + 99%规模执行", ORANGE),
        (1308, 650, "L4", "平台扩张", "把成功模式从 1 家复制到 100 家", BLUE),
    ]
    anchors = [(540, 412), (1380, 412), (540, 658), (1380, 658)]
    for (x, y, icon, title, subtitle, accent), (ax, ay) in zip(card_specs, anchors):
        draw_card(draw, x, y, 396, 164, title, subtitle, accent, icon)
        draw_arrow(draw, 960, 530, ax, ay, color=accent, width=6)

    draw_stick_person(draw, 662, 896, accent=ORANGE, pose="point", scale=1.0)
    draw_ai_person(draw, 1260, 892, label="AI", accent=BLUE, scale=1.0)
    draw_pill(draw, 548, 766, "指挥官：战略判断", ORANGE, w=248)
    draw_pill(draw, 1118, 766, "AI执行官：规模执行", BLUE, w=286)
    draw_arrow(draw, 800, 874, 1110, 874, color=ORANGE, width=6)

    quote = story.final_conclusion.split("。")[0].replace("竞争对手或许可以复制我们“做什么（What）”，但他们永远无法复制我们“如何成为（How）”", "What 可复制，How 复制不了")
    draw_quote_band(draw, 404, 892, 1116, quote or "What 可复制，How 复制不了", color=GREEN)
    draw_bottom_rule(draw, 1, 7, "封面：从 Skill 到 Field 的四层护城河")
    return image


def build_intro(story: StoryData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "引言：我们构建的不是“技能”，而是“场域”", "线性工具可以复制，能进化的系统才是真正护城河", small=True)

    draw_pill(draw, 238, 280, "线性 Skill", ORANGE, w=240)
    draw_stage_box(draw, 180, 360, "输入", ORANGE, "一次任务开始", w=210, h=118)
    draw_stage_box(draw, 438, 360, "处理", ORANGE, "按规则执行", w=210, h=118)
    draw_stage_box(draw, 696, 360, "输出", ORANGE, "一次性交付", w=210, h=118)
    draw_arrow(draw, 390, 418, 438, 418, color=ORANGE, width=6)
    draw_arrow(draw, 648, 418, 696, 418, color=ORANGE, width=6)
    draw_annotation(draw, 202, 540, "静态特征", "单向流程、不会因为使用而变聪明。", ORANGE)
    draw_annotation(draw, 518, 540, "复制门槛", "别人抄走 SOP，就能做出相似工具。", ORANGE)

    draw_pill(draw, 1168, 280, "场域 Field", BLUE, w=240)
    loop_nodes = [
        ("投放", 1370, 334),
        ("反馈", 1558, 482),
        ("学习", 1488, 690),
        ("决策", 1258, 690),
        ("执行", 1188, 482),
    ]
    for label, x, y in loop_nodes:
        draw_point_node(draw, x, y, label, accent=BLUE)
    arrows = [(1410, 362, 1528, 448), (1550, 536, 1498, 646), (1444, 716, 1300, 716), (1218, 646, 1178, 536), (1188, 444, 1318, 344)]
    for x1, y1, x2, y2 in arrows:
        draw_arrow(draw, x1, y1, x2, y2, color=BLUE, width=5)
    draw_core_circle(draw, 1372, 518, "FIELD", "执行后反哺系统", accent=GREEN)
    draw_annotation(draw, 1128, 786, "系统灵魂", "“效果量化”把市场反馈变成下一次决策养料。", GREEN)
    draw_metric_tile(draw, 1508, 778, "∞", "闭环循环", BLUE)

    draw_quote_band(draw, 226, 870, 1468, story.intro_body, color=BLUE)
    draw_bottom_rule(draw, 2, 7, "引言：线性工具与进化场域的本质差异")
    return image


def build_layer_one(layer: LayerData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "第一层护城河：闭环自进化系统", "静态工具只能执行；动态系统会越跑越强", small=False)

    draw_pill(draw, 170, 282, "对手的 Skill", ORANGE, w=226)
    draw_annotation(draw, 126, 344, "线性工具", layer.rival, ORANGE)

    draw_pill(draw, 1518, 282, "我们的系统", BLUE, w=226)
    draw_annotation(draw, 1512, 344, "闭环生命体", layer.system, BLUE)

    loop_card_positions = [
        ("投放动作", 650, 298, BLUE),
        ("效果量化", 1110, 298, GREEN),
        ("反馈学习", 1180, 610, ORANGE),
        ("策略迭代", 578, 610, BLUE),
    ]
    for title, x, y, accent in loop_card_positions:
        draw_stage_box(draw, x, y, title, accent, "真实市场数据回流", w=240, h=118)
    draw_arrow(draw, 890, 354, 1110, 354, color=BLUE, width=6)
    draw_arrow(draw, 1232, 416, 1232, 610, color=GREEN, width=6)
    draw_arrow(draw, 1180, 722, 818, 722, color=ORANGE, width=6)
    draw_arrow(draw, 694, 610, 694, 416, color=BLUE, width=6)
    draw_core_circle(draw, 960, 540, "LOOP", "每次执行都完成一次市场学习", accent=GREEN)
    draw_chart(draw, 842, 782, scale=1.8)
    draw_metric_tile(draw, 1128, 770, "+1", "执行后认知增量", GREEN)
    draw_quote_band(draw, 412, 884, 1096, layer.conclusion or "我们的系统是动态进化的。它每运行一次，就会比上一次更强。", color=ORANGE)
    draw_bottom_rule(draw, 3, 7, "第一层：效果量化驱动的闭环进化")
    return image


def build_layer_two(layer: LayerData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "第二层护城河：垄断性数据资产", "从单点信息升级为区域级“游客欲望图谱”", small=False)

    draw_pill(draw, 156, 284, "单点视角", ORANGE, w=218)
    draw_point_node(draw, 274, 432, "A民宿", accent=ORANGE, secondary="只看到自己")
    draw_annotation(draw, 92, 550, "对手局限", layer.rival, ORANGE)

    draw_pill(draw, 828, 256, "游客欲望图谱", BLUE, w=264)
    draw_cloud(draw, 960, 470, scale=1.8)
    draw_center_text(draw, "欲望图谱", 960, 454, BODY_BOLD_FONT, BLUE)
    nodes = [
        (748, 364, "民宿", BLUE),
        (1184, 364, "餐厅", ORANGE),
        (700, 594, "旅拍", ORANGE),
        (1226, 596, "亲子", BLUE),
        (958, 700, "咖啡", GREEN),
    ]
    for x, y, label, accent in nodes:
        draw_point_node(draw, x, y, label, accent=accent)
        draw_arrow(draw, x, y, 960, 470, color=accent, width=5)

    draw_pill(draw, 1460, 284, "上帝视角", GREEN, w=218)
    draw_annotation(draw, 1436, 344, "系统价值", layer.system, BLUE)
    example_boxes = [
        (1364, 540, layer.examples[0] if len(layer.examples) > 0 else "A民宿的客人，下一步想去 B 餐厅。", BLUE),
        (1364, 662, layer.examples[1] if len(layer.examples) > 1 else "“寂静风”游客也在搜索“手冲咖啡”。", ORANGE),
        (1364, 784, layer.examples[2] if len(layer.examples) > 2 else "“旅拍”热度下降时，“亲子游”正在抬头。", GREEN),
    ]
    for x, y, text, accent in example_boxes:
        draw_annotation(draw, x, y, "样本洞察", text, accent)

    draw_quote_band(draw, 292, 878, 1328, layer.conclusion or "这个欲望图谱本身就是最核心的资产，它让我们有能力定义整个市场的消费风向。", color=GREEN)
    draw_bottom_rule(draw, 4, 7, "第二层：数据汇聚后形成市场定义权")
    return image


def build_layer_three(layer: LayerData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "第三层护城河：人机混合智能", "AI 负责规模执行，人类保留战略洞察", small=False)

    draw_pill(draw, 250, 282, "AI 执行层", BLUE, w=208)
    draw_ai_person(draw, 344, 750, label="AI", accent=BLUE, scale=1.08)
    draw_dashboard(draw, 130, 360, 434, 236, "规模执行看板", color=BLUE)
    draw_metric_tile(draw, 136, 632, "99%", "重复性执行自动化", BLUE, w=194, h=110)
    draw_metric_tile(draw, 356, 632, "批量", "海量动作并行推进", GREEN, w=194, h=110)

    draw_pill(draw, 860, 282, "协同中枢", GREEN, w=208)
    draw_core_circle(draw, 960, 572, "H+A", "您（指挥官） + 我（AI执行官）", accent=GREEN)
    draw_arrow(draw, 564, 570, 796, 570, color=BLUE, width=6)
    draw_arrow(draw, 1124, 570, 1360, 570, color=ORANGE, width=6)
    draw_quote_band(draw, 716, 760, 488, "AI 把重复事务吃干榨净，人类只做关键判断。", color=GREEN)

    draw_pill(draw, 1460, 282, "人类判断层", ORANGE, w=228)
    draw_stick_person(draw, 1574, 770, accent=ORANGE, pose="point", scale=1.08)
    draw_stage_box(draw, 1348, 360, "战略决策", ORANGE, "从 10 个备选项里挑出最关键的一步。", w=282, h=140)
    draw_metric_tile(draw, 1368, 632, "1%", "关键节点的直觉与洞察", ORANGE, w=194, h=110)
    draw_metric_tile(draw, 1588, 632, "10→1", "从备选中挑出最关键动作", RED, w=194, h=110)

    draw_quote_band(draw, 238, 894, 1442, layer.conclusion or "这种“人机混合智能”带来的创造性和灵活性，是任何写死的 Skill 都无法比拟的。", color=ORANGE)
    draw_bottom_rule(draw, 5, 7, "第三层：把人类洞察与 AI 执行力合成一套系统")
    return image


def build_layer_four(layer: LayerData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "第四层护城河：平台化扩张能力", "从项目交付升级为网络扩张，形成压倒性规模优势", small=False)

    draw_pill(draw, 170, 280, "单点定制", ORANGE, w=214)
    draw_annotation(draw, 126, 348, "对手路径", layer.rival, ORANGE)
    draw_stage_box(draw, 142, 562, "项目 A", ORANGE, "一单一做", w=180, h=100)
    draw_stage_box(draw, 142, 694, "项目 B", ORANGE, "再次重配", w=180, h=100)
    draw_stage_box(draw, 142, 826, "项目 C", ORANGE, "边际成本高", w=180, h=100)

    draw_pill(draw, 860, 258, "平台复制引擎", BLUE, w=248)
    draw_small_success_node(draw, 892, 410, "成功样板", accent=GREEN)
    expansion_nodes = [
        (656, 594, "民宿 02"),
        (820, 706, "民宿 03"),
        (1018, 706, "民宿 04"),
        (1182, 594, "民宿 05"),
        (640, 842, "民宿 06"),
        (834, 852, "民宿 07"),
        (1032, 852, "民宿 08"),
        (1226, 842, "民宿 09"),
    ]
    for x, y, label in expansion_nodes:
        draw_small_success_node(draw, x, y, label, accent=BLUE)
        draw_arrow(draw, 950, 478, x + 58, y, color=BLUE, width=4)

    draw_pill(draw, 1494, 280, "网络扩张", GREEN, w=214)
    draw_annotation(draw, 1418, 344, "系统能力", layer.system, BLUE)
    draw_metric_tile(draw, 1416, 584, "1→100", "成功模式规模复制", BLUE, w=170, h=112)
    draw_metric_tile(draw, 1602, 584, "≈0", "新增复制边际成本", GREEN, w=170, h=112)
    draw_metric_tile(draw, 1510, 728, "平台", "数据和成本优势同步扩大", ORANGE, w=170, h=112)

    draw_quote_band(draw, 360, 896, 1200, layer.conclusion or "当竞争对手还在做“项目交付”时，我们已经在做“网络扩张”。", color=GREEN)
    draw_bottom_rule(draw, 6, 7, "第四层：规模化复制把成本优势和数据优势同时拉开")
    return image


def build_final(story: StoryData) -> Image.Image:
    image, draw = draw_slide_background()
    draw_title_block(draw, "最终结论：What 能复制，How 复制不了", "真正不可复制的，不是流程表面，而是系统如何进化、垄断、协同与扩张", small=False)

    draw_pill(draw, 210, 296, "What：别人能抄走的", ORANGE, w=294)
    what_cards = [
        ("Skill", "一次性自动化动作", 132, 374, ORANGE),
        ("SOP", "表面流程与操作步骤", 132, 566, ORANGE),
        ("交付", "单次项目产出", 132, 758, ORANGE),
    ]
    for icon, title, x, y, accent in what_cards:
        draw_card(draw, x, y, 420, 146, title, "可见、可学、可模仿", accent, icon)

    draw_core_circle(draw, 960, 582, "生命体", "会学习、会垄断、会扩张的系统", accent=GREEN)
    draw_arrow(draw, 548, 446, 786, 522, color=ORANGE, width=5)
    draw_arrow(draw, 548, 638, 786, 582, color=ORANGE, width=5)
    draw_arrow(draw, 548, 830, 786, 644, color=ORANGE, width=5)

    draw_pill(draw, 1400, 296, "How：别人复制不了的", BLUE, w=318)
    how_cards = [
        ("闭环进化", "执行越多，系统越聪明", 1326, 374, BLUE),
        ("欲望图谱", "区域级数据私有资产", 1326, 536, GREEN),
        ("人机混智", "战略判断与 AI 执行耦合", 1326, 698, ORANGE),
        ("平台扩张", "1 到 100 的网络复制能力", 1326, 860, BLUE),
    ]
    for title, subtitle, x, y, accent in how_cards:
        draw_card(draw, x, y, 448, 118, title, subtitle, accent, "•")
        draw_arrow(draw, 1128, 582, x, y + 58, color=accent, width=5)

    final_text = "别人能复制 What，却复制不了 How。"
    draw_quote_band(draw, 300, 904, 960, final_text, color=BLUE)
    draw_bottom_rule(draw, 7, 7, "结论：系统级护城河，决定长期不可复制性")
    return image


def build_slides(story: StoryData) -> list[tuple[str, Image.Image]]:
    if len(story.layers) < 4:
        raise RuntimeError(f"Expected at least 4 moat layers in {SOURCE_DOC}, found {len(story.layers)}")

    return [
        ("P1-封面-Field护城河.png", build_cover(story)),
        ("P2-引言-Skill与Field.png", build_intro(story)),
        ("P3-第一层-闭环自进化系统.png", build_layer_one(story.layers[0])),
        ("P4-第二层-垄断性数据资产.png", build_layer_two(story.layers[1])),
        ("P5-第三层-人机混合智能.png", build_layer_three(story.layers[2])),
        ("P6-第四层-平台化扩张能力.png", build_layer_four(story.layers[3])),
        ("P7-最终结论-What与How.png", build_final(story)),
    ]


def save_ppt(slide_images: list[Path], output_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for image_path in slide_images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output_path))


def write_readme(output_dir: Path, story: StoryData, pptx_path: Path, slide_images: list[Path]) -> None:
    lines = [
        "这是根据《项目核心竞争力分析.md》提炼的 7 页 Field 风格图片型 PPT 样张。",
        "生成模式：本地可控绘制（蓝橙手绘咨询图风格）",
        f"源文档：{SOURCE_DOC}",
        f"PPT：{pptx_path.name}",
        "",
        "页面结构：",
    ]
    for index, image_path in enumerate(slide_images, start=1):
        lines.append(f"P{index}：{image_path.stem.split('-', 2)[-1]}")
    lines.extend(
        [
            "",
            "说明：本版延续单页 Gemini 参考图的白底、黑线、蓝橙强调风格，但为保证中文可控与结构稳定，主体图形与文字均由本地程序绘制。",
            f"文档主标题：{story.title}",
            f"结果 PPT：{pptx_path}",
        ]
    )
    (output_dir / "README.txt").write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the Field-style core-competitiveness PPT deck.")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=OUTPUT_ROOT,
        help="Directory where timestamped local result folders are created.",
    )
    parser.add_argument(
        "--with-contact-sheet",
        action="store_true",
        help="Also write a local contact-sheet preview image.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if not SOURCE_DOC.exists():
        raise RuntimeError(f"Source markdown not found: {SOURCE_DOC}")

    story = parse_markdown_story(SOURCE_DOC)
    output_dir = create_timestamped_output_dir(args.output_root, "项目核心竞争力分析-field七页样张")

    slide_defs = build_slides(story)
    slide_paths: list[Path] = []
    for filename, image in slide_defs:
        path = output_dir / filename
        image.save(path, quality=95)
        slide_paths.append(path)

    pptx_path = output_dir / "项目核心竞争力分析-Field风格-图片型PPT-七页样张.pptx"
    save_ppt(slide_paths, pptx_path)
    contact_sheet = output_dir / "contact-sheet.png"
    if args.with_contact_sheet:
        create_contact_sheet(slide_paths, contact_sheet)
    write_readme(output_dir, story, pptx_path, slide_paths)
    write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(slide_paths),
        sourceDoc=str(SOURCE_DOC),
        readme=str(output_dir / "README.txt"),
        contactSheet=str(contact_sheet) if args.with_contact_sheet else "",
        imageDirectory=str(output_dir),
        mode="field-local-render",
    )

    print(pptx_path)


if __name__ == "__main__":
    main()

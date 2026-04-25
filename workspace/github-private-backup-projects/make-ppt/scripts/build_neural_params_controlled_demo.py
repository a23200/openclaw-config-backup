from __future__ import annotations

import argparse
from pathlib import Path
import math

from PIL import Image, ImageDraw
from local_result import create_timestamped_output_dir, repo_runtime_dir, write_result_manifest

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import (
    WIDTH,
    HEIGHT,
    WHITE,
    BLUE,
    ORANGE,
    NAVY,
    GRAY,
    BLUE_LIGHT,
    ORANGE_LIGHT,
    load_font,
    draw_arrow,
    draw_speech,
)


TITLE_FONT = load_font(62, heavy=True)
BODY_FONT = load_font(30)
BODY_SMALL_FONT = load_font(24)
SCREEN_FONT = load_font(24, heavy=True)
SCREEN_AI_FONT = load_font(54, heavy=True)

DESKTOP = Path("/Users/mac/Desktop")
OUTPUT_ROOT = repo_runtime_dir("single-slide-demo")


def text_width(text: str, font) -> float:
    probe = ImageDraw.Draw(Image.new("RGB", (8, 8), WHITE))
    return probe.textlength(text, font=font)


def wrap_text(text: str, font, max_width: int) -> str:
    lines: list[str] = []
    current = ""
    for char in text:
        if char == "\n":
            if current:
                lines.append(current)
                current = ""
            continue
        candidate = current + char
        if not current or text_width(candidate, font) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = char
    if current:
        lines.append(current)
    return "\n".join(lines)


def draw_centered(draw: ImageDraw.ImageDraw, text: str, y: int, font, fill: str, max_width: int | None = None, spacing: int = 8) -> None:
    rendered = wrap_text(text, font, max_width) if max_width else text
    box = draw.multiline_textbbox((0, 0), rendered, font=font, spacing=spacing, align="center")
    x = (WIDTH - (box[2] - box[0])) / 2
    draw.multiline_text((x, y), rendered, font=font, fill=fill, spacing=spacing, align="center")


def draw_husky(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    body_box = (
        x - int(115 * scale),
        y - int(148 * scale),
        x + int(86 * scale),
        y + int(28 * scale),
    )
    draw.ellipse(body_box, fill="#E9EEF5", outline=NAVY, width=5)
    draw.ellipse(
        (
            x - int(52 * scale),
            y - int(112 * scale),
            x + int(42 * scale),
            y - int(4 * scale),
        ),
        fill=WHITE,
        outline=None,
    )

    head_x = x + int(82 * scale)
    head_y = y - int(162 * scale)
    head_r = int(62 * scale)
    draw.ellipse((head_x - head_r, head_y - head_r, head_x + head_r, head_y + head_r), fill="#E9EEF5", outline=NAVY, width=5)
    draw.polygon(
        [(head_x - int(42 * scale), head_y - int(42 * scale)), (head_x - int(28 * scale), head_y - int(104 * scale)), (head_x + int(4 * scale), head_y - int(44 * scale))],
        fill="#CBD5E1",
        outline=NAVY,
    )
    draw.polygon(
        [(head_x + int(14 * scale), head_y - int(42 * scale)), (head_x + int(42 * scale), head_y - int(100 * scale)), (head_x + int(56 * scale), head_y - int(30 * scale))],
        fill="#CBD5E1",
        outline=NAVY,
    )
    draw.polygon(
        [
            (head_x - int(28 * scale), head_y - int(28 * scale)),
            (head_x + int(26 * scale), head_y - int(24 * scale)),
            (head_x + int(48 * scale), head_y + int(24 * scale)),
            (head_x + int(8 * scale), head_y + int(54 * scale)),
            (head_x - int(28 * scale), head_y + int(24 * scale)),
        ],
        fill=WHITE,
        outline=None,
    )
    draw.ellipse(
        (
            head_x + int(24 * scale),
            head_y + int(2 * scale),
            head_x + int(94 * scale),
            head_y + int(48 * scale),
        ),
        fill=WHITE,
        outline=NAVY,
        width=4,
    )
    draw.ellipse((head_x - int(24 * scale), head_y - int(8 * scale), head_x - int(10 * scale), head_y + int(6 * scale)), fill=BLUE)
    draw.ellipse((head_x + int(20 * scale), head_y - int(6 * scale), head_x + int(34 * scale), head_y + int(8 * scale)), fill=BLUE)
    draw.ellipse((head_x + int(78 * scale), head_y + int(18 * scale), head_x + int(94 * scale), head_y + int(34 * scale)), fill=NAVY)
    draw.arc((head_x + int(34 * scale), head_y + int(28 * scale), head_x + int(78 * scale), head_y + int(58 * scale)), 8, 156, fill=NAVY, width=3)
    draw.line((head_x + int(56 * scale), head_y + int(24 * scale), head_x + int(58 * scale), head_y + int(42 * scale)), fill=NAVY, width=3)

    for dx in [-70, -20, 42]:
        draw.line((x + int(dx * scale), y + int(6 * scale), x + int(dx * scale), y + int(92 * scale)), fill=NAVY, width=6)
        draw.line((x + int(dx * scale), y + int(92 * scale), x + int((dx + 26) * scale), y + int(92 * scale)), fill=NAVY, width=5)

    tail_base = (x - int(104 * scale), y - int(92 * scale))
    tail_tip = (x - int(188 * scale), y - int(170 * scale))
    draw.line((tail_base, tail_tip), fill=NAVY, width=8)
    draw.arc((tail_tip[0] - int(22 * scale), tail_tip[1] - int(26 * scale), tail_tip[0] + int(72 * scale), tail_tip[1] + int(58 * scale)), 210, 46, fill=NAVY, width=8)
    for radius in [22, 42]:
        draw.arc((tail_tip[0] - radius, tail_tip[1] - radius, tail_tip[0] + radius + 42, tail_tip[1] + radius), 230, 45, fill=ORANGE, width=3)


def draw_computer_person(draw: ImageDraw.ImageDraw, cx: int, bottom_y: int, accent: str, screen_text: str, font, fill: str = NAVY) -> None:
    head_y1 = bottom_y - 356
    head_y2 = bottom_y - 226
    screen_x1 = cx - 106
    screen_y1 = head_y1 + 18
    screen_x2 = cx + 106
    screen_y2 = head_y2 - 18
    draw.rounded_rectangle((cx - 126, head_y1, cx + 126, head_y2), radius=28, fill=WHITE, outline=accent, width=6)
    draw.rounded_rectangle((screen_x1, screen_y1, screen_x2, screen_y2), radius=18, fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT, outline=accent, width=4)

    rendered = wrap_text(screen_text, font, 178)
    box = draw.multiline_textbbox((0, 0), rendered, font=font, spacing=3, align="center")
    text_x = screen_x1 + ((screen_x2 - screen_x1) - (box[2] - box[0])) / 2
    text_y = screen_y1 + ((screen_y2 - screen_y1) - (box[3] - box[1])) / 2 - 2
    draw.multiline_text((text_x, text_y), rendered, font=font, fill=fill, spacing=3, align="center")

    body_top = head_y2
    hip_y = bottom_y - 82
    draw.line((cx, body_top, cx, hip_y), fill=NAVY, width=8)
    draw.line((cx, body_top + 52, cx - 58, body_top + 84), fill=NAVY, width=7)
    draw.line((cx, body_top + 52, cx + 58, body_top + 84), fill=NAVY, width=7)
    draw.line((cx, hip_y, cx - 58, bottom_y), fill=NAVY, width=8)
    draw.line((cx, hip_y, cx + 58, bottom_y), fill=NAVY, width=8)


def build_slide() -> Image.Image:
    image = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw = ImageDraw.Draw(image)

    draw_centered(draw, "神经网络学习的重要参数", 74, TITLE_FONT, BLUE)
    draw_centered(
        draw,
        "• 激活函数：决定神经元是否“激活”\n• 常见函数：ReLU，Sigmoid，Tanh",
        198,
        BODY_FONT,
        NAVY,
        max_width=1080,
        spacing=14,
    )

    draw_husky(draw, 260, 760, scale=1.0)
    draw_computer_person(draw, 720, 912, BLUE, "耳朵*权重+\n颜色*权重+\n...=2.7", SCREEN_FONT)
    draw_computer_person(draw, 1080, 912, ORANGE, "激活函数*\n(2.7)=1", SCREEN_FONT)
    draw_computer_person(draw, 1440, 912, BLUE, "AI", SCREEN_AI_FONT, fill=BLUE)

    draw_speech(draw, 1376, 420, "这是狗", outline=ORANGE, max_width=140)

    draw_arrow(draw, 440, 694, 560, 694, color=ORANGE, width=6)
    draw_arrow(draw, 854, 694, 946, 694, color=ORANGE, width=6)
    draw_arrow(draw, 1214, 694, 1306, 694, color=ORANGE, width=6)

    return image


def save_ppt(image_path: Path, ppt_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(ppt_path))


def save_prompt(output_path: Path) -> None:
    prompt = "\n".join(
        [
            "严格生图提示词（适合 Gemini / Nano Banana 一类模型）：",
            "",
            "请生成 16:9 横版教学信息图，纯白背景，火柴人动画风格，线条简洁，黑色线稿，蓝色和橙色作为唯一强调色，整体干净清晰。",
            "非常重要：只保留我明确指定的元素，不要增加任何额外角色、图标、装饰、齿轮、书本、山、表格、背景场景、边框花纹、公式、英文、数字、水印。",
            "",
            "布局要求：",
            "1）画面上半部分留白，供后续叠加标题和说明文字。",
            "2）画面下半部分从左到右仅包含四个主体，并保持同一水平线排布：",
            "   - 一只坐在地上、摇尾巴的哈士奇；",
            "   - 一个正面朝向的计算机头火柴人，屏幕留白，供后续叠字；",
            "   - 第二个正面朝向的计算机头火柴人，屏幕留白，供后续叠字；",
            "   - 第三个正面朝向的计算机头火柴人，屏幕中仅保留一个简洁 AI 标识或留白；其头顶上方保留一个空白对话框，供后续叠字。",
            "3）四个主体之间仅用橙色箭头连接。",
            "",
            "硬性限制：",
            "- 不要生成任何可读文字；",
            "- 不要生成任何公式；",
            "- 不要生成任何多余背景；",
            "- 画面必须简洁；",
            "- 所有主体都要完整显示，不要裁切。",
            "",
            "说明：这类页建议只让模型负责“无文字主体构图”，所有中文标题、公式、对话框内容都由本地程序后期叠加，避免错字和乱加元素。",
        ]
    )
    output_path.write_text(prompt, encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the controlled neural-parameters single-slide PPT locally.")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=OUTPUT_ROOT,
        help="Directory where timestamped local result folders are created.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_dir = create_timestamped_output_dir(args.output_root, "神经网络学习-严格版样张")
    out_png = output_dir / "神经网络学习的重要参数-严格版样张.png"
    out_pptx = output_dir / "神经网络学习的重要参数-严格版样张.pptx"
    out_prompt = output_dir / "神经网络学习-严格生图提示词.txt"
    image = build_slide()
    image.save(out_png, quality=95)
    save_ppt(out_png, out_pptx)
    save_prompt(out_prompt)
    write_result_manifest(
        output_dir,
        out_pptx,
        previewImage=str(out_png),
        promptFile=str(out_prompt),
        mode="controlled-local-render",
    )
    print(out_pptx)


if __name__ == "__main__":
    main()

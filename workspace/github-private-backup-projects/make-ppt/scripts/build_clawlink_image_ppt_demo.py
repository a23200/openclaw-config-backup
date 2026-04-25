from __future__ import annotations

import argparse
import base64
import json
import math
import os
import subprocess
import textwrap
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    load_env_candidates,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None


WIDTH = 1920
HEIGHT = 1080

WHITE = "#FFFFFF"
NAVY = "#0F172A"
INK = "#111827"
BLUE = "#1F6FB2"
BLUE_BRIGHT = "#2F7FD1"
BLUE_DEEP = "#154B82"
BLUE_LIGHT = "#DCEEFE"
BLUE_PALE = "#F0F7FF"
ORANGE = "#F29A2E"
ORANGE_DEEP = "#C76C13"
ORANGE_LIGHT = "#FDE9CF"
GRAY = "#64748B"
GRAY_LIGHT = "#CBD5E1"
GRAY_PALE = "#F1F5F9"
GREEN = "#69C27D"
GREEN_LIGHT = "#E8F8EC"
RED = "#EF4444"
RED_LIGHT = "#FEE2E2"

FONT_SANS = "/System/Library/Fonts/Hiragino Sans GB.ttc"
FONT_HEAVY = "/System/Library/Fonts/STHeiti Medium.ttc"

DEFAULT_OUTPUT_ROOT = repo_runtime_dir("image-ppt-demo")
DOC_PATH = Path("/Users/mac/Desktop/文旅引爆点项目方案.txt")
DEFAULT_GEMINI_IMAGE_MODEL = "gemini-3-pro-image-preview"
FAST_GEMINI_IMAGE_MODEL = "gemini-3.1-flash-image-preview"
ARCHIVE_DIR = latest_archive_dir()
CLAWLINK_LAYOUT_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "文旅引爆点-Gemini单页样张-P2.png" if ARCHIVE_DIR else None,
    ARCHIVE_DIR / "runtime-image-ppt-demo" / "20260414-232000-文旅引爆点-clawlink六页增强样张" / "P2-机会判断-系统增长.png" if ARCHIVE_DIR else None,
)
CLAWLINK_STYLE_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    CLAWLINK_LAYOUT_REFERENCE,
)
CLAWLINK_PROMPT_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "runtime-image-ppt-demo" / "20260414-002010-文旅引爆点-火柴人样张" / "prompt.txt" if ARCHIVE_DIR else None,
)


def load_font(size: int, heavy: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    try:
        return ImageFont.truetype(FONT_HEAVY if heavy else FONT_SANS, size=size)
    except Exception:
        return ImageFont.load_default()


TITLE_FONT = load_font(58, heavy=True)
TITLE_SMALL_FONT = load_font(50, heavy=True)
SUBTITLE_FONT = load_font(26)
BODY_FONT = load_font(28)
BODY_BOLD_FONT = load_font(30, heavy=True)
SECTION_FONT = load_font(34, heavy=True)
SMALL_FONT = load_font(22)
SMALL_BOLD_FONT = load_font(24, heavy=True)
TINY_FONT = load_font(18)
TINY_BOLD_FONT = load_font(19, heavy=True)
LABEL_FONT = load_font(20, heavy=True)
AI_FONT = load_font(42, heavy=True)
HUGE_FONT = load_font(82, heavy=True)


@dataclass(frozen=True)
class SlideSpec:
    title: str
    subtitle: str
    filename: str
    scene_prompt: str
    builder: Callable[[Image.Image | None], Image.Image]


@dataclass(frozen=True)
class GeminiConfig:
    api_key: str
    model: str
    base_url: str


def gemini_config_from_env() -> GeminiConfig | None:
    load_env_candidates()
    api_key = (
        os.environ.get("GEMINI_API_KEY")
        or os.environ.get("GOOGLE_API_KEY")
        or os.environ.get("GOOGLE_GENAI_API_KEY")
    )
    if not api_key:
        return None
    return GeminiConfig(
        api_key=api_key,
        model=os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL),
        base_url=os.environ.get("GEMINI_IMAGE_BASE_URL", "https://generativelanguage.googleapis.com/v1beta"),
    )


def load_reference_prompt_text() -> str:
    if not CLAWLINK_PROMPT_REFERENCE or not CLAWLINK_PROMPT_REFERENCE.exists():
        return ""
    return CLAWLINK_PROMPT_REFERENCE.read_text(encoding="utf-8").strip()


def reference_image_paths() -> list[Path]:
    paths: list[Path] = []
    for candidate in [CLAWLINK_LAYOUT_REFERENCE, CLAWLINK_STYLE_REFERENCE]:
        if candidate and candidate.exists() and candidate not in paths:
            paths.append(candidate)
    return paths


def mime_type_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    return "application/octet-stream"


def text_width(text: str, font: ImageFont.ImageFont) -> float:
    probe = ImageDraw.Draw(Image.new("RGB", (8, 8), WHITE))
    return probe.textlength(text, font=font)


def wrap_text(text: str, font: ImageFont.ImageFont, max_width: int) -> str:
    lines: list[str] = []
    current = ""
    for char in text:
        if char == "\n":
            if current:
                lines.append(current)
                current = ""
            continue
        candidate = f"{current}{char}"
        if text_width(candidate, font) <= max_width or not current:
            current = candidate
        else:
            lines.append(current)
            current = char
    if current:
        lines.append(current)
    return "\n".join(lines)


def draw_centered_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    y: int,
    font: ImageFont.ImageFont,
    fill: str,
    max_width: int | None = None,
    spacing: int = 10,
) -> None:
    rendered = wrap_text(text, font, max_width) if max_width else text
    box = draw.multiline_textbbox((0, 0), rendered, font=font, spacing=spacing, align="center")
    x = (WIDTH - (box[2] - box[0])) / 2
    draw.multiline_text((x, y), rendered, font=font, fill=fill, spacing=spacing, align="center")


def draw_text_in_box(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.ImageFont,
    fill: str,
    max_width: int,
    spacing: int = 8,
) -> int:
    rendered = wrap_text(text, font, max_width)
    draw.multiline_text(xy, rendered, font=font, fill=fill, spacing=spacing)
    bbox = draw.multiline_textbbox(xy, rendered, font=font, spacing=spacing)
    return bbox[3]


def resize_cover(image: Image.Image, width: int = WIDTH, height: int = HEIGHT) -> Image.Image:
    image = image.convert("RGB")
    scale = max(width / image.width, height / image.height)
    resized = image.resize((int(image.width * scale), int(image.height * scale)))
    left = (resized.width - width) // 2
    top = (resized.height - height) // 2
    return resized.crop((left, top, left + width, top + height))


def draw_sketch_line(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    fill: str = INK,
    width: int = 3,
    repeats: int = 2,
) -> None:
    x1, y1, x2, y2 = xy
    for index in range(repeats):
        offset = index - repeats // 2
        draw.line((x1 + offset, y1 - offset, x2 + offset, y2 + offset), fill=fill, width=width)


def draw_sketch_round_rect(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    radius: int,
    fill: str,
    outline: str,
    width: int = 4,
    repeats: int = 2,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)
    if repeats > 1:
        x1, y1, x2, y2 = xy
        draw.rounded_rectangle((x1 + 3, y1 - 2, x2 - 2, y2 + 3), radius=radius, outline=outline, width=1)


def draw_slide_background(backdrop: Image.Image | None = None) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    if backdrop:
        image = resize_cover(backdrop)
        white_layer = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
        image = Image.blend(image, white_layer, 0.68).filter(ImageFilter.SMOOTH)
    else:
        image = Image.new("RGB", (WIDTH, HEIGHT), WHITE)

    draw = ImageDraw.Draw(image)

    for x in range(110, WIDTH - 80, 110):
        for y in range(184, HEIGHT - 118, 88):
            if (x + y) % 3 == 0:
                draw.ellipse((x - 2, y - 2, x + 2, y + 2), fill="#E2E8F0")

    for x, y, color in [
        (78, 118, BLUE_LIGHT),
        (1718, 106, ORANGE_LIGHT),
        (150, 884, ORANGE_LIGHT),
        (1710, 840, BLUE_LIGHT),
    ]:
        draw.rounded_rectangle((x, y, x + 128, y + 58), radius=28, fill=color, outline=None)

    draw.arc((62, 226, 300, 420), 190, 352, fill="#DBEAFE", width=4)
    draw.arc((1632, 226, 1864, 426), 18, 168, fill="#FED7AA", width=4)
    draw.arc((1548, 772, 1838, 1014), 190, 350, fill="#DBEAFE", width=4)
    return image, draw


def draw_title_block(draw: ImageDraw.ImageDraw, title: str, subtitle: str, small: bool = False) -> None:
    draw_centered_text(draw, title, 52, TITLE_SMALL_FONT if small else TITLE_FONT, BLUE, max_width=1460)
    draw_centered_text(draw, subtitle, 132, SUBTITLE_FONT, GRAY, max_width=1380)
    draw.rounded_rectangle((748, 196, 1172, 206), radius=5, fill=ORANGE, outline=None)
    draw.rounded_rectangle((810, 214, 1110, 220), radius=3, fill=BLUE_LIGHT, outline=None)


def draw_bottom_rule(draw: ImageDraw.ImageDraw, page: int, total: int, caption: str) -> None:
    y = 1000
    draw.line((110, y, 1810, y), fill=GRAY_LIGHT, width=3)
    for x, color in [(190, BLUE_BRIGHT), (294, ORANGE), (424, BLUE), (1640, GREEN)]:
        draw.ellipse((x - 10, y - 10, x + 10, y + 10), fill=color, outline=NAVY, width=2)
    for x in [150, 232, 366, 1510, 1720]:
        draw.line((x, y - 12, x + 26, y + 12), fill=NAVY, width=3)
    draw.text((118, 1018), caption, font=TINY_BOLD_FONT, fill=GRAY)
    draw.text((1718, 1018), f"{page:02d}/{total:02d}", font=TINY_BOLD_FONT, fill=BLUE)


def draw_bulb(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    r = 14 * scale
    draw.ellipse((x - r, y - r, x + r, y + r), outline=BLUE, width=max(2, int(3 * scale)))
    draw.line((x, y + r, x, y + r + 14 * scale), fill=ORANGE, width=max(2, int(3 * scale)))
    draw.line((x - 8 * scale, y + r + 14 * scale, x + 8 * scale, y + r + 14 * scale), fill=NAVY, width=max(2, int(3 * scale)))
    for dx, dy in [(-20, -18), (0, -24), (20, -18)]:
        draw.line(
            (x + dx * scale, y + dy * scale, x + dx * scale * 1.1, y + dy * scale * 1.1 - 10 * scale),
            fill=ORANGE,
            width=max(2, int(2 * scale)),
        )


def draw_gear(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    r = 16 * scale
    for angle in range(0, 360, 60):
        rad = math.radians(angle)
        x1 = x + math.cos(rad) * (r + 2 * scale)
        y1 = y + math.sin(rad) * (r + 2 * scale)
        x2 = x + math.cos(rad) * (r + 10 * scale)
        y2 = y + math.sin(rad) * (r + 10 * scale)
        draw.line((x1, y1, x2, y2), fill=GRAY, width=max(2, int(4 * scale)))
    draw.ellipse((x - r, y - r, x + r, y + r), outline=GRAY, width=max(2, int(3 * scale)))
    draw.ellipse((x - 5 * scale, y - 5 * scale, x + 5 * scale, y + 5 * scale), fill=WHITE, outline=GRAY, width=max(2, int(2 * scale)))


def draw_arrow(
    draw: ImageDraw.ImageDraw,
    x1: int,
    y1: int,
    x2: int,
    y2: int,
    color: str = ORANGE,
    width: int = 7,
    dashed: bool = False,
) -> None:
    if dashed:
        segments = 16
        for index in range(segments):
            if index % 2 == 0:
                sx = x1 + (x2 - x1) * index / segments
                sy = y1 + (y2 - y1) * index / segments
                ex = x1 + (x2 - x1) * (index + 1) / segments
                ey = y1 + (y2 - y1) * (index + 1) / segments
                draw.line((sx, sy, ex, ey), fill=color, width=width)
    else:
        draw_sketch_line(draw, (x1, y1, x2, y2), fill=color, width=width, repeats=2)
    angle = math.atan2(y2 - y1, x2 - x1)
    head = 22
    p1 = (x2 + head * math.cos(angle + math.pi * 0.84), y2 + head * math.sin(angle + math.pi * 0.84))
    p2 = (x2 + head * math.cos(angle - math.pi * 0.84), y2 + head * math.sin(angle - math.pi * 0.84))
    draw.polygon([(x2, y2), p1, p2], fill=color)


def draw_chip(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, color: str = BLUE) -> None:
    padding_x = 18
    box = draw.textbbox((0, 0), text, font=LABEL_FONT)
    width = box[2] - box[0] + padding_x * 2
    fill = BLUE_LIGHT if color == BLUE else ORANGE_LIGHT if color == ORANGE else GREEN_LIGHT
    draw.rounded_rectangle((x, y, x + width, y + 42), radius=21, fill=fill, outline=color, width=3)
    draw.text((x + padding_x, y + 8), text, font=LABEL_FONT, fill=color)


def draw_annotation(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, body: str, color: str) -> None:
    draw_sketch_round_rect(draw, (x, y, x + 280, y + 116), radius=24, fill=WHITE, outline=color, width=4)
    draw.text((x + 22, y + 16), title, font=SMALL_BOLD_FONT, fill=color)
    draw_text_in_box(draw, (x + 22, y + 50), body, SMALL_FONT, NAVY, 232, spacing=5)


def draw_metric_tile(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    number: str,
    label: str,
    color: str,
    w: int = 186,
    h: int = 118,
) -> None:
    fill = BLUE_PALE if color == BLUE else ORANGE_LIGHT if color == ORANGE else GREEN_LIGHT
    draw_sketch_round_rect(draw, (x, y, x + w, y + h), radius=22, fill=fill, outline=color, width=3)
    draw.text((x + 22, y + 18), number, font=SECTION_FONT, fill=color)
    draw_text_in_box(draw, (x + 22, y + 68), label, TINY_BOLD_FONT, NAVY, w - 44, spacing=3)


def draw_chart(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    sx = lambda value: x + int(value * scale)
    sy = lambda value: y + int(value * scale)
    draw.line((sx(0), sy(0), sx(0), sy(-58)), fill=GRAY, width=max(2, int(2 * scale)))
    draw.line((sx(0), sy(0), sx(100), sy(0)), fill=GRAY, width=max(2, int(2 * scale)))
    bars = [20, 34, 48, 38]
    for idx, height in enumerate(bars):
        bx = sx(12 + idx * 20)
        draw.rectangle((bx, sy(-height), bx + int(11 * scale), sy(0)), fill=BLUE_BRIGHT if idx % 2 == 0 else ORANGE, outline=BLUE)
    draw_arrow(draw, sx(18), sy(-20), sx(90), sy(-76), color=BLUE, width=max(3, int(4 * scale)))


def draw_cloud(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    circles = [
        (x - 36 * scale, y, 26 * scale),
        (x - 4 * scale, y - 16 * scale, 30 * scale),
        (x + 28 * scale, y, 24 * scale),
    ]
    for cx, cy, r in circles:
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=WHITE, outline=GRAY, width=max(2, int(3 * scale)))
    draw.rounded_rectangle((x - 60 * scale, y, x + 56 * scale, y + 24 * scale), radius=12, fill=WHITE, outline=GRAY, width=max(2, int(3 * scale)))


def draw_mountains(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    def p(dx: int, dy: int) -> tuple[int, int]:
        return (x + int(dx * scale), y + int(dy * scale))

    draw.polygon([p(0, 120), p(80, 0), p(160, 120)], outline=BLUE, fill="#EDF6FF")
    draw.polygon([p(90, 130), p(180, 24), p(280, 130)], outline=BLUE, fill="#F4FAFF")
    draw.line((p(60, 120), p(260, 120)), fill=GRAY, width=max(2, int(3 * scale)))
    draw.arc((x + int(28 * scale), y + int(90 * scale), x + int(206 * scale), y + int(168 * scale)), 0, 180, fill=ORANGE, width=max(2, int(3 * scale)))


def draw_inn(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    def p(dx: int, dy: int) -> tuple[int, int]:
        return (x + int(dx * scale), y + int(dy * scale))

    draw.rectangle((p(40, 80), p(220, 240)), fill="#F8FBFF", outline=BLUE, width=max(3, int(5 * scale)))
    draw.polygon([p(20, 100), p(130, 28), p(240, 100)], fill=ORANGE_LIGHT, outline=ORANGE)
    draw.rectangle((p(110, 160), p(150, 240)), fill=WHITE, outline=BLUE, width=max(2, int(4 * scale)))
    for wx in [70, 165]:
        draw.rectangle((p(wx, 130), p(wx + 28, 158)), fill=BLUE_LIGHT, outline=BLUE, width=max(2, int(3 * scale)))
    draw.text(p(82, 56), "民宿", font=SMALL_BOLD_FONT, fill=BLUE)


def draw_social_card(draw: ImageDraw.ImageDraw, x: int, y: int, label: str, color: str) -> None:
    draw_sketch_round_rect(draw, (x, y, x + 164, y + 98), radius=18, fill=WHITE, outline=color, width=4)
    draw.ellipse((x + 16, y + 16, x + 38, y + 38), fill=color)
    draw.rectangle((x + 48, y + 18, x + 126, y + 27), fill=GRAY_LIGHT)
    draw.rectangle((x + 48, y + 36, x + 110, y + 43), fill=GRAY_LIGHT)
    draw.rectangle((x + 18, y + 58, x + 146, y + 68), fill=BLUE_LIGHT if color == BLUE else ORANGE_LIGHT)
    draw.rectangle((x + 18, y + 78, x + 126, y + 86), fill=BLUE_LIGHT if color == BLUE else ORANGE_LIGHT)
    draw.text((x + 104, y + 12), label, font=LABEL_FONT, fill=color)


def draw_joint_line(
    draw: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    fill: str = NAVY,
    width: int = 7,
    joint_radius: int = 4,
) -> None:
    draw.line((*start, *end), fill=fill, width=width)
    for center_x, center_y in [start, end]:
        draw.ellipse(
            (center_x - joint_radius, center_y - joint_radius, center_x + joint_radius, center_y + joint_radius),
            fill=fill,
        )


def draw_stick_person(draw: ImageDraw.ImageDraw, x: int, y: int, accent: str = BLUE, pose: str = "stand", scale: float = 1.0) -> None:
    stroke_width = max(4, int(7 * scale))
    joint_radius = max(3, int(4 * scale))
    head_radius = int(31 * scale)
    body_width = int(46 * scale)
    shoulder_y = y - int(18 * scale)
    hip_y = y + int(70 * scale)
    neck_y = y - int(48 * scale)
    shoe_y = y + int(126 * scale)

    shadow_box = (
        x - int(42 * scale),
        shoe_y + int(8 * scale),
        x + int(42 * scale),
        shoe_y + int(24 * scale),
    )
    draw.ellipse(shadow_box, fill="#DDE5EE")
    draw.ellipse(
        (x - head_radius, y - int(122 * scale), x + head_radius, y - int(60 * scale)),
        outline=NAVY,
        width=stroke_width,
        fill=WHITE,
    )
    eye_radius = max(2, int(3 * scale))
    for eye_x in [x - int(10 * scale), x + int(10 * scale)]:
        draw.ellipse((eye_x - eye_radius, y - int(96 * scale), eye_x + eye_radius, y - int(90 * scale)), fill=NAVY)
    draw.arc(
        (x - int(12 * scale), y - int(94 * scale), x + int(12 * scale), y - int(76 * scale)),
        20,
        160,
        fill=NAVY,
        width=max(2, int(2 * scale)),
    )
    draw.rounded_rectangle(
        (x - body_width, neck_y, x + body_width, hip_y),
        radius=int(18 * scale),
        fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT,
        outline=NAVY,
        width=stroke_width,
    )
    draw.rounded_rectangle(
        (x - int(25 * scale), neck_y + int(18 * scale), x + int(25 * scale), hip_y - int(12 * scale)),
        radius=int(12 * scale),
        fill=accent,
        outline=None,
    )
    draw.line((x, hip_y, x, hip_y + int(10 * scale)), fill=NAVY, width=stroke_width)

    if pose == "point":
        draw_joint_line(draw, (x - body_width, shoulder_y), (x - int(78 * scale), shoulder_y + int(30 * scale)), width=stroke_width, joint_radius=joint_radius)
        draw_joint_line(draw, (x + body_width, shoulder_y), (x + int(80 * scale), shoulder_y - int(42 * scale)), width=stroke_width, joint_radius=joint_radius)
        draw.line((x + int(80 * scale), shoulder_y - int(42 * scale), x + int(104 * scale), shoulder_y - int(54 * scale)), fill=NAVY, width=max(3, int(5 * scale)))
    elif pose == "walk":
        draw_joint_line(draw, (x - body_width, shoulder_y), (x - int(80 * scale), shoulder_y - int(18 * scale)), width=stroke_width, joint_radius=joint_radius)
        draw_joint_line(draw, (x + body_width, shoulder_y), (x + int(76 * scale), shoulder_y + int(28 * scale)), width=stroke_width, joint_radius=joint_radius)
    else:
        draw_joint_line(draw, (x - body_width, shoulder_y), (x - int(74 * scale), shoulder_y + int(28 * scale)), width=stroke_width, joint_radius=joint_radius)
        draw_joint_line(draw, (x + body_width, shoulder_y), (x + int(74 * scale), shoulder_y + int(28 * scale)), width=stroke_width, joint_radius=joint_radius)
    if pose == "walk":
        draw_joint_line(draw, (x - int(18 * scale), hip_y), (x - int(68 * scale), shoe_y), width=stroke_width, joint_radius=joint_radius)
        draw_joint_line(draw, (x + int(18 * scale), hip_y), (x + int(70 * scale), shoe_y - int(16 * scale)), width=stroke_width, joint_radius=joint_radius)
    else:
        draw_joint_line(draw, (x - int(18 * scale), hip_y), (x - int(54 * scale), shoe_y), width=stroke_width, joint_radius=joint_radius)
        draw_joint_line(draw, (x + int(18 * scale), hip_y), (x + int(54 * scale), shoe_y), width=stroke_width, joint_radius=joint_radius)
    for foot_x in [x - int(54 * scale), x + int(54 * scale)]:
        draw.rounded_rectangle(
            (foot_x - int(20 * scale), shoe_y - int(4 * scale), foot_x + int(20 * scale), shoe_y + int(10 * scale)),
            radius=int(7 * scale),
            fill=NAVY,
            outline=None,
        )


def draw_ai_person(draw: ImageDraw.ImageDraw, x: int, y: int, label: str = "AI", accent: str = BLUE, scale: float = 1.0) -> None:
    stroke_width = max(4, int(7 * scale))
    joint_radius = max(3, int(4 * scale))
    head_width = int(80 * scale)
    head_height = int(56 * scale)
    body_top = y - int(36 * scale)
    body_bottom = y + int(88 * scale)
    draw.ellipse(
        (x - int(42 * scale), y + int(150 * scale), x + int(42 * scale), y + int(166 * scale)),
        fill="#DDE5EE",
    )
    draw.rounded_rectangle(
        (x - head_width, y - int(136 * scale), x + head_width, y - int(28 * scale)),
        radius=int(24 * scale),
        fill=WHITE,
        outline=NAVY,
        width=stroke_width,
    )
    draw.rounded_rectangle(
        (x - int(62 * scale), y - int(116 * scale), x + int(62 * scale), y - int(48 * scale)),
        radius=int(16 * scale),
        fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT,
        outline=accent,
        width=max(3, int(4 * scale)),
    )
    for pin_x in [x - int(92 * scale), x + int(92 * scale)]:
        draw.line((pin_x, y - int(82 * scale), pin_x + (8 if pin_x < x else -8), y - int(82 * scale)), fill=NAVY, width=max(2, int(3 * scale)))
    font = load_font(max(20, int(42 * scale)), heavy=True)
    box = draw.textbbox((0, 0), label, font=font)
    draw.text((x - (box[2] - box[0]) / 2, y - int(104 * scale)), label, font=font, fill=NAVY)
    draw.rounded_rectangle(
        (x - int(38 * scale), body_top, x + int(38 * scale), body_bottom),
        radius=int(18 * scale),
        fill=WHITE,
        outline=NAVY,
        width=stroke_width,
    )
    draw.rounded_rectangle(
        (x - int(22 * scale), body_top + int(16 * scale), x + int(22 * scale), body_bottom - int(16 * scale)),
        radius=int(10 * scale),
        fill=accent,
        outline=None,
    )
    draw_joint_line(draw, (x - int(38 * scale), y + int(8 * scale)), (x - int(76 * scale), y + int(28 * scale)), width=stroke_width, joint_radius=joint_radius)
    draw_joint_line(draw, (x + int(38 * scale), y + int(8 * scale)), (x + int(76 * scale), y + int(28 * scale)), width=stroke_width, joint_radius=joint_radius)
    draw_joint_line(draw, (x - int(18 * scale), body_bottom), (x - int(52 * scale), y + int(146 * scale)), width=stroke_width, joint_radius=joint_radius)
    draw_joint_line(draw, (x + int(18 * scale), body_bottom), (x + int(52 * scale), y + int(146 * scale)), width=stroke_width, joint_radius=joint_radius)
    for foot_x in [x - int(52 * scale), x + int(52 * scale)]:
        draw.rounded_rectangle(
            (foot_x - int(20 * scale), y + int(142 * scale), foot_x + int(20 * scale), y + int(156 * scale)),
            radius=int(7 * scale),
            fill=NAVY,
            outline=None,
        )


def draw_speech(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, outline: str = BLUE, max_width: int = 260) -> None:
    rendered = wrap_text(text, SMALL_BOLD_FONT, max_width)
    box = draw.multiline_textbbox((0, 0), rendered, font=SMALL_BOLD_FONT, spacing=5)
    width = (box[2] - box[0]) + 34
    height = (box[3] - box[1]) + 26
    draw.rounded_rectangle((x, y, x + width, y + height), radius=18, fill=WHITE, outline=outline, width=4)
    draw.polygon([(x + 24, y + height), (x + 42, y + height), (x + 31, y + height + 18)], fill=WHITE, outline=outline)
    draw.multiline_text((x + 17, y + 12), rendered, font=SMALL_BOLD_FONT, fill=NAVY, spacing=5)


def draw_clipboard(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, lines: list[str], color: str = BLUE) -> None:
    draw_sketch_round_rect(draw, (x, y, x + 308, y + 318), radius=24, fill=WHITE, outline=color, width=5)
    draw.rounded_rectangle((x + 100, y - 18, x + 204, y + 30), radius=12, fill=ORANGE_LIGHT, outline=ORANGE, width=4)
    draw.text((x + 54, y + 26), title, font=BODY_BOLD_FONT, fill=color)
    current_y = y + 88
    for line in lines:
        draw.ellipse((x + 28, current_y + 8, x + 40, current_y + 20), fill=ORANGE)
        draw_text_in_box(draw, (x + 54, current_y), line, SMALL_BOLD_FONT, NAVY, 220, spacing=4)
        current_y += 58


def draw_money_tag(draw: ImageDraw.ImageDraw, x: int, y: int, text: str) -> None:
    draw.rounded_rectangle((x, y, x + 128, y + 66), radius=18, fill=ORANGE_LIGHT, outline=ORANGE, width=4)
    draw.text((x + 22, y + 16), text, font=BODY_BOLD_FONT, fill=ORANGE_DEEP)


def draw_card(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, title: str, subtitle: str, accent: str, icon: str) -> None:
    draw_sketch_round_rect(draw, (x, y, x + w, y + h), radius=26, fill=WHITE, outline=accent, width=4)
    draw.rounded_rectangle((x + 18, y + 18, x + 94, y + 94), radius=18, fill=BLUE_LIGHT if accent == BLUE else ORANGE_LIGHT, outline=accent, width=3)
    draw.text((x + 38, y + 36), icon, font=BODY_BOLD_FONT, fill=accent)
    draw.text((x + 116, y + 24), title, font=BODY_BOLD_FONT, fill=BLUE)
    draw_text_in_box(draw, (x + 116, y + 70), subtitle, SMALL_FONT, NAVY, w - 142, spacing=7)


def draw_phone(draw: ImageDraw.ImageDraw, x: int, y: int, label: str, color: str = BLUE) -> None:
    draw.rounded_rectangle((x, y, x + 118, y + 210), radius=24, fill=WHITE, outline=NAVY, width=4)
    draw.rounded_rectangle((x + 12, y + 24, x + 106, y + 184), radius=12, fill=BLUE_PALE if color == BLUE else ORANGE_LIGHT, outline=color, width=3)
    draw.ellipse((x + 50, y + 188, x + 68, y + 206), fill=WHITE, outline=NAVY, width=2)
    draw.text((x + 28, y + 44), label, font=LABEL_FONT, fill=color)
    for index, line_width in enumerate([60, 78, 48, 68]):
        draw.rectangle((x + 28, y + 82 + index * 24, x + 28 + line_width, y + 92 + index * 24), fill=WHITE, outline=None)


def draw_dashboard(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, title: str, color: str = BLUE) -> None:
    draw_sketch_round_rect(draw, (x, y, x + w, y + h), radius=24, fill=WHITE, outline=color, width=4)
    draw.rounded_rectangle((x + 16, y + 16, x + w - 16, y + 58), radius=16, fill=BLUE_PALE if color == BLUE else ORANGE_LIGHT, outline=None)
    draw.text((x + 34, y + 24), title, font=SMALL_BOLD_FONT, fill=color)
    draw_chart(draw, x + 44, y + 176, scale=1.35)
    for idx, value in enumerate(["入住率", "均价", "热度"]):
        draw.rounded_rectangle((x + 230, y + 84 + idx * 60, x + w - 34, y + 124 + idx * 60), radius=18, fill=GRAY_PALE, outline=GRAY_LIGHT, width=2)
        draw.text((x + 248, y + 92 + idx * 60), value, font=TINY_BOLD_FONT, fill=NAVY)
        draw.rectangle((x + 328, y + 100 + idx * 60, x + 450 + idx * 30, y + 112 + idx * 60), fill=color)


def draw_megaphone(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0) -> None:
    draw.polygon(
        [
            (x, y),
            (x + int(128 * scale), y - int(58 * scale)),
            (x + int(128 * scale), y + int(58 * scale)),
        ],
        fill=ORANGE_LIGHT,
        outline=ORANGE,
    )
    draw.rounded_rectangle((x - int(54 * scale), y - int(26 * scale), x + int(8 * scale), y + int(26 * scale)), radius=int(12 * scale), fill=WHITE, outline=ORANGE, width=max(3, int(4 * scale)))
    draw.line((x - int(36 * scale), y + int(28 * scale), x - int(18 * scale), y + int(92 * scale)), fill=NAVY, width=max(3, int(5 * scale)))
    for index, radius in enumerate([78, 118, 158]):
        draw.arc((x + int(radius * scale), y - int(radius * scale), x + int((radius + 120) * scale), y + int(radius * scale)), 320, 40, fill=BLUE if index % 2 == 0 else ORANGE, width=max(2, int(4 * scale)))


def draw_magnifier(draw: ImageDraw.ImageDraw, x: int, y: int, scale: float = 1.0, color: str = BLUE) -> None:
    r = int(42 * scale)
    draw.ellipse((x - r, y - r, x + r, y + r), fill=WHITE, outline=color, width=max(3, int(5 * scale)))
    draw.line((x + int(32 * scale), y + int(32 * scale), x + int(78 * scale), y + int(78 * scale)), fill=NAVY, width=max(3, int(7 * scale)))
    draw.text((x - int(20 * scale), y - int(18 * scale)), "数", font=load_font(max(18, int(30 * scale)), heavy=True), fill=color)


def draw_review_bubble(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, color: str, tag: str) -> None:
    draw.rounded_rectangle((x, y, x + 280, y + 92), radius=24, fill=WHITE, outline=color, width=4)
    draw.ellipse((x + 18, y + 18, x + 48, y + 48), fill=color)
    draw.text((x + 62, y + 16), tag, font=LABEL_FONT, fill=color)
    draw_text_in_box(draw, (x + 62, y + 48), text, TINY_BOLD_FONT, NAVY, 190, spacing=2)


def draw_funnel(draw: ImageDraw.ImageDraw, x: int, y: int) -> None:
    levels = [
        ("社交数据", BLUE, 0),
        ("兴趣标签", ORANGE, 42),
        ("体验动作", GREEN, 84),
    ]
    for label, color, offset in levels:
        draw.polygon(
            [
                (x + offset, y + offset),
                (x + 360 - offset, y + offset),
                (x + 292 - offset // 2, y + offset + 56),
                (x + 68 + offset // 2, y + offset + 56),
            ],
            fill=BLUE_PALE if color == BLUE else ORANGE_LIGHT if color == ORANGE else GREEN_LIGHT,
            outline=color,
        )
        draw.text((x + 142, y + offset + 14), label, font=SMALL_BOLD_FONT, fill=color)


def draw_cover_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "文旅引爆点——AI合伙人运营系统", "大理站｜把潜力民宿打造成“引爆点”")

    draw_chip(draw, 118, 238, "零前期投入", BLUE)
    draw_chip(draw, 304, 238, "纯增量分成", ORANGE)
    draw_chip(draw, 512, 238, "24小时体检报告", BLUE)
    draw_chip(draw, 1438, 238, "试点 → 复制 → 扩张", ORANGE)

    draw_stick_person(draw, 270, 464, accent=BLUE, pose="point", scale=1.08)
    draw_speech(draw, 112, 322, "选一家有潜力的民宿", outline=BLUE, max_width=230)
    draw_chart(draw, 418, 526, scale=1.2)
    draw_magnifier(draw, 520, 394, 0.86)

    draw_ai_person(draw, 952, 540, label="AI", accent=BLUE, scale=1.2)
    draw_speech(draw, 846, 320, "把运营动作自动化", outline=ORANGE, max_width=230)
    draw_cloud(draw, 932, 702, 1.25)
    draw_gear(draw, 774, 382, 1.1)
    draw_bulb(draw, 1138, 336, 1.1)

    draw_inn(draw, 1302, 376, scale=1.1)
    draw_mountains(draw, 1376, 340, scale=1.0)
    draw_money_tag(draw, 1472, 690, "分成")
    draw_arrow(draw, 392, 454, 730, 470, color=BLUE, width=6)
    draw_arrow(draw, 1060, 506, 1308, 496, color=ORANGE, width=7)
    draw_arrow(draw, 1448, 656, 1240, 710, color=BLUE, width=5, dashed=True)

    draw.text((196, 802), "发现潜力", font=BODY_BOLD_FONT, fill=BLUE)
    draw.text((804, 808), "AI 自动化运营", font=BODY_BOLD_FONT, fill=ORANGE)
    draw.text((1336, 802), "商家收入增长", font=BODY_BOLD_FONT, fill=BLUE)
    draw_social_card(draw, 128, 854, "抖音", BLUE)
    draw_social_card(draw, 326, 854, "小红书", ORANGE)
    draw_social_card(draw, 1126, 854, "舆情", BLUE)
    draw_social_card(draw, 1324, 854, "定价", ORANGE)
    draw_social_card(draw, 1522, 854, "洞察", BLUE)

    draw_bottom_rule(draw, 1, 6, "项目定位：用 AI 把文旅资源做成可运营、可分成、可复制的增长系统")
    return image


def draw_opportunity_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "为什么现在要做文旅引爆点？", "本地好资源不缺，缺的是持续内容、数据运营和可复制打法")

    draw_sketch_round_rect(draw, (138, 292, 650, 826), radius=34, fill=WHITE, outline=ORANGE, width=5)
    draw.text((186, 326), "传统民宿运营", font=SECTION_FONT, fill=ORANGE_DEEP)
    left_items = [
        ("内容靠灵感", "断更、同质化、没人看"),
        ("定价靠经验", "旺季涨不够，淡季降太慢"),
        ("评论靠人工", "好评没放大，差评响应慢"),
        ("客群不清楚", "不知道真正喜欢什么"),
    ]
    y = 394
    for title, body in left_items:
        draw.ellipse((188, y + 8, 204, y + 24), fill=ORANGE)
        draw.text((224, y), title, font=BODY_BOLD_FONT, fill=NAVY)
        draw.text((224, y + 40), body, font=SMALL_FONT, fill=GRAY)
        y += 100
    draw_stick_person(draw, 520, 770, accent=ORANGE, pose="stand", scale=0.78)
    draw_speech(draw, 402, 618, "忙，但不系统", outline=ORANGE, max_width=180)

    draw_sketch_round_rect(draw, (1268, 292, 1780, 826), radius=34, fill=WHITE, outline=BLUE, width=5)
    draw.text((1318, 326), "AI合伙人系统", font=SECTION_FONT, fill=BLUE)
    right_items = [
        ("内容潮引擎", "批量生成游客视角内容"),
        ("市场洞察脚本", "竞品、天气、活动自动监控"),
        ("口碑哨兵", "评论分级、回复建议、风险预警"),
        ("洞察周报", "把人群画像变成服务动作"),
    ]
    y = 394
    for title, body in right_items:
        draw.ellipse((1318, y + 8, 1334, y + 24), fill=BLUE)
        draw.text((1352, y), title, font=BODY_BOLD_FONT, fill=NAVY)
        draw.text((1352, y + 40), body, font=SMALL_FONT, fill=GRAY)
        y += 100
    draw_ai_person(draw, 1640, 770, accent=BLUE, scale=0.78)
    draw_speech(draw, 1512, 618, "持续自动化", outline=BLUE, max_width=190)

    draw_arrow(draw, 696, 550, 1214, 550, color=BLUE, width=8)
    draw.text((812, 480), "从“人肉运营”", font=BODY_BOLD_FONT, fill=NAVY)
    draw.text((858, 548), "升级为", font=SECTION_FONT, fill=ORANGE)
    draw.text((812, 616), "“系统增长”", font=BODY_BOLD_FONT, fill=NAVY)
    draw_bulb(draw, 920, 394, 1.25)
    draw_gear(draw, 1044, 702, 1.2)
    draw_metric_tile(draw, 760, 746, "7×24", "自动监控", BLUE)
    draw_metric_tile(draw, 978, 746, "24h", "生成报告", ORANGE)

    draw_bottom_rule(draw, 2, 6, "机会判断：不是做更多人工动作，而是把增长环节系统化、自动化")
    return image


def draw_services_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "AI合伙人的四大核心赋能", "同一套系统，解决客流、价格、口碑、体验四个关键问题")

    draw_ai_person(draw, 960, 516, label="AI", accent=BLUE, scale=1.18)
    draw_speech(draw, 856, 320, "一套系统\n四个增长杠杆", outline=ORANGE, max_width=220)

    cards = [
        (142, 292, "内容引爆矩阵", "自动生成游客风格内容\n把商家打造成引爆点", BLUE, "爆"),
        (1156, 292, "智能动态定价", "监控竞品、天气、活动\n给出每日调价建议", ORANGE, "价"),
        (142, 676, "口碑舆情维护", "好评自动放大\n差评草拟高情商回复", ORANGE, "评"),
        (1156, 676, "客户洞察分析", "识别人群画像与兴趣\n每周输出行动建议", BLUE, "客"),
    ]
    for x, y, title, subtitle, accent, icon in cards:
        draw_card(draw, x, y, 430, 186, title, subtitle, accent, icon)

    draw_arrow(draw, 586, 392, 796, 458, color=BLUE, width=5)
    draw_arrow(draw, 1150, 392, 1046, 458, color=ORANGE, width=5)
    draw_arrow(draw, 586, 786, 796, 612, color=ORANGE, width=5)
    draw_arrow(draw, 1150, 786, 1046, 612, color=BLUE, width=5)
    draw_bulb(draw, 728, 282, 0.96)
    draw_gear(draw, 1228, 262, 0.96)
    draw_bulb(draw, 1262, 868, 0.96)
    draw_gear(draw, 704, 870, 0.96)

    draw_annotation(draw, 706, 716, "运营飞轮", "内容带客 → 数据反馈 → 再优化内容与价格", BLUE)
    draw_bottom_rule(draw, 3, 6, "系统架构：内容、价格、口碑、洞察互相反馈，形成运营飞轮")
    return image


def draw_content_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "内容引爆矩阵如何带来客流？", "用游客视角批量生产内容，再把注意力精准导回商家")

    draw_inn(draw, 1400, 520, scale=1.08)
    draw_mountains(draw, 1464, 486, scale=0.92)
    draw_stick_person(draw, 190, 742, accent=ORANGE, pose="walk", scale=0.82)
    draw_speech(draw, 118, 574, "游客正在被种草", outline=ORANGE, max_width=192)
    draw_megaphone(draw, 720, 506, scale=1.18)
    draw_ai_person(draw, 660, 786, label="AI", accent=BLUE, scale=0.82)

    phone_cards = [
        (248, 322, "抖音", BLUE),
        (404, 386, "小红书", ORANGE),
        (562, 322, "视频号", BLUE),
    ]
    for x, y, label, color in phone_cards:
        draw_phone(draw, x, y, label, color)
        draw_arrow(draw, x + 70, y + 218, 716, 550, color=color, width=4, dashed=True)

    draw_arrow(draw, 938, 530, 1348, 592, color=ORANGE, width=8)
    draw_arrow(draw, 1542, 514, 1552, 394, color=BLUE, width=5)
    draw_speech(draw, 1346, 330, "内容精准植入\n商家信息", outline=BLUE, max_width=208)

    draw_sketch_round_rect(draw, (1054, 732, 1698, 890), radius=28, fill=WHITE, outline=BLUE, width=4)
    draw.text((1094, 762), "内容生产 SOP", font=BODY_BOLD_FONT, fill=BLUE)
    sop = ["1. 提炼民宿亮点", "2. 生成游客口吻素材", "3. 分平台改写发布", "4. 追踪互动与转化"]
    for index, item in enumerate(sop):
        draw_chip(draw, 1094 + (index % 2) * 292, 812 + (index // 2) * 48, item, ORANGE if index % 2 else BLUE)

    draw_metric_tile(draw, 170, 842, "批量", "持续内容潮", BLUE)
    draw_metric_tile(draw, 388, 842, "种草", "真实游客视角", ORANGE)
    draw_metric_tile(draw, 606, 842, "导流", "精准植入商家", GREEN)
    draw_bottom_rule(draw, 4, 6, "获客逻辑：不是硬广，而是把真实游客视角内容做成持续内容潮")
    return image


def draw_dashboard_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "动态定价 + 口碑哨兵", "价格、天气、活动、评论四类信号，变成每天可执行的动作")

    draw_dashboard(draw, 126, 304, 540, 352, "大理市场洞察", BLUE)
    draw_dashboard(draw, 1254, 304, 540, 352, "全网口碑哨兵", ORANGE)
    draw_ai_person(draw, 960, 612, label="AI", accent=BLUE, scale=1.02)
    draw_speech(draw, 848, 392, "把信号整理成\n行动建议", outline=ORANGE, max_width=230)

    signals = [
        (740, 302, "竞品价格", BLUE),
        (902, 268, "机票热度", ORANGE),
        (1066, 302, "天气活动", BLUE),
        (810, 796, "好评放大", GREEN),
        (1002, 796, "差评预警", RED),
    ]
    for x, y, text, color in signals:
        draw_chip(draw, x, y, text, color if color in [BLUE, ORANGE, GREEN] else ORANGE)

    draw_arrow(draw, 666, 484, 820, 536, color=BLUE, width=5)
    draw_arrow(draw, 1254, 484, 1100, 536, color=ORANGE, width=5)
    draw_arrow(draw, 960, 728, 960, 834, color=BLUE, width=6)
    draw_sketch_round_rect(draw, (684, 844, 1236, 932), radius=28, fill=WHITE, outline=GREEN, width=4)
    draw.text((734, 866), "每日 8:00 推送：分房型调价建议 + 评论处理清单", font=SMALL_BOLD_FONT, fill=NAVY)

    draw_review_bubble(draw, 1308, 696, "房间景观很棒，推荐！", GREEN, "好评")
    draw_review_bubble(draw, 1490, 814, "入住等待偏久，需要安抚", ORANGE, "差评")
    draw_money_tag(draw, 302, 694, "涨价")
    draw_money_tag(draw, 466, 736, "促销")
    draw_magnifier(draw, 208, 746, 0.74, BLUE)

    draw_bottom_rule(draw, 5, 6, "执行机制：复杂数据不直接给商家，只给“今天该怎么做”的建议")
    return image


def draw_launch_slide(backdrop: Image.Image | None = None) -> Image.Image:
    image, draw = draw_slide_background(backdrop)
    draw_title_block(draw, "第一阶段如何启动合作？", "先给一份商家无法拒绝的 AI 体检报告，再谈合作与分成")

    step_y = 474
    centers = [218, 638, 1080, 1512]
    labels = [
        ("1", "锁定目标民宿"),
        ("2", "24小时生成体检报告"),
        ("3", "拿报告建立信任"),
        ("4", "启动合作并分成"),
    ]
    for idx, (cx, (num, text)) in enumerate(zip(centers, labels)):
        color = BLUE if idx % 2 == 0 else ORANGE
        draw_sketch_round_rect(draw, (cx - 112, step_y - 62, cx + 112, step_y + 62), radius=30, fill=WHITE, outline=color, width=4)
        draw.ellipse((cx - 142, step_y - 20, cx - 98, step_y + 24), fill=ORANGE if idx % 2 == 0 else BLUE, outline=None)
        draw.text((cx - 129, step_y - 9), num, font=LABEL_FONT, fill=WHITE)
        rendered = wrap_text(text, SMALL_BOLD_FONT, 164)
        box = draw.multiline_textbbox((0, 0), rendered, font=SMALL_BOLD_FONT, spacing=5, align="center")
        draw.multiline_text((cx - (box[2] - box[0]) / 2, step_y - 21), rendered, font=SMALL_BOLD_FONT, fill=NAVY, spacing=5, align="center")
        if idx < len(centers) - 1:
            draw_arrow(draw, cx + 122, step_y, centers[idx + 1] - 136, step_y, color=ORANGE, width=6)

    draw_stick_person(draw, 214, 750, accent=BLUE, pose="point", scale=0.82)
    draw_inn(draw, 96, 714, scale=0.78)
    draw_clipboard(draw, 498, 642, "AI体检报告", ["全网口碑舆情汇总", "5公里竞品价格对比", "建议动作清单"], BLUE)
    draw_stick_person(draw, 1126, 752, accent=ORANGE, pose="stand", scale=0.84)
    draw_speech(draw, 1018, 622, "这份报告先免费给您", outline=ORANGE, max_width=230)
    draw_ai_person(draw, 1546, 764, label="AI", accent=BLUE, scale=0.88)
    draw_money_tag(draw, 1642, 632, "分成")
    draw_chart(draw, 1280, 822, scale=1.1)

    draw_annotation(draw, 120, 248, "敲门砖", "不先卖方案，先拿出情报能力证明", BLUE)
    draw_annotation(draw, 1486, 248, "成交点", "商家只为纯增量收益分成，阻力更低", ORANGE)
    draw_bottom_rule(draw, 6, 6, "落地策略：先用免费体检报告建立信任，再推进 AI 合伙人合作")
    return image


SLIDES: list[SlideSpec] = [
    SlideSpec(
        "文旅引爆点——AI合伙人运营系统",
        "大理站｜把潜力民宿打造成“引爆点”",
        "P1-文旅引爆点-封面战略图.png",
        "封面总览页：左侧是发现潜力民宿与内容机会，中央是 AI 合伙人系统，右侧是被点亮的民宿与收入增长结果；要有山景、小民宿、人物、箭头、图表和增长符号，整体像高级业务战略图。",
        draw_cover_slide,
    ),
    SlideSpec(
        "为什么现在要做文旅引爆点？",
        "本地好资源不缺，缺的是持续内容、数据运营和可复制打法",
        "P2-机会判断-系统增长.png",
        "左右对比布局：左侧是传统民宿运营的混乱低效场景，右侧是 AI 合伙人系统驱动的整洁增长场景；中间有从“人肉运营”升级为“系统增长”的桥接箭头；保留两张大圆角信息卡片容器与少量背景手绘线稿。",
        draw_opportunity_slide,
    ),
    SlideSpec(
        "AI合伙人的四大核心赋能",
        "同一套系统，解决客流、价格、口碑、体验四个关键问题",
        "P3-AI合伙人-四大赋能.png",
        "中心辐射布局：中央是 AI 合伙人，四周四张大卡片分别代表内容引爆、动态定价、口碑维护、客户洞察，四条连接线形成运营飞轮，画面稳定、对称、清爽。",
        draw_services_slide,
    ),
    SlideSpec(
        "内容引爆矩阵如何带来客流？",
        "用游客视角批量生产内容，再把注意力精准导回商家",
        "P4-内容引爆矩阵-获客路径.png",
        "获客路径布局：左侧是多个内容平台与游客，中部是内容引爆矩阵和扩音器，右侧是民宿与转化路径，下方有 SOP 或流程卡片，表达从内容种草到导流成交的全过程。",
        draw_content_slide,
    ),
    SlideSpec(
        "动态定价 + 口碑哨兵",
        "价格、天气、活动、评论四类信号，变成每天可执行的动作",
        "P5-动态定价-口碑哨兵.png",
        "双看板布局：左侧是市场洞察仪表盘，右侧是口碑哨兵仪表盘，中间 AI 助手汇总信号，下方是每日行动建议；可有评价气泡、图表、信号标签和提醒框。",
        draw_dashboard_slide,
    ),
    SlideSpec(
        "第一阶段如何启动合作？",
        "先给一份商家无法拒绝的 AI 体检报告，再谈合作与分成",
        "P6-启动合作-体检报告路径.png",
        "路线图布局：从目标民宿、免费体检报告、商家沟通到 AI 合伙人合作与分成，形成一条清晰的推进路径；有诊断报告、握手、AI 角色、收益符号和阶段箭头。",
        draw_launch_slide,
    ),
]


def build_gemini_prompt(spec: SlideSpec) -> str:
    reference_prompt = load_reference_prompt_text()
    return textwrap.dedent(
        f"""
        Create one finished 16:9 slide illustration for a Chinese business PowerPoint.
        请把随附参考图当成强约束参考，不是普通灵感图。

        视觉必须尽量贴近参考图：
        - 白底，大面积留白，整体轻盈
        - 粗黑线条 + 蓝橙主色 + 少量浅色块
        - 手绘火柴人 / 商业信息图混合风格
        - 咨询汇报级别的整洁感，不要低幼，不要照片感，不要 3D
        - 保留参考图里那种大圆角卡片、流程箭头、背景浅线稿、轻微发光与业务示意图感觉

        当前页标题：{spec.title}
        当前页副标题：{spec.subtitle}
        当前页要表达的内容：{spec.scene_prompt}

        版式要求：
        - 顶部必须预留完整标题区和副标题区，不要被插画占满
        - 主画面集中在中部和下部，构成一个完整的可讲述业务场景
        - 如需卡片、屏幕、对话框、流程块，请直接画出容器，但里面不要出现可读文字
        - 可以出现人物、民宿、山景、图表、箭头、齿轮、扩音器、金币、信号标签等业务元素
        - 不要直接复制参考图内容本身，要复用的是视觉语言、版式气质、线条风格与信息图组织方式

        角色造型要求（非常重要）：
        - 小人必须是清晰、稳定、可爱的商业火柴人：圆头、简单表情、胶囊形身体、正常比例四肢
        - AI 角色使用电脑屏幕头 + 简洁身体，比例端正，不要恐怖、畸形或机械怪物感
        - 人物数量要少而明确，避免画一堆很小、很乱、粘连、扭曲的人
        - 手脚关节要自然，不要多手、多腿、断肢、歪脸、怪表情、过长四肢
        - 背景里的小人保持浅淡和辅助性，不要抢主视觉

        硬约束：
        - 不要任何可读文字、汉字、英文、数字、水印、logo、品牌标志、UI 标签
        - 不要脏乱拥挤，不要密集小元素堆满画面
        - 不要边框，不要裁切，不要黑底

        用户原始补充提示（同样需要吸收）：
        {reference_prompt or "火柴人动画风格，线条简洁、色彩明快（蓝色和橙色为主，白色背景）。"}
        """
    ).strip()


def summarize_http_error(error_body: str) -> str:
    try:
        payload = json.loads(error_body)
        message = payload.get("error", {}).get("message")
        status = payload.get("error", {}).get("status")
        if message and status:
            return f"{status}: {message}"
        if message:
            return message
    except Exception:
        pass
    return error_body[:480]


def generate_gemini_image(prompt: str, output_path: Path, config: GeminiConfig, reference_images: list[Path]) -> tuple[bool, str]:
    model = urllib.parse.quote(config.model, safe="")
    url = f"{config.base_url.rstrip('/')}/models/{model}:generateContent"
    parts: list[dict[str, object]] = [{"text": prompt}]
    for reference_image in reference_images:
        parts.append(
            {
                "inline_data": {
                    "mime_type": mime_type_for_path(reference_image),
                    "data": base64.b64encode(reference_image.read_bytes()).decode("utf-8"),
                }
            }
        )
    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
            "imageConfig": {"aspectRatio": "16:9", "imageSize": "2K"},
        },
    }
    payload_json = json.dumps(payload, ensure_ascii=False).encode("utf-8")

    try:
        response = subprocess.run(
            [
                "curl",
                "-sS",
                "-X",
                "POST",
                url,
                "-H",
                f"x-goog-api-key: {config.api_key}",
                "-H",
                "Content-Type: application/json",
                "--data-binary",
                "@-",
            ],
            input=payload_json,
            capture_output=True,
            timeout=240,
            check=False,
        )
    except FileNotFoundError:
        response = None
    except subprocess.TimeoutExpired:
        return False, "Gemini image request timed out after 240 seconds."
    except Exception as exc:
        return False, str(exc)

    if response is not None:
        if response.returncode != 0:
            stderr = response.stderr.decode("utf-8", errors="replace").strip()
            return False, stderr or f"curl exited with code {response.returncode}"
        body_text = response.stdout.decode("utf-8", errors="replace")
        try:
            response_payload = json.loads(body_text)
        except json.JSONDecodeError:
            return False, body_text[:480]
    else:
        request = urllib.request.Request(
            url,
            data=payload_json,
            headers={"Content-Type": "application/json", "x-goog-api-key": config.api_key},
            method="POST",
        )
        try:
            with urllib.request.urlopen(request, timeout=180) as response_fallback:
                response_payload = json.loads(response_fallback.read().decode("utf-8"))
        except urllib.error.HTTPError as exc:
            body = exc.read().decode("utf-8", errors="replace")
            return False, summarize_http_error(body)
        except Exception as exc:
            return False, str(exc)

    if response_payload.get("error"):
        return False, summarize_http_error(json.dumps(response_payload, ensure_ascii=False))

    parts = response_payload.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    for part in parts:
        inline_data = part.get("inlineData") or part.get("inline_data")
        if inline_data and inline_data.get("data"):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(base64.b64decode(inline_data["data"]))
            return True, "ok"

    text_parts = [part.get("text", "") for part in parts if part.get("text")]
    message = "；".join(text_parts) if text_parts else "Gemini response did not contain inline image data."
    return False, message


def maybe_generate_gemini_backdrops(output_dir: Path, use_gemini: bool, verbose: bool = False) -> tuple[dict[str, Image.Image], list[str]]:
    if not use_gemini:
        return {}, []

    config = gemini_config_from_env()
    if not config:
        raise RuntimeError("已显式启用 Gemini 生图，但未检测到 GEMINI_API_KEY / GOOGLE_API_KEY / GOOGLE_GENAI_API_KEY。")
    refs = reference_image_paths()
    if not refs:
        raise RuntimeError("已显式启用 Gemini 生图，但未在归档中找到可用参考图。")

    backdrops: dict[str, Image.Image] = {}
    errors: list[str] = []
    backdrop_dir = output_dir / "gemini-backdrops"
    for index, spec in enumerate(SLIDES, start=1):
        if verbose:
            print(f"[Gemini] P{index}: {spec.title}")
        backdrop_path = backdrop_dir / f"P{index}-gemini-backdrop.png"
        ok, message = generate_gemini_image(build_gemini_prompt(spec), backdrop_path, config, refs)
        if ok:
            backdrops[spec.filename] = Image.open(backdrop_path).convert("RGB")
            if verbose:
                print(f"[Gemini] P{index}: ok")
        else:
            errors.append(f"P{index} {config.model}: {message}")
            if verbose:
                print(f"[Gemini] P{index}: {message}")
    return backdrops, errors


def create_contact_sheet(images: list[Path], output_path: Path) -> None:
    thumb_w, thumb_h = 500, 281
    padding = 28
    label_h = 48
    columns = 3
    rows = math.ceil(len(images) / columns)
    sheet = Image.new(
        "RGB",
        (padding * (columns + 1) + thumb_w * columns, padding * (rows + 1) + (thumb_h + label_h) * rows),
        NAVY,
    )
    draw = ImageDraw.Draw(sheet)
    font = load_font(24, heavy=True)
    for index, image_path in enumerate(images):
        image = Image.open(image_path).convert("RGB").resize((thumb_w, thumb_h))
        row = index // columns
        col = index % columns
        x = padding + col * (thumb_w + padding)
        y = padding + row * (thumb_h + label_h + padding)
        sheet.paste(image, (x, y))
        draw.rounded_rectangle((x - 2, y - 2, x + thumb_w + 2, y + thumb_h + 2), radius=8, outline="#475569", width=2)
        draw.text((x, y + thumb_h + 14), f"P{index + 1}  {images[index].stem.split('-', 1)[-1]}", font=font, fill=WHITE)
    sheet.save(output_path, quality=92)


def build_demo(use_gemini: bool = False, output_root: Path = DEFAULT_OUTPUT_ROOT, verbose: bool = False) -> tuple[Path, Path]:
    output_dir = create_timestamped_output_dir(output_root, "文旅引爆点-clawlink六页增强样张")
    refs = reference_image_paths()

    backdrops, gemini_errors = maybe_generate_gemini_backdrops(output_dir, use_gemini, verbose=verbose)
    if use_gemini and not backdrops:
        joined = "；".join(gemini_errors[:3]) if gemini_errors else "Gemini 未返回可用图片。"
        raise RuntimeError(f"Gemini 生图未成功产出任何页面：{joined}")
    slide_images = [output_dir / spec.filename for spec in SLIDES]

    for image_path, spec in zip(slide_images, SLIDES):
        backdrop = backdrops.get(spec.filename)
        spec.builder(backdrop).save(image_path, quality=95)

    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for image_path in slide_images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)

    mode_suffix = "Gemini底图增强" if backdrops else "本地绘制增强"
    pptx_path = output_dir / f"文旅引爆点-ClawLink风格-图片型PPT-六页样张-{mode_suffix}.pptx"
    presentation.save(str(pptx_path))

    contact_sheet = output_dir / "contact-sheet.png"
    create_contact_sheet(slide_images, contact_sheet)

    readme_lines = [
        "这是根据《文旅引爆点项目方案.txt》提炼的 6 页 ClawLink 风格图片型 PPT 增强样张。",
        f"生成模式：{mode_suffix}",
        f"源文档：{DOC_PATH}",
        f"PPT：{pptx_path.name}",
        "",
        "页面结构：",
    ]
    for index, spec in enumerate(SLIDES, start=1):
        readme_lines.append(f"P{index}：{spec.title}")
    if gemini_errors:
        readme_lines.extend(["", "Gemini 生成提示：", *gemini_errors])
    if refs:
        readme_lines.extend(["", "参考图：", *[str(path) for path in refs]])
    if CLAWLINK_PROMPT_REFERENCE and CLAWLINK_PROMPT_REFERENCE.exists():
        readme_lines.extend(["", f"提示词参考：{CLAWLINK_PROMPT_REFERENCE}"])
    readme_lines.extend(
        [
            "",
            "说明：Gemini 模式会把归档中的参考图一并喂给模型；中文标题、标签、说明文字仍由本地程序叠加。",
        ]
    )
    (output_dir / "README.txt").write_text("\n".join(readme_lines), encoding="utf-8")
    (output_dir / "gemini-prompts.json").write_text(
        json.dumps(
            [
                {
                    "page": index,
                    "title": spec.title,
                    "model": os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL),
                    "prompt": build_gemini_prompt(spec),
                    "referenceImages": [str(path) for path in refs],
                    "referencePrompt": str(CLAWLINK_PROMPT_REFERENCE) if CLAWLINK_PROMPT_REFERENCE else "",
                }
                for index, spec in enumerate(SLIDES, start=1)
            ],
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    manifest_path = write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(slide_images),
        mode=mode_suffix,
        contactSheet=str(contact_sheet),
        readme=str(output_dir / "README.txt"),
        promptCatalog=str(output_dir / "gemini-prompts.json"),
        imageDirectory=str(output_dir),
        geminiBackdropDirectory=str((output_dir / "gemini-backdrops") if backdrops else ""),
        geminiErrors=gemini_errors,
        referenceImages=[str(path) for path in refs],
        referencePrompt=str(CLAWLINK_PROMPT_REFERENCE) if CLAWLINK_PROMPT_REFERENCE else "",
    )

    return pptx_path, manifest_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a ClawLink-style image-based PPT demo.")
    parser.add_argument(
        "--use-gemini",
        action="store_true",
        help="Use Gemini image generation with archived reference images; requires a valid Gemini/Google API key.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=DEFAULT_OUTPUT_ROOT,
        help="Directory where timestamped demo folders are created.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Print per-page Gemini progress logs.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    pptx_path, _ = build_demo(use_gemini=args.use_gemini, output_root=args.output_root, verbose=args.verbose)
    print(pptx_path)

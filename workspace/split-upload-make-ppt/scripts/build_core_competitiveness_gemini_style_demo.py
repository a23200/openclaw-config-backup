from __future__ import annotations

import argparse
import base64
import json
import os
import re
import subprocess
import textwrap
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    latest_matching_path,
    load_env_candidates,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import BLUE, NAVY, WHITE, load_font, wrap_text
from build_neural_params_controlled_demo import draw_centered


DESKTOP = Path("/Users/mac/Desktop")
SOURCE_DOC = DESKTOP / "项目核心竞争力分析.md"
ARCHIVE_DIR = latest_archive_dir()
OUTPUT_ROOT = repo_runtime_dir("image-ppt-demo")
DEFAULT_GEMINI_IMAGE_MODEL = "gemini-3-pro-image-preview"

TEXT_FREE_STYLE_REFERENCE = first_existing_path(
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini真实生图七页样张/gemini-backdrops/P1-gemini-backdrop.png"),
    latest_matching_path(ARCHIVE_DIR, "runtime-image-ppt-demo/*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png") if ARCHIVE_DIR else None,
)
TEXT_RICH_STYLE_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    DESKTOP / "参考图-大模型能力突然变强.png",
)
STYLE_REFERENCE = first_existing_path(TEXT_FREE_STYLE_REFERENCE, TEXT_RICH_STYLE_REFERENCE)
LAYOUT_REFERENCE_ROOT = None
NATIVE_TEXT_REFERENCE_ROOT = first_existing_path(
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini融合文字*页样张"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini图文融合生图*页样张"),
)

WIDTH = 1920
HEIGHT = 1080

TITLE_FONT = load_font(58, heavy=True)
TITLE_SMALL_FONT = load_font(50, heavy=True)
BODY_FONT = load_font(30)
BODY_SMALL_FONT = load_font(26)
CAPTION_FONT = load_font(19, heavy=True)


@dataclass(frozen=True)
class SlidePage:
    title: str
    bullets: list[str]
    scene_prompt: str
    filename: str


@dataclass(frozen=True)
class MarkdownSection:
    level: int
    title: str
    lines: list[str]


def clean_markdown_text(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^\s*[-*+]\s*", "", text.strip())
    text = re.sub(r"\s+", " ", text)
    return text.strip(" ：:")


def shorten(text: str, limit: int = 42) -> str:
    text = clean_markdown_text(text)
    return text if len(text) <= limit else f"{text[: limit - 1]}…"


def parse_markdown_sections(markdown: str) -> tuple[str, list[MarkdownSection]]:
    deck_title = SOURCE_DOC.stem
    sections: list[MarkdownSection] = []
    current_level = 0
    current_title = ""
    current_lines: list[str] = []
    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*$")

    for raw_line in markdown.splitlines():
        match = heading_re.match(raw_line.strip())
        if match:
            if current_title:
                sections.append(MarkdownSection(current_level, current_title, current_lines))
            current_level = len(match.group(1))
            current_title = clean_markdown_text(match.group(2))
            current_lines = []
            if current_level == 1:
                deck_title = current_title
            continue
        if current_title:
            current_lines.append(raw_line)

    if current_title:
        sections.append(MarkdownSection(current_level, current_title, current_lines))
    return deck_title, sections


def extract_points(section: MarkdownSection, max_points: int = 2) -> list[str]:
    candidates: list[str] = []
    for raw_line in section.lines:
        line = clean_markdown_text(raw_line)
        if not line or line == "---":
            continue
        if line.startswith("结论："):
            candidates.insert(0, line)
        elif "：" in line or len(line) >= 18:
            candidates.append(line)

    unique: list[str] = []
    for item in candidates:
        item = shorten(item.replace("结论：", ""), 38)
        if item and item not in unique:
            unique.append(item)
        if len(unique) >= max_points:
            break
    return unique or [shorten(section.title, 46)]


def sanitize_filename_part(text: str) -> str:
    text = re.sub(r"[\\/:*?\"<>|`]+", "", text)
    text = re.sub(r"\s+", "", text)
    return text[:26] or "页面"


def scene_prompt_for(title: str, bullets: list[str], deck_title: str) -> str:
    joined = "；".join(bullets)
    if "闭环" in title or "自进化" in title:
        return "闭环自进化系统页：用飞轮、反馈回路、数据反哺、决策引擎和下一轮优化表达系统越跑越强。"
    if "数据" in title or "图谱" in title:
        return "数据资产页：多个分散节点和用户信号汇聚成区域级图谱网络，表达从单点数据升级到市场风向定义能力。"
    if "人机" in title or "混合智能" in title:
        return "人机混合智能页：AI 执行官处理规模化动作，人类指挥官做关键判断，形成互补协作的组合智能。"
    if "平台" in title or "扩张" in title:
        return "平台扩张页：从一个样板节点复制到大量节点，呈现标准化平台、网络扩张、规模效应和低边际成本。"
    if "结论" in title or "What" in title or "How" in title:
        return "最终结论页：用冰山、隐形引擎或生命体内核表达表面的 What 可复制，但深层 How 无法复制。"
    if "引言" in title or "技能" in title or "场域" in title:
        return "总览页：对比一次性 Skill 和持续进化 Field，中心呈现会学习、会积累、会扩张的系统生命体。"
    return f"战略信息图页：围绕《{deck_title}》中的“{title}”，用清晰模块、箭头、抽象系统图和业务符号表达：{joined}"


def build_pages_from_markdown(markdown: str) -> tuple[str, list[SlidePage]]:
    deck_title, sections = parse_markdown_sections(markdown)
    h2_or_h3 = [section for section in sections if section.level in {2, 3}]
    layer_sections = [section for section in h2_or_h3 if section.level == 3 and ("护城河" in section.title or "层" in section.title)]
    intro_sections = [section for section in h2_or_h3 if "引言" in section.title or "技能" in section.title or "场域" in section.title]
    final_sections = [section for section in h2_or_h3 if "最终" in section.title or "结论" in section.title]

    selected: list[MarkdownSection] = []
    if intro_sections:
        selected.append(intro_sections[0])
    if layer_sections:
        selected.extend(layer_sections)
    else:
        selected.extend(section for section in h2_or_h3 if section not in selected and section not in final_sections)
    selected.extend(section for section in final_sections if section not in selected)

    if not selected:
        selected = [section for section in sections if section.level != 1] or sections[:1]

    pages: list[SlidePage] = []
    for index, section in enumerate(selected, start=1):
        page_title = deck_title.split("：", 1)[0] if index == 1 and "引言" in section.title else section.title
        points = extract_points(section, max_points=2)
        filename = f"P{index}-{sanitize_filename_part(page_title)}.png"
        pages.append(
            SlidePage(
                title=page_title,
                bullets=points,
                scene_prompt=scene_prompt_for(section.title, points, deck_title),
                filename=filename,
            )
        )
    return deck_title, pages


def get_api_key() -> str:
    load_env_candidates()
    for key in ["GEMINI_API_KEY", "GOOGLE_API_KEY", "GOOGLE_GENAI_API_KEY"]:
        value = os.environ.get(key)
        if value:
            return value
    raise RuntimeError("Gemini API key not found in env.")


def mime_type_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    return "application/octet-stream"


def layout_reference_for(page: SlidePage, native_text: bool = False) -> Path | None:
    root = NATIVE_TEXT_REFERENCE_ROOT if native_text else LAYOUT_REFERENCE_ROOT
    if not root:
        return None
    candidate = root / page.filename
    return candidate if candidate.exists() else None


def find_latest_output_dir(output_root: Path) -> Path | None:
    if not output_root.exists():
        return None
    candidates = [
        path
        for path in output_root.iterdir()
        if path.is_dir() and "项目核心竞争力分析-gemini动态生图" in path.name
    ]
    return sorted(candidates)[-1] if candidates else None


def bullets_text(page: SlidePage) -> str:
    return "\n".join(f"• {point}" for point in page.bullets[:2])


def build_prompt(page: SlidePage, native_text: bool = False) -> str:
    context_points = "\n".join(f"- {point}" for point in page.bullets)
    exact_lines = "\n".join([page.title, *[f"• {point}" for point in page.bullets[:2]]])
    if native_text:
        return textwrap.dedent(
            f"""
            Create one finished 16:9 slide illustration for a Chinese strategy PowerPoint.
            The first reference image defines the overall visual language and style.
            If a second reference image is provided, use it for text placement rhythm and integrated editorial composition.

            Visual style requirements:
            - white background with large clean negative space
            - thick black outlines, deep blue accents, orange accents, light blue and light orange blocks
            - polished business-infographic style, not childish, not photorealistic, not 3D
            - integrated editorial composition where text and graphics feel designed together

            Main scene to generate:
            {page.scene_prompt}

            You must render these exact Chinese lines directly inside the image, with clean typography and correct spelling:
            {exact_lines}

            Text layout rules:
            - place all required text in one natural editorial zone near the upper-left or upper-middle area
            - the text zone must feel integrated into the illustration, like a designed info card or embedded layout block
            - title larger, bullets smaller, all typography crisp and readable
            - do not add any extra text beyond the exact lines above

            Hard constraints:
            - no wrong characters, no garbled Chinese, no fake text, no pseudo text, no English, no numbers unless already in the required Chinese lines
            - no watermark, no logo, no extra labels
            - no malformed people, no creepy faces, no extra hands or legs
            """
        ).strip()
    return textwrap.dedent(
        f"""
        Create one finished 16:9 slide illustration for a Chinese strategy PowerPoint.
        The first reference image defines the overall visual language and style.
        If a second reference image is provided, use it only for page composition, layout rhythm, and information density.

        Visual style requirements:
        - white background with large clean negative space
        - thick black outlines, deep blue titles, orange accents, light blue and light orange blocks
        - polished business-infographic style, not childish, not photorealistic, not 3D
        - look like a premium consulting deck mixed with hand-drawn system thinking graphics
        - subtle geometric ornaments, arrows, loops, cards, dashboards, abstract business symbols

        Content context only. Do not render these words in the image.
        Page title context: {page.title}
        Key point context:
        {context_points}

        Main scene to generate:
        {page.scene_prompt}

        Composition constraints:
        - design a natural upper-left editorial information zone that feels built into the illustration
        - this zone should look like part of the composition: a light glass card, paper card, or calm low-detail area integrated with the scene
        - do not create a giant horizontal title banner across the whole page
        - keep the main illustration mainly in the middle and lower half
        - create a distinct page-level composition, not a generic background
        - use 3 to 6 major visual groups only; keep hierarchy clear
        - if reference layout is provided, keep a similar page structure but do not copy existing text

        Hard constraints:
        - no readable text, no fake text, no pseudo text, no scribbled labels, no Chinese, no English, no numbers
        - no labels, no logos, no watermark, no UI text
        - no crowded tiny icons everywhere
        - no black background, no border
        - no malformed people, no creepy faces, no extra hands or legs
        """
    ).strip()


def summarize_http_error(error_body: str) -> str:
    try:
        payload = json.loads(error_body)
        message = payload.get("error", {}).get("message")
        status = payload.get("error", {}).get("status")
        if message and status:
            return f"{status}: {message}"
        if message:
            return message
    except Exception:
        pass
    return error_body[:480]


def generate_gemini_image(page: SlidePage, output_path: Path, native_text: bool = False) -> str:
    api_key = get_api_key()
    model = urllib.parse.quote(os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL), safe="")
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    references: list[Path] = []
    if STYLE_REFERENCE and STYLE_REFERENCE.exists():
        references.append(STYLE_REFERENCE)
    layout_reference = layout_reference_for(page, native_text=native_text)
    if layout_reference and layout_reference.exists():
        references.append(layout_reference)

    parts: list[dict[str, object]] = [{"text": build_prompt(page, native_text=native_text)}]
    for reference in references:
        parts.append(
            {
                "inline_data": {
                    "mime_type": mime_type_for_path(reference),
                    "data": base64.b64encode(reference.read_bytes()).decode("utf-8"),
                }
            }
        )

    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
            "imageConfig": {"aspectRatio": "16:9", "imageSize": "2K"},
        },
    }
    payload_json = json.dumps(payload, ensure_ascii=False).encode("utf-8")

    try:
        response = subprocess.run(
            [
                "curl",
                "-sS",
                "-X",
                "POST",
                url,
                "-H",
                f"x-goog-api-key: {api_key}",
                "-H",
                "Content-Type: application/json",
                "--data-binary",
                "@-",
            ],
            input=payload_json,
            capture_output=True,
            timeout=300,
            check=False,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError("Gemini image request timed out after 300 seconds.")

    if response.returncode != 0:
        raise RuntimeError(response.stderr.decode("utf-8", errors="replace") or f"curl exited {response.returncode}")

    try:
        response_payload = json.loads(response.stdout.decode("utf-8", errors="replace"))
    except json.JSONDecodeError:
        raise RuntimeError(response.stdout.decode("utf-8", errors="replace")[:480])

    if response_payload.get("error"):
        raise RuntimeError(summarize_http_error(json.dumps(response_payload, ensure_ascii=False)))

    parts = response_payload.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    for part in parts:
        inline_data = part.get("inlineData") or part.get("inline_data")
        if inline_data and inline_data.get("data"):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(base64.b64decode(inline_data["data"]))
            return "gemini-api"

    text_parts = [part.get("text", "") for part in parts if part.get("text")]
    raise RuntimeError("；".join(text_parts) if text_parts else "Gemini response did not contain inline image data.")


def draw_left_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    x: int,
    y: int,
    font,
    fill,
    max_width: int,
    spacing: int = 10,
    shadow: bool = True,
) -> tuple[int, int, int, int]:
    rendered = wrap_text(text, font, max_width)
    if shadow:
        draw.multiline_text((x + 2, y + 2), rendered, font=font, fill=(255, 255, 255, 180), spacing=spacing)
    draw.multiline_text((x, y), rendered, font=font, fill=fill, spacing=spacing)
    return draw.multiline_textbbox((x, y), rendered, font=font, spacing=spacing)


def draw_glass_card(
    overlay_draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    accent: tuple[int, int, int, int],
) -> None:
    x1, y1, x2, y2 = xy
    overlay_draw.rounded_rectangle(xy, radius=28, fill=(255, 255, 255, 176), outline=(255, 255, 255, 120), width=2)
    overlay_draw.rounded_rectangle((x1 + 16, y1 + 14, x1 + 28, y2 - 14), radius=6, fill=accent, outline=None)


def render_page(page: SlidePage, raw_path: Path, final_path: Path, index: int, total: int, native_text: bool = False) -> None:
    if native_text:
        image = Image.open(raw_path).convert("RGB").resize((WIDTH, HEIGHT))
        image.save(final_path, quality=95)
        return
    image = Image.open(raw_path).convert("RGB").resize((WIDTH, HEIGHT))
    image = Image.blend(image, Image.new("RGB", (WIDTH, HEIGHT), WHITE), 0.06).filter(ImageFilter.SMOOTH)
    image_rgba = image.convert("RGBA")
    overlay = Image.new("RGBA", (WIDTH, HEIGHT), (255, 255, 255, 0))
    overlay_draw = ImageDraw.Draw(overlay)

    accent_rgba = (31, 111, 178, 188) if index % 2 == 0 else (242, 154, 46, 188)
    title_x = 124
    title_y = 96
    probe = ImageDraw.Draw(Image.new("RGBA", (8, 8), (255, 255, 255, 0)))
    title_font = TITLE_SMALL_FONT if len(page.title) > 22 else TITLE_FONT
    title_rendered = wrap_text(page.title, title_font, 590)
    title_bbox = probe.multiline_textbbox((0, 0), title_rendered, font=title_font, spacing=8)
    title_h = title_bbox[3] - title_bbox[1]

    bullet_cards: list[tuple[int, str, tuple[int, int, int, int]]] = []
    current_y = title_y + title_h + 44
    for bullet_index, point in enumerate(page.bullets[:2]):
        rendered = wrap_text(f"• {point}", BODY_SMALL_FONT, 520)
        bbox = probe.multiline_textbbox((0, 0), rendered, font=BODY_SMALL_FONT, spacing=10)
        card_h = (bbox[3] - bbox[1]) + 30
        card_box = (124, current_y, 696, current_y + card_h)
        bullet_cards.append((bullet_index, rendered, card_box))
        current_y += card_h + 16

    panel_box = (88, 64, 740, current_y + 20)
    overlay_draw.rounded_rectangle(panel_box, radius=34, fill=(255, 255, 255, 150), outline=(255, 255, 255, 96), width=2)
    overlay_draw.rounded_rectangle((104, 84, 118, panel_box[3] - 24), radius=6, fill=accent_rgba, outline=None)

    for bullet_index, _, card_box in bullet_cards:
        fill_alpha = 122 if bullet_index == 0 else 102
        overlay_draw.rounded_rectangle(card_box, radius=22, fill=(255, 255, 255, fill_alpha), outline=(255, 255, 255, 68), width=1)
        overlay_draw.ellipse((144, card_box[1] + 18, 160, card_box[1] + 34), fill=accent_rgba)

    overlay_draw.rounded_rectangle((1598, 72, 1768, 128), radius=22, fill=(255, 255, 255, 120), outline=(255, 255, 255, 90), width=1)
    overlay_draw.rounded_rectangle((1618, 90, 1636, 108), radius=9, fill=accent_rgba, outline=None)
    overlay_draw.rounded_rectangle((1650, 90, 1668, 108), radius=9, fill=(31, 111, 178, 150), outline=None)

    image_rgba = Image.alpha_composite(image_rgba, overlay)
    draw = ImageDraw.Draw(image_rgba)

    title_box = draw_left_text(draw, page.title, title_x, title_y, title_font, BLUE, 590, spacing=8)
    draw.rounded_rectangle((title_x, title_box[3] + 18, title_x + 280, title_box[3] + 28), radius=5, fill=(242, 154, 46, 255), outline=None)
    draw.rounded_rectangle((title_x + 42, title_box[3] + 38, title_x + 228, title_box[3] + 44), radius=3, fill=(220, 238, 254, 255), outline=None)

    for _, rendered, card_box in bullet_cards:
        draw.multiline_text((172, card_box[1] + 12), rendered, font=BODY_SMALL_FONT, fill=NAVY, spacing=10)

    draw.text((1708, 86), f"{index + 1:02d}", font=CAPTION_FONT, fill=NAVY)
    image_rgba.convert("RGB").save(final_path, quality=95)


def create_contact_sheet(images: list[Path], output_path: Path) -> None:
    columns = 3
    rows = (len(images) + columns - 1) // columns
    thumb_w = 480
    thumb_h = 270
    padding = 28
    label_h = 56
    sheet_w = padding + columns * (thumb_w + padding)
    sheet_h = padding + rows * (thumb_h + label_h + padding)
    sheet = Image.new("RGB", (sheet_w, sheet_h), "#111827")
    draw = ImageDraw.Draw(sheet)
    font = load_font(20, heavy=True)

    for idx, image_path in enumerate(images):
        image = Image.open(image_path).convert("RGB").resize((thumb_w, thumb_h))
        row = idx // columns
        col = idx % columns
        x = padding + col * (thumb_w + padding)
        y = padding + row * (thumb_h + label_h + padding)
        sheet.paste(image, (x, y))
        draw.rounded_rectangle((x - 2, y - 2, x + thumb_w + 2, y + thumb_h + 2), radius=8, outline="#475569", width=2)
        draw.text((x, y + thumb_h + 14), image_path.stem, font=font, fill="#FFFFFF")

    sheet.save(output_path, quality=92)


def save_ppt(images: list[Path], output_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output_path))


def write_readme(output_dir: Path, raw_paths: list[Path], deck_title: str, pages: list[SlidePage]) -> None:
    lines = [
        f"这是根据《{SOURCE_DOC.name}》动态分析后生成的 {len(pages)} 页 Gemini 实际生图图片型 PPT。",
        "生成模式：逐页调用 Gemini 生图；本地将中文文字嵌入到画面内的轻量信息卡中，减少突兀感。",
        f"文档标题：{deck_title}",
        f"源文档：{SOURCE_DOC}",
        f"风格参考图：{STYLE_REFERENCE}",
        f"布局参考目录：{LAYOUT_REFERENCE_ROOT}",
        "",
        "页面规划：",
        *[f"P{index + 1}: {page.title}" for index, page in enumerate(pages)],
        "",
        "原始生图底图：",
        *[str(path) for path in raw_paths],
    ]
    (output_dir / "README.txt").write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the Gemini-generated core-competitiveness PPT deck.")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=OUTPUT_ROOT,
        help="Directory where timestamped local result folders are created.",
    )
    parser.add_argument(
        "--with-contact-sheet",
        action="store_true",
        help="Also write a local contact-sheet preview image.",
    )
    parser.add_argument(
        "--reuse-latest-raw",
        action="store_true",
        help="Reuse the latest generated raw Gemini backdrops and only rerender local integrated text layout.",
    )
    parser.add_argument(
        "--native-text",
        action="store_true",
        help="Ask Gemini to render the Chinese title and bullets directly in the image; disables local text overlay.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if not SOURCE_DOC.exists():
        raise RuntimeError(f"Source markdown not found: {SOURCE_DOC}")
    if STYLE_REFERENCE is None or not STYLE_REFERENCE.exists():
        raise RuntimeError(f"Style reference not found: {STYLE_REFERENCE}")

    deck_title, pages = build_pages_from_markdown(SOURCE_DOC.read_text(encoding="utf-8"))
    if args.native_text:
        label = f"项目核心竞争力分析-gemini原生出字生图{len(pages)}页样张"
    elif args.reuse_latest_raw:
        label = f"项目核心竞争力分析-gemini融合文字{len(pages)}页样张"
    else:
        label = f"项目核心竞争力分析-gemini图文融合生图{len(pages)}页样张"
    output_dir = create_timestamped_output_dir(args.output_root, label)
    raw_dir = output_dir / "gemini-backdrops"
    latest_output_dir = find_latest_output_dir(args.output_root) if args.reuse_latest_raw else None
    latest_raw_dir = latest_output_dir / "gemini-backdrops" if latest_output_dir else None

    image_paths: list[Path] = []
    raw_paths: list[Path] = []
    prompt_rows: list[dict[str, object]] = []
    total = len(pages)
    for idx, page in enumerate(pages):
        raw_path = raw_dir / f"P{idx + 1}-gemini-backdrop.png"
        final_path = output_dir / page.filename
        if args.reuse_latest_raw:
            if latest_raw_dir is None:
                raise RuntimeError("No latest raw Gemini backdrop directory found to reuse.")
            reuse_path = latest_raw_dir / f"P{idx + 1}-gemini-backdrop.png"
            if not reuse_path.exists():
                raise RuntimeError(f"Reusable raw Gemini backdrop not found: {reuse_path}")
            raw_path.parent.mkdir(parents=True, exist_ok=True)
            raw_path.write_bytes(reuse_path.read_bytes())
            mode = "reuse-latest-raw"
        else:
            mode = generate_gemini_image(page, raw_path, native_text=args.native_text)
        render_page(page, raw_path, final_path, idx, total, native_text=args.native_text)
        image_paths.append(final_path)
        raw_paths.append(raw_path)
        prompt_rows.append(
            {
                "page": idx + 1,
                "title": page.title,
                "bullets": page.bullets,
                "model": os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL),
                "prompt": build_prompt(page, native_text=args.native_text),
                "styleReference": str(STYLE_REFERENCE),
                "layoutReference": str(layout_reference_for(page, native_text=args.native_text) or ""),
                "mode": mode,
                "rawBackdrop": str(raw_path),
            }
        )

    pptx_path = output_dir / f"项目核心竞争力分析-Gemini动态真实生图-{len(pages)}页样张.pptx"
    save_ppt(image_paths, pptx_path)
    contact_sheet = output_dir / "contact-sheet.png"
    if args.with_contact_sheet:
        create_contact_sheet(image_paths, contact_sheet)

    (output_dir / "slide-plan.json").write_text(
        json.dumps(
            {
                "deckTitle": deck_title,
                "slideCount": len(pages),
                "slides": [
                    {"page": index + 1, "title": page.title, "bullets": page.bullets, "scenePrompt": page.scene_prompt, "filename": page.filename}
                    for index, page in enumerate(pages)
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    (output_dir / "gemini-prompts.json").write_text(
        json.dumps(prompt_rows, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    write_readme(output_dir, raw_paths, deck_title, pages)
    write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(image_paths),
        deckTitle=deck_title,
        sourceDoc=str(SOURCE_DOC),
        styleReference=str(STYLE_REFERENCE),
        textRichStyleReference=str(TEXT_RICH_STYLE_REFERENCE or ""),
        layoutReferenceRoot=str(LAYOUT_REFERENCE_ROOT or ""),
        readme=str(output_dir / "README.txt"),
        slidePlan=str(output_dir / "slide-plan.json"),
        promptCatalog=str(output_dir / "gemini-prompts.json"),
        contactSheet=str(contact_sheet) if args.with_contact_sheet else "",
        imageDirectory=str(output_dir),
        geminiBackdropDirectory=str(raw_dir),
        mode="gemini-native-text-per-page" if args.native_text else "gemini-generated-per-page",
    )
    print(pptx_path)


if __name__ == "__main__":
    main()

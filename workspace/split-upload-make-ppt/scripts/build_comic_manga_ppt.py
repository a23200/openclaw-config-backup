from __future__ import annotations

import argparse
import base64
import json
import os
import re
import subprocess
import textwrap
import urllib.parse
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw
from local_result import (
    create_timestamped_output_dir,
    first_existing_path,
    latest_archive_dir,
    latest_matching_path,
    load_env_candidates,
    repo_runtime_dir,
    write_result_manifest,
)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError:
    Presentation = None
    Inches = None

from build_clawlink_image_ppt_demo import BLUE, NAVY, WHITE, load_font, wrap_text


DEFAULT_GEMINI_IMAGE_MODEL = "gemini-3-pro-image-preview"
WIDTH = 1920
HEIGHT = 1080
ARCHIVE_DIR = latest_archive_dir()
OUTPUT_ROOT = repo_runtime_dir("image-ppt-demo")

STYLE_REFERENCE = first_existing_path(
    ARCHIVE_DIR / "desktop-generated" / "参考图-大模型能力突然变强.png" if ARCHIVE_DIR else None,
    Path("/Users/mac/Desktop/参考图-大模型能力突然变强.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png"),
)
LAYOUT_REFERENCE = first_existing_path(
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini原生出字生图*页样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-项目核心竞争力分析-gemini图文融合生图*页样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/P2-*.png"),
    latest_matching_path(repo_runtime_dir("image-ppt-demo"), "*-文旅引爆点-clawlink六页增强样张/gemini-backdrops/P2-gemini-backdrop.png"),
)

TITLE_FONT = load_font(58, heavy=True)
BODY_FONT = load_font(30)
SMALL_FONT = load_font(24)

IGNORED_SECTION_TITLES = {
    "目录",
    "contents",
    "附录",
    "参考资料",
    "致谢",
}


@dataclass(frozen=True)
class MarkdownSection:
    index: int
    level: int
    title: str
    lines: list[str]


@dataclass(frozen=True)
class SlidePage:
    page: int
    title: str
    bullets: list[str]
    context_title: str
    context_points: list[str]
    scene_prompt: str
    filename: str


def emit_progress(stage: str, message: str, percent: int, **extra: object) -> None:
    payload: dict[str, object] = {
        "stage": stage,
        "message": message,
        "percent": max(0, min(100, int(percent))),
    }
    payload.update(extra)
    print(f"__PROGRESS__ {json.dumps(payload, ensure_ascii=False)}", flush=True)


def emit_log(kind: str, title: str, body: str, **extra: object) -> None:
    payload: dict[str, object] = {
        "kind": kind,
        "title": title,
        "body": body,
    }
    if extra:
        payload["extra"] = extra
    print(f"__LOG__ {json.dumps(payload, ensure_ascii=False)}", flush=True)


def clean_markdown_text(text: str) -> str:
    cleaned = str(text or "")
    cleaned = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"`([^`]+)`", r"\1", cleaned)
    cleaned = re.sub(r"\*\*([^*]+)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"\*([^*]+)\*", r"\1", cleaned)
    cleaned = re.sub(r"^\s*[-*+]\s*", "", cleaned.strip())
    cleaned = re.sub(r"^\s*\d+[.)]\s*", "", cleaned.strip())
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip(" ：:|")


def shorten(text: str, limit: int) -> str:
    cleaned = clean_markdown_text(text)
    return cleaned if len(cleaned) <= limit else f"{cleaned[: limit - 1]}…"


def sanitize_filename_part(text: str) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|`]+", "", text)
    cleaned = re.sub(r"\s+", "", cleaned)
    return cleaned[:24] or "页面"


def mime_type_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    return "application/octet-stream"


def get_api_key() -> str:
    load_env_candidates()
    for key in ["GEMINI_API_KEY", "GOOGLE_API_KEY", "GOOGLE_GENAI_API_KEY"]:
        value = os.environ.get(key)
        if value:
            return value
    raise RuntimeError("Gemini API key not found in env.")


def parse_markdown_sections(markdown: str, source_path: Path) -> tuple[str, list[MarkdownSection]]:
    deck_title = source_path.stem
    sections: list[MarkdownSection] = []
    current_level = 0
    current_title = ""
    current_lines: list[str] = []
    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*$")

    for raw_line in markdown.splitlines():
        stripped = raw_line.strip()
        match = heading_re.match(stripped)
        if match:
            if current_title:
                sections.append(
                    MarkdownSection(
                        index=len(sections),
                        level=current_level,
                        title=current_title,
                        lines=current_lines,
                    )
                )
            current_level = len(match.group(1))
            current_title = clean_markdown_text(match.group(2))
            current_lines = []
            if current_level == 1 and current_title:
                deck_title = current_title
            continue
        if current_title:
            current_lines.append(raw_line)

    if current_title:
        sections.append(
            MarkdownSection(
                index=len(sections),
                level=current_level,
                title=current_title,
                lines=current_lines,
            )
        )

    return deck_title or source_path.stem, sections


def extract_section_points(section: MarkdownSection, max_points: int = 4) -> list[str]:
    candidates: list[str] = []
    for raw_line in section.lines:
        stripped = raw_line.strip()
        if not stripped or stripped == "---" or stripped.startswith("```"):
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            line = clean_markdown_text(stripped.replace("|", " | "))
            if line:
                candidates.append(line)
            continue

        if re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+[.)]\s+", stripped):
            line = clean_markdown_text(stripped)
            if line:
                candidates.append(line)
            continue

        line = clean_markdown_text(stripped)
        if not line:
            continue
        parts = re.split(r"[。；;]", line)
        for part in parts:
            part = clean_markdown_text(part)
            if len(part) >= 8:
                candidates.append(part)

    deduped: list[str] = []
    for item in candidates:
        normalized = item.lower()
        if not item or normalized in {existing.lower() for existing in deduped}:
            continue
        deduped.append(item)
        if len(deduped) >= max_points:
            break

    return deduped


def score_section(section: MarkdownSection, points: list[str]) -> int:
    score = len(points) * 6
    if section.level == 2:
        score += 8
    if re.search(r"(总结|结论|建议|展望|下一步|规划|策略|方案|路径|能力|优势|竞争力|洞察|分析|架构|系统|数据)", section.title):
        score += 10
    if re.search(r"\d|%|万|亿|年|月|日", section.title):
        score += 4
    if any(re.search(r"\d|%|万|亿|年|月|日", point) for point in points):
        score += 6
    if len(section.title) >= 6:
        score += 2
    return score


def build_global_points(sections: list[MarkdownSection]) -> list[str]:
    collected: list[str] = []
    for section in sections:
        points = extract_section_points(section, max_points=2)
        for point in points:
            short_point = shorten(point, 22)
            if short_point and short_point not in collected:
                collected.append(short_point)
            if len(collected) >= 2:
                return collected
    return collected or ["核心观点拆解", "关键动作落地"]


def preferred_sections(sections: list[MarkdownSection]) -> list[MarkdownSection]:
    level2 = [section for section in sections if section.level == 2]
    if level2:
        return level2
    level3 = [section for section in sections if section.level == 3]
    if level3:
        return level3
    return [section for section in sections if section.level >= 1]


def select_story_sections(sections: list[MarkdownSection], max_sections: int = 8) -> list[MarkdownSection]:
    candidates = [
        section
        for section in preferred_sections(sections)
        if section.title and section.title.lower() not in IGNORED_SECTION_TITLES
    ]
    if not candidates:
        return []

    intro_like = [section for section in candidates if re.search(r"(引言|概述|摘要|总览|背景)", section.title)]
    closing_like = [section for section in candidates if re.search(r"(结论|总结|建议|展望|下一步)", section.title)]
    excluded = {section.index for section in intro_like[:1]}
    middle_pool = [section for section in candidates if section.index not in excluded]

    if len(middle_pool) <= max_sections:
        return middle_pool

    reserved_ids: set[int] = set()
    if closing_like:
        reserved_ids.add(closing_like[0].index)

    slots = max_sections - len(reserved_ids)
    ranked = sorted(
        (
            (score_section(section, extract_section_points(section, max_points=4)), section)
            for section in middle_pool
            if section.index not in reserved_ids
        ),
        key=lambda item: (-item[0], item[1].index),
    )
    picked_ids = {section.index for _, section in ranked[: max(0, slots)]}
    picked_ids.update(reserved_ids)
    selected = [section for section in middle_pool if section.index in picked_ids]
    return selected[:max_sections]


def visual_scene_prompt(title: str, points: list[str], deck_title: str) -> str:
    joined = "；".join(points[:3])
    if re.search(r"(竞争力|优势|护城河|能力)", title):
        return f"hero capability map with layered shields, business engines, arrows and abstract modules, showing how the organization creates durable advantage around {joined}"
    if re.search(r"(分析|洞察|诊断|问题|现状)", title):
        return f"diagnostic comic storyboard with comparison cards, magnifier, dashboards and cause-effect flow, highlighting {joined}"
    if re.search(r"(策略|方案|路径|举措|规划|路线)", title):
        return f"action roadmap board with milestones, branching arrows, key actions and a clear next-step flow, focused on {joined}"
    if re.search(r"(架构|系统|平台|生态|协同)", title):
        return f"system blueprint scene with layered platforms, nodes, data links and modular capability blocks, illustrating {joined}"
    if re.search(r"(数据|指标|增长|收入|市场)", title):
        return f"growth dashboard scene with charts, market signals, data nodes and business momentum, centered on {joined}"
    if re.search(r"(结论|总结|建议|展望|下一步)", title):
        return f"closing editorial hero panel with concise decision cards, priority arrows and a forward-looking action path, summarizing {joined}"
    return f"premium business comic infographic for 《{deck_title}》, using editorial cards, arrows, dashboards and abstract symbols to explain {joined}"


def build_pages_from_markdown(markdown: str, source_path: Path) -> tuple[str, list[SlidePage]]:
    deck_title, sections = parse_markdown_sections(markdown, source_path)
    story_sections = select_story_sections(sections, max_sections=8)
    pages: list[SlidePage] = []

    cover_points = build_global_points(story_sections or sections)
    cover_title = shorten(deck_title.split("：", 1)[0], 18)
    pages.append(
        SlidePage(
            page=1,
            title=cover_title or "漫画总览",
            bullets=[shorten(point, 18) for point in cover_points[:2]],
            context_title=deck_title,
            context_points=cover_points[:4],
            scene_prompt=(
                f"editorial comic cover scene for a Chinese business PowerPoint about {deck_title}, "
                "with layered abstract modules, flowing arrows, key themes and a strong management-deck feel"
            ),
            filename=f"P1-{sanitize_filename_part(cover_title or deck_title)}.png",
        )
    )

    for index, section in enumerate(story_sections, start=2):
        context_points = extract_section_points(section, max_points=4) or [section.title]
        render_title = shorten(section.title, 18)
        render_bullets = [shorten(point, 18) for point in context_points[:2]]
        pages.append(
            SlidePage(
                page=index,
                title=render_title,
                bullets=render_bullets,
                context_title=section.title,
                context_points=context_points,
                scene_prompt=visual_scene_prompt(section.title, context_points, deck_title),
                filename=f"P{index}-{sanitize_filename_part(render_title)}.png",
            )
        )

    return deck_title, pages


def build_analysis_summary(markdown: str, source_path: Path, deck_title: str, pages: list[SlidePage]) -> str:
    _, sections = parse_markdown_sections(markdown, source_path)
    section_titles = [section.title for section in sections[:6] if section.title]
    selected_titles = [page.context_title for page in pages[1:5]]
    return (
        f"识别文档标题《{deck_title}》，共解析 {len(sections)} 个章节；"
        f"本次选取 {len(pages)} 页内容进行漫画化表达。"
        + (f" 重点章节：{' / '.join(section_titles)}。" if section_titles else "")
        + (f" 当前入选页面：{' / '.join(selected_titles)}。" if selected_titles else "")
    )


def build_design_summary(pages: list[SlidePage]) -> str:
    slide_titles = " / ".join(page.title for page in pages[:4])
    return (
        "设计上采用商业漫画信息图路线：白底留白、蓝橙强调色、粗描边、卡片化信息区，"
        "并让标题与要点直接融合进插画。"
        + (f" 首页到前几页的叙事顺序为：{slide_titles}。" if slide_titles else "")
    )


def build_prompt(page: SlidePage) -> str:
    context_points = "\n".join(f"- {point}" for point in page.context_points[:4])
    exact_lines = "\n".join([page.title, *[f"• {point}" for point in page.bullets[:2]]])
    return textwrap.dedent(
        f"""
        Create one finished 16:9 slide illustration for a Chinese business comic PowerPoint.
        The first reference image defines the target visual language.
        If a second reference image is provided, use it only for composition rhythm, text integration style and page density.

        Visual style requirements:
        - premium editorial comic infographic, not childish anime, not photorealistic, not 3D
        - white background, large clean negative space, crisp black outlines
        - deep blue accents, orange accents, light blue and light orange blocks
        - business storyboard feeling, designed for management presentation instead of entertainment poster
        - text must feel embedded inside the design, not pasted on top later
        - prefer diagrams, cards, dashboards, arrows and abstract symbols over detailed people
        - avoid cute mascots and avoid malformed human figures

        Narrative context:
        - page focus: {page.context_title}
        - key facts:
        {context_points}

        Main scene to generate:
        {page.scene_prompt}

        You must render these exact Chinese lines directly inside the image, with clean readable typography and correct spelling:
        {exact_lines}

        Text layout rules:
        - place all required text in one natural editorial zone near the upper-left or upper-middle area
        - the text zone must feel like part of the illustration, such as a designed information card or embedded panel
        - title larger, bullets smaller, all text crisp and aligned
        - do not add any extra text beyond the exact lines above
        - keep the text short, balanced and visually fused with the illustration

        Hard constraints:
        - no wrong characters, no garbled Chinese, no fake text, no pseudo text
        - no watermark, no logo, no extra labels
        - no giant disconnected title banner across the whole page
        - no creepy faces, no extra hands or legs
        """
    ).strip()


def summarize_http_error(error_body: str) -> str:
    try:
        payload = json.loads(error_body)
        message = payload.get("error", {}).get("message")
        status = payload.get("error", {}).get("status")
        if message and status:
            return f"{status}: {message}"
        if message:
            return message
    except Exception:
        pass
    return error_body[:480]


def generate_gemini_image(page: SlidePage, output_path: Path) -> str:
    api_key = get_api_key()
    model = urllib.parse.quote(os.environ.get("GEMINI_IMAGE_MODEL", DEFAULT_GEMINI_IMAGE_MODEL), safe="")
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    references: list[Path] = []
    for candidate in [STYLE_REFERENCE, LAYOUT_REFERENCE]:
        if candidate and candidate.exists() and candidate not in references:
            references.append(candidate)

    if not references:
        raise RuntimeError("At least one reference image is required for comic manga generation.")

    parts: list[dict[str, object]] = [{"text": build_prompt(page)}]
    for reference in references:
        parts.append(
            {
                "inline_data": {
                    "mime_type": mime_type_for_path(reference),
                    "data": base64.b64encode(reference.read_bytes()).decode("utf-8"),
                }
            }
        )

    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
            "imageConfig": {"aspectRatio": "16:9", "imageSize": "2K"},
        },
    }
    payload_json = json.dumps(payload, ensure_ascii=False).encode("utf-8")

    response = subprocess.run(
        [
            "curl",
            "-sS",
            "-X",
            "POST",
            url,
            "-H",
            f"x-goog-api-key: {api_key}",
            "-H",
            "Content-Type: application/json",
            "--data-binary",
            "@-",
        ],
        input=payload_json,
        capture_output=True,
        timeout=300,
        check=False,
    )

    if response.returncode != 0:
        raise RuntimeError(response.stderr.decode("utf-8", errors="replace") or f"curl exited {response.returncode}")

    try:
        response_payload = json.loads(response.stdout.decode("utf-8", errors="replace"))
    except json.JSONDecodeError:
        raise RuntimeError(response.stdout.decode("utf-8", errors="replace")[:480])

    if response_payload.get("error"):
        raise RuntimeError(summarize_http_error(json.dumps(response_payload, ensure_ascii=False)))

    parts = response_payload.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    for part in parts:
        inline_data = part.get("inlineData") or part.get("inline_data")
        if inline_data and inline_data.get("data"):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(base64.b64decode(inline_data["data"]))
            return "gemini-api"

    text_parts = [part.get("text", "") for part in parts if part.get("text")]
    raise RuntimeError("；".join(text_parts) if text_parts else "Gemini response did not contain inline image data.")


def generate_dry_run_image(page: SlidePage, output_path: Path) -> str:
    image = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((120, 90, 880, 420), radius=42, fill="#F8FBFF", outline="#DCEEFE", width=4)
    draw.rounded_rectangle((150, 130, 170, 380), radius=10, fill="#F29A2E")
    draw.rounded_rectangle((960, 120, 1760, 920), radius=52, fill="#F5FAFF", outline="#E2E8F0", width=3)
    draw.ellipse((1120, 240, 1540, 660), fill="#DCEEFE", outline="#1F6FB2", width=6)
    draw.rounded_rectangle((1120, 700, 1600, 820), radius=32, fill="#FDE9CF", outline="#F29A2E", width=4)
    draw.line((980, 760, 1100, 760), fill="#1F6FB2", width=8)
    draw.line((1088, 746, 1120, 760), fill="#1F6FB2", width=8)
    draw.line((1088, 774, 1120, 760), fill="#1F6FB2", width=8)

    title = wrap_text(page.title, TITLE_FONT, 620)
    draw.multiline_text((210, 132), title, font=TITLE_FONT, fill=BLUE, spacing=12)
    bullet_y = 250
    for bullet in page.bullets[:2]:
        rendered = wrap_text(f"• {bullet}", BODY_FONT, 560)
        draw.rounded_rectangle((210, bullet_y - 16, 780, bullet_y + 96), radius=24, fill="#FFFFFF", outline="#E2E8F0", width=2)
        draw.multiline_text((240, bullet_y), rendered, font=BODY_FONT, fill=NAVY, spacing=10)
        bullet_y += 124

    caption = wrap_text("Dry Run 仅用于联调进度与产物结构，正式模板会改走 Gemini 真生图。", SMALL_FONT, 700)
    draw.multiline_text((210, 870), caption, font=SMALL_FONT, fill="#64748B", spacing=8)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, quality=94)
    return "dry-run"


def build_outline_markdown(deck_title: str, pages: list[SlidePage]) -> str:
    lines = [f"# {deck_title}", ""]
    for page in pages:
        lines.append(f"## {page.page:02d}. {page.context_title}")
        for bullet in page.context_points[:4]:
            lines.append(f"- {bullet}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def save_ppt(images: list[Path], output_path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed in the active Python environment.")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def write_readme(
    output_dir: Path,
    source_path: Path,
    deck_title: str,
    pages: list[SlidePage],
    raw_paths: list[Path],
    outline_path: Path,
    pptx_path: Path,
) -> None:
    lines = [
        f"这是根据《{source_path.name}》自动分析后生成的漫画风图片型 PPT，共 {len(pages)} 页。",
        "生成模式：Google Gemini 真生图 + 参考图喂模 + 中文文字直接在图内生成。",
        f"文档标题：{deck_title}",
        f"源文档：{source_path}",
        f"风格参考图：{STYLE_REFERENCE}",
        f"布局参考图：{LAYOUT_REFERENCE}",
        f"大纲：{outline_path}",
        f"PPT：{pptx_path}",
        "",
        "页面规划：",
        *[f"P{page.page}: {page.context_title}" for page in pages],
        "",
        "原始生图底图：",
        *[str(path) for path in raw_paths],
    ]
    (output_dir / "README.txt").write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a comic-style PPT deck from Markdown using Gemini image generation.")
    parser.add_argument("--source", type=Path, required=True, help="Markdown source path.")
    parser.add_argument("--project-name", type=str, default="", help="Optional project title override.")
    parser.add_argument("--output-dir", type=Path, default=None, help="Exact output directory for generated assets.")
    parser.add_argument("--output-root", type=Path, default=OUTPUT_ROOT, help="Fallback root when output-dir is omitted.")
    parser.add_argument("--pptx-path", type=Path, default=None, help="Exact final PPTX output path.")
    parser.add_argument("--outline-path", type=Path, default=None, help="Exact outline markdown output path.")
    parser.add_argument("--dry-run", action="store_true", help="Skip Gemini and render local placeholders for integration testing.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    source_path = args.source.expanduser().resolve()
    if not source_path.exists():
        raise RuntimeError(f"Source markdown not found: {source_path}")
    if not args.dry_run and (STYLE_REFERENCE is None or not STYLE_REFERENCE.exists()):
        raise RuntimeError(f"Style reference not found: {STYLE_REFERENCE}")

    markdown = source_path.read_text(encoding="utf-8")
    output_dir = args.output_dir.resolve() if args.output_dir else create_timestamped_output_dir(args.output_root, f"{sanitize_filename_part(source_path.stem)}-漫画风PPT")
    output_dir.mkdir(parents=True, exist_ok=True)
    raw_dir = output_dir / "gemini-backdrops"
    raw_dir.mkdir(parents=True, exist_ok=True)

    emit_progress("prepare_markdown", "正在解析文档结构", 8)
    deck_title, pages = build_pages_from_markdown(markdown, source_path)
    if args.project_name.strip():
        deck_title = args.project_name.strip()
    emit_log("analysis", "文档分析结果", build_analysis_summary(markdown, source_path, deck_title, pages))
    emit_progress("plan_story", f"已规划 {len(pages)} 页漫画分镜", 18, total=len(pages))
    emit_log("design", "设计思路", build_design_summary(pages))

    prompt_rows: list[dict[str, object]] = []
    image_paths: list[Path] = []
    raw_paths: list[Path] = []
    total = len(pages)

    for index, page in enumerate(pages, start=1):
        percent = 18 + int(index / max(total, 1) * 66)
        emit_log(
            "page",
            f"第 {index}/{total} 页设计",
            f"页面主题《{page.context_title}》；画面重点：{page.scene_prompt}；文案：{' / '.join(page.bullets[:2])}",
            page=index,
            total=total,
        )
        emit_progress(
            "generate_pages",
            f"正在生成第 {index}/{total} 页漫画插画",
            percent,
            current=index,
            total=total,
        )
        raw_path = raw_dir / f"P{index}-gemini-backdrop.png"
        final_path = output_dir / page.filename
        mode = generate_dry_run_image(page, raw_path) if args.dry_run else generate_gemini_image(page, raw_path)
        image = Image.open(raw_path).convert("RGB").resize((WIDTH, HEIGHT))
        image.save(final_path, quality=95)
        image_paths.append(final_path)
        raw_paths.append(raw_path)
        prompt_rows.append(
            {
                "page": index,
                "title": page.title,
                "bullets": page.bullets,
                "contextTitle": page.context_title,
                "contextPoints": page.context_points,
                "scenePrompt": page.scene_prompt,
                "prompt": build_prompt(page),
                "styleReference": str(STYLE_REFERENCE or ""),
                "layoutReference": str(LAYOUT_REFERENCE or ""),
                "rawBackdrop": str(raw_path),
                "mode": mode,
            }
        )

    outline_path = args.outline_path.resolve() if args.outline_path else output_dir / "outline.md"
    outline_path.parent.mkdir(parents=True, exist_ok=True)
    outline_path.write_text(build_outline_markdown(deck_title, pages), encoding="utf-8")

    pptx_path = args.pptx_path.resolve() if args.pptx_path else output_dir / f"{sanitize_filename_part(deck_title)}-漫画风PPT.pptx"
    emit_progress("package_ppt", "正在打包 PPT 文件", 92)
    emit_log("progress", "打包导出", "全部页面已完成，正在打包成最终 PPT 文件。")
    save_ppt(image_paths, pptx_path)

    slide_plan_path = output_dir / "slide-plan.json"
    slide_plan_path.write_text(
        json.dumps(
            {
                "deckTitle": deck_title,
                "slideCount": len(pages),
                "slides": [
                    {
                        "page": page.page,
                        "title": page.title,
                        "contextTitle": page.context_title,
                        "bullets": page.context_points,
                        "scenePrompt": page.scene_prompt,
                        "filename": page.filename,
                    }
                    for page in pages
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    prompt_catalog_path = output_dir / "gemini-prompts.json"
    prompt_catalog_path.write_text(json.dumps(prompt_rows, ensure_ascii=False, indent=2), encoding="utf-8")
    write_readme(output_dir, source_path, deck_title, pages, raw_paths, outline_path, pptx_path)
    manifest_path = write_result_manifest(
        output_dir,
        pptx_path,
        slideCount=len(pages),
        deckTitle=deck_title,
        sourceDoc=str(source_path),
        outlinePath=str(outline_path),
        slidePlan=str(slide_plan_path),
        promptCatalog=str(prompt_catalog_path),
        imageDirectory=str(output_dir),
        geminiBackdropDirectory=str(raw_dir),
        styleReference=str(STYLE_REFERENCE or ""),
        layoutReference=str(LAYOUT_REFERENCE or ""),
        mode="comic-gemini-native-text",
        dryRun=args.dry_run,
    )

    result_payload = {
        "deckTitle": deck_title,
        "slideCount": len(pages),
        "pptxPath": str(pptx_path),
        "outlinePath": str(outline_path),
        "outputDir": str(output_dir),
        "slidePlanPath": str(slide_plan_path),
        "promptCatalogPath": str(prompt_catalog_path),
        "manifestPath": str(manifest_path),
        "dryRun": args.dry_run,
    }
    emit_progress("completed", "漫画风 PPT 已生成完成", 100, total=len(pages))
    emit_log("success", "完成", f"已输出 {len(pages)} 页 PPT，结果目录：{output_dir}")
    print(f"__RESULT__ {json.dumps(result_payload, ensure_ascii=False)}", flush=True)


if __name__ == "__main__":
    main()
